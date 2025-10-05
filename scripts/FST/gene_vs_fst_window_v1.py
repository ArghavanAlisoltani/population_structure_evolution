#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
GENE vs FST windows — full pipeline analogous to the TE version.

Inputs
------
1) FST windows CSV  (must include CHROM, WIN_START, WIN_END, WIN_MID, q_poi_3vs4/3vs5/4vs5)
2) Gene Excel       (Aria_curated_annotation_both_1ab_1.xlsx) — we use the first 3 columns:
                     Scaffold, start, end   (others are ignored for window counts)
3) GWAS Excel       (Significant_GWAS_list_Aria_v1.xlsx) — for optional GWAS overlays

Outputs
-------
• fst_windows_with_GENES.tsv      (FST windows augmented with GENE_count/density + GWAS columns)
• gene_presence_fisher_chisq.tsv  (Fisher & Chi-square presence tests for genes)
• Manhattan figures (Gb axis, scaffold-ordered):
   A) manhattan_q_panels_GENE.png                 (−log10(q), top-K both-low-q & high-count)
   B) manhattan_q_panels_GENE_GWASannot.png       (−log10(q) with GWAS-only annotations)
   C) manhattan_counts_GENE_count.png             (Y = raw GENE_count per window, 3 panels)
• A small PPTX listing key output plots

Tweak appearance in the CONFIG section below (fonts, colors, alpha, panel sizes, etc.).
"""

import os, re, math
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

from scipy.stats import norm, fisher_exact, chi2_contingency
import statsmodels.api as sm
from statsmodels.api import GLM
from statsmodels.genmod.families import links
from statsmodels.genmod.families.family import Poisson, NegativeBinomial
from pptx import Presentation
from pptx.util import Inches

# =========================
# CONFIG (easy to tweak)
# =========================
FST_CSV   = "fst_window_counts_tests_win10000000_ov2000000.csv"
GENE_XLSX = "Aria_curated_annotation_both_1ab_1.xlsx"  # same folder as script (or absolute path)
GWAS_XLSX = "Significant_GWAS_list_Aria_v1.xlsx"       # optional; used for GWAS overlays

Q_CUTOFF = 0.05    # significance threshold for Poisson q-values

OUTDIR = Path("gene_vs_fst_out")
OUTDIR.mkdir(parents=True, exist_ok=True)

# ------- Figure styling knobs -------
BASE_FONTSIZE    = 11         # axis/label font size
PANEL_ROW_HEIGHT = 3.6        # inches per row (each figure has 3 rows: 3vs4/3vs5/4vs5)
POINT_SIZE       = 10         # scatter point size
ALPHA_SIG        = 1.0        # alpha for q < cutoff (bold)
ALPHA_NONSIG     = 0.18       # alpha for q >= cutoff (faded)
COLOR_SCAF_A     = "#6BAED6"  # blue
COLOR_SCAF_B     = "#BDBDBD"  # gray
COLOR_TOP_LABELS = "red"
TOP_K_ANNOTATE   = 5          # labels per panel for the −log10(q) feature figure
SPACER_BP        = 1_000_000  # gap between scaffolds in genome coords (bp)

np.random.seed(1)

# =========================
# IO helpers
# =========================
def read_fst_windows(path: str) -> pd.DataFrame:
    df = pd.read_csv(path)
    need = ["CHROM","WIN_START","WIN_END","WIN_MID",
            "q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]
    miss = [c for c in need if c not in df.columns]
    if miss:
        raise ValueError(f"FST file missing required cols: {miss}")
    df["CHROM"]     = df["CHROM"].astype(str)
    df["WIN_START"] = pd.to_numeric(df["WIN_START"], errors="coerce").astype("Int64")
    df["WIN_END"]   = pd.to_numeric(df["WIN_END"],   errors="coerce").astype("Int64")
    df["WIN_MID"]   = pd.to_numeric(df["WIN_MID"],   errors="coerce").astype("Int64")
    df["WIN_LEN"]   = (df["WIN_END"] - df["WIN_START"] + 1).astype("Int64")
    # Clamp q to (0,1] so −log10 is finite and alpha rule works robustly
    for c in ["q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]:
        df[c] = pd.to_numeric(df[c], errors="coerce").clip(lower=1e-300, upper=1.0).fillna(1.0)
    return df

def read_gene_excel(path: str) -> pd.DataFrame:
    """
    We only need the first three columns (by name) for overlap:
    Scaffold, start, end
    """
    if not Path(path).exists():
        raise FileNotFoundError(f"Gene Excel file not found: {path}")
    gene = pd.read_excel(path, engine="openpyxl")
    need = ["Scaffold","start","end"]
    miss = [c for c in need if c not in gene.columns[:len(need)]]
    # If the sheet contains many columns but not in the first three, still try by names:
    if miss:
        if not all(col in gene.columns for col in need):
            raise ValueError(f"Gene Excel missing required columns: {need}. "
                             f"Found columns: {list(gene.columns)[:10]} ...")
    # Use columns by name (robust even if not first physically)
    g = gene[["Scaffold","start","end"]].copy()
    g["Scaffold"] = g["Scaffold"].astype(str)
    g["start"]    = pd.to_numeric(g["start"], errors="coerce").astype("Int64")
    g["end"]      = pd.to_numeric(g["end"],   errors="coerce").astype("Int64")
    g = g.dropna(subset=["Scaffold","start","end"]).reset_index(drop=True)
    g = g.sort_values(["Scaffold","start","end"])
    return g

def read_gwas(path: str) -> pd.DataFrame:
    """Expected columns: Scaffold, Pos, traits"""
    if not Path(path).exists():
        print(f"[WARN] GWAS file not found at {path}; GWAS-only panel will be skipped.")
        return pd.DataFrame(columns=["Scaffold","Pos","traits"])
    df = pd.read_excel(path, engine="openpyxl")
    need = ["Scaffold","Pos","traits"]
    if not all(c in df.columns for c in need):
        raise ValueError(f"GWAS Excel missing required cols: {need}")
    df["Scaffold"] = df["Scaffold"].astype(str)
    df["Pos"]      = pd.to_numeric(df["Pos"], errors="coerce").astype("Int64")
    df["traits"]   = df["traits"].astype(str)
    df = df.dropna(subset=["Scaffold","Pos"]).reset_index(drop=True)
    return df

# =========================
# Overlap counting (GENES per window)
# =========================
def _count_overlaps_one_chrom(wins: pd.DataFrame, feats: pd.DataFrame):
    """Count how many features (genes) intersect each window on a single scaffold."""
    w_st = wins["WIN_START"].to_numpy(dtype=np.int64)
    w_en = wins["WIN_END"].to_numpy(dtype=np.int64)
    nW   = len(wins)

    counts = np.zeros(nW, dtype=np.int64)
    if feats.empty:
        return counts

    f_st = feats["start"].to_numpy(dtype=np.int64)
    f_en = feats["end"].to_numpy(dtype=np.int64)

    j = 0
    for i in range(nW):
        ws, we = w_st[i], w_en[i]
        while j < len(f_st) and f_en[j] < ws:
            j += 1
        k = j
        while k < len(f_st) and f_st[k] <= we:
            if not (f_en[k] < ws or f_st[k] > we):
                counts[i] += 1
            k += 1
    return counts

def count_genes_per_window(fst_df: pd.DataFrame, gene_df: pd.DataFrame) -> pd.DataFrame:
    """Return per-window GENE_count and GENE_density."""
    base = fst_df[["CHROM","WIN_START","WIN_END","WIN_LEN","WIN_MID"]].copy().reset_index(drop=True)
    all_counts = np.zeros(len(base), dtype=np.int64)

    for chrom, wins in fst_df.groupby("CHROM", sort=False):
        rows = wins.index.to_numpy()
        wins_sorted = wins.sort_values("WIN_START")
        order_back  = np.argsort(wins_sorted.index.to_numpy())
        feats = gene_df[gene_df["Scaffold"] == chrom].sort_values("start")
        cts = _count_overlaps_one_chrom(wins_sorted, feats)
        all_counts[rows] = cts[order_back]

    out = base.copy()
    out["GENE_count"]   = all_counts.astype(int)
    out["GENE_density"] = out["GENE_count"] / out["WIN_LEN"].astype(float)
    return out

# =========================
# GLM helpers (Poisson + NegBin fallback)
# =========================
def dispersion_ratio(y, mu, df_adj=1):
    y = np.asarray(y, float); mu = np.asarray(mu, float)
    ok = mu > 0
    if ok.sum() <= 1:
        return np.nan
    return float(np.sum(((y[ok]-mu[ok])**2)/mu[ok]) / max(ok.sum()-df_adj, 1))

def clamp_exp(x):  # safe exp for CI
    return float(np.exp(np.clip(x, -700, 700)))

def fit_glm_poisson(y, X, offset=None):
    fam = Poisson(link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    return model.fit()

def estimate_nb_alpha(y, mu):
    numer = np.sum((y - mu) ** 2 - mu)
    denom = np.sum(mu ** 2)
    if denom <= 0:
        return 1.0
    alpha = numer / denom
    return float(max(alpha, 1e-8))

def fit_glm_nb(y, X, offset=None, alpha=1.0):
    fam = NegativeBinomial(alpha=alpha, link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    return model.fit()

def glm_compare_counts(y, group, offset_log=None, prefer_nb_if_overdispersed=True):
    X = pd.DataFrame({"intercept": 1.0, "group": group.astype(float)})
    res_p = fit_glm_poisson(y, X, offset=offset_log)
    mu_p  = res_p.fittedvalues
    disp  = dispersion_ratio(y, mu_p)

    res_used = res_p
    family_used = "Poisson"
    if prefer_nb_if_overdispersed and np.isfinite(disp) and disp > 1.5:
        alpha = estimate_nb_alpha(y, mu_p)
        res_nb = fit_glm_nb(y, X, offset=offset_log, alpha=alpha)
        res_used = res_nb
        family_used = f"NegBin(alpha={alpha:.3g})"

    b  = res_used.params.get("group", np.nan)
    se = res_used.bse.get("group", np.nan)
    z  = (b / se) if se and np.isfinite(se) and se > 0 else np.nan
    p  = 2 * (1 - norm.cdf(abs(z))) if np.isfinite(z) else np.nan
    rr = math.exp(b) if np.isfinite(b) else np.nan
    ci_lo = clamp_exp(b - 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan
    ci_hi = clamp_exp(b + 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan

    return {"family": family_used, "coef": b, "se": se, "z": z, "p": p,
            "RR": rr, "RR_CI_lo": ci_lo, "RR_CI_hi": ci_hi,
            "dispersion": disp, "res": res_used}

# =========================
# Presence tests (Fisher/Chi²) — presence = window has ≥1 gene
# =========================
def presence_tests(df, is_sig, count_col="GENE_count"):
    a = int((df.loc[ is_sig, count_col] >= 1).sum())
    b = int((df.loc[ is_sig, count_col] == 0).sum())
    c = int((df.loc[~is_sig, count_col] >= 1).sum())
    d = int((df.loc[~is_sig, count_col] == 0).sum())
    table = np.array([[a,b],[c,d]], dtype=int)
    try:
        from scipy.stats import fisher_exact, chi2_contingency
        odds, p_f = fisher_exact(table, alternative="two-sided")
        chi2, p_c, dof, _ = chi2_contingency(table, correction=False)
    except Exception:
        odds, p_f, chi2, p_c, dof = (np.nan, np.nan, np.nan, np.nan, np.nan)
    return {"a_sig_ge1":a, "b_sig_eq0":b, "c_not_ge1":c, "d_not_eq0":d,
            "odds_ratio_fisher":odds, "p_fisher":p_f, "chi2":chi2, "p_chi2":p_c, "df":dof}

# =========================
# Genome coordinates (Gb) & scaffold sort
# =========================
_scaf_pat = re.compile(r"scaffold[_\-]?(\d+)([ab])?$", re.IGNORECASE)
def scaffold_sort_key(s: str):
    m = _scaf_pat.search(s)
    if not m:
        return (1e9, s)
    num = int(m.group(1))
    suf = m.group(2) or ""
    suf_rank = {"a":0, "b":1}.get(suf.lower(), 2)
    return (num, suf_rank)

def build_genome_coords(df):
    scafs = sorted(df["CHROM"].unique(), key=scaffold_sort_key)
    offsets = {}
    current = 0
    for s in scafs:
        offsets[s] = current
        s_len = int(df.loc[df["CHROM"]==s, "WIN_END"].max())
        current += s_len + SPACER_BP
    df2 = df.copy()
    df2["X_bp"] = df2.apply(lambda r: offsets[r["CHROM"]] + float(r["WIN_MID"]), axis=1)
    df2["X_Gb"] = df2["X_bp"] / 1e9
    return df2, offsets

# =========================
# GWAS overlap columns to TSV
# =========================
def add_gwas_overlap_columns(df_aug: pd.DataFrame, gwas_df: pd.DataFrame) -> pd.DataFrame:
    if gwas_df is None or len(gwas_df) == 0:
        out = df_aug.copy()
        out["GWAS_n_hits"] = 0
        out["GWAS_positions"] = ""
        out["GWAS_traits"] = ""
        out["GWAS_has_hit"] = 0
        return out

    gw_by_scaf = {s: sub[["Pos","traits"]].astype({"Pos":int}).sort_values("Pos").reset_index(drop=True)
                  for s, sub in gwas_df.groupby("Scaffold")}
    n = len(df_aug)
    g_n   = np.zeros(n, dtype=int)
    g_pos = []
    g_trt = []
    g_fl  = np.zeros(n, dtype=int)

    for idx, r in df_aug.iterrows():
        sc = r["CHROM"]
        if sc not in gw_by_scaf:
            g_pos.append(""); g_trt.append(""); continue
        sub = gw_by_scaf[sc]
        ok  = (sub["Pos"] >= int(r["WIN_START"])) & (sub["Pos"] <= int(r["WIN_END"]))
        hits = sub.loc[ok]
        g_n[idx]  = len(hits)
        g_fl[idx] = 1 if len(hits) > 0 else 0
        if len(hits):
            pos_list = hits["Pos"].astype(int).tolist()
            trt_list = hits["traits"].astype(str).tolist()
            g_pos.append(",".join(map(str, pos_list[:10])) + ("..." if len(pos_list)>10 else ""))
            uniq_tr = []
            seen = set()
            for t in trt_list:
                if t not in seen:
                    uniq_tr.append(t); seen.add(t)
            g_trt.append(";".join(uniq_tr[:10]) + ("..." if len(uniq_tr)>10 else ""))
        else:
            g_pos.append(""); g_trt.append("")
    out = df_aug.copy()
    out["GWAS_n_hits"]   = g_n
    out["GWAS_positions"] = g_pos
    out["GWAS_traits"]    = g_trt
    out["GWAS_has_hit"]   = g_fl
    return out

# =========================
# Manhattan helpers
# =========================
def _alpha_from_q(qseries):
    return np.where(qseries < Q_CUTOFF, ALPHA_SIG, ALPHA_NONSIG)

def _axes_common_style(ax):
    ax.grid(axis="y", alpha=0.25)
    ax.tick_params(labelsize=BASE_FONTSIZE-1)

def _scatter_by_scaffold(ax, df, color_a=COLOR_SCAF_A, color_b=COLOR_SCAF_B, y_col="Y", alpha=None):
    scafs = sorted(df["CHROM"].unique(), key=scaffold_sort_key)
    for i, s in enumerate(scafs):
        sub = df[df["CHROM"]==s]
        col = color_a if (i % 2 == 0) else color_b
        a = alpha[sub.index] if isinstance(alpha, pd.Series) else alpha
        ax.scatter(sub["X_Gb"], sub[y_col], c=col, s=POINT_SIZE, alpha=a, edgecolors="none", zorder=2)

def _select_top_both(dfc, qcol, count_col, k):
    """Pick top-k among significant rows combining small q and large counts."""
    sig = dfc[dfc[qcol] < Q_CUTOFF].copy()
    if len(sig) == 0:
        return sig
    sig["r_q"] = sig[qcol].rank(method="min", ascending=True)
    sig["r_c"] = (-sig[count_col]).rank(method="min", ascending=True)
    sig["score"] = sig["r_q"] + sig["r_c"]
    return sig.nsmallest(k, "score")

def _annotate_labs(ax, rows, label_with_count=True, text_color=COLOR_TOP_LABELS):
    for _, r in rows.iterrows():
        lab = f"{r['CHROM']}:{int(r['WIN_START'])}-{int(r['WIN_END'])}"
        if label_with_count and "COUNT_FEAT" in rows.columns:
            lab += f" [{int(r['COUNT_FEAT'])}]"
        ax.scatter(r["X_Gb"], r["Y"], s=POINT_SIZE*1.4, c=text_color, zorder=5)
        ax.text(r["X_Gb"], r["Y"], lab, fontsize=BASE_FONTSIZE-1, color=text_color,
                rotation=45, ha="left", va="bottom")

# ---------- (A) −log10(q) panels for the GENE feature ----------
def manhattan_q_panels_gene(df_aug, outdir, gwas_df=None):
    """3-row −log10(q) panels (3vs4/3vs5/4vs5) for the GENE feature."""
    comps = ["3vs4","3vs5","4vs5"]
    qmap  = {c: f"q_poi_{c}" for c in comps}
    feature_col = "GENE_count"   # count feature used for composite ranking

    base = df_aug[["CHROM","WIN_START","WIN_END","WIN_MID", feature_col,
                   qmap["3vs4"], qmap["3vs5"], qmap["4vs5"]]].copy()
    base = base.rename(columns={feature_col: "COUNT_FEAT"})
    base, _ = build_genome_coords(base)
    for c in comps:
        base[f"neglog10_{c}"] = -np.log10(base[qmap[c]].astype(float).clip(lower=1e-300, upper=1.0))

    # (A1) top-K composite annotations (small q & high count)
    fig, axes = plt.subplots(3, 1, figsize=(14, PANEL_ROW_HEIGHT*3), sharex=True)
    for i, comp in enumerate(comps):
        ax = axes[i]
        qcol = qmap[comp]
        ycol = f"neglog10_{comp}"
        dfc = base.copy()
        dfc["Y"] = dfc[ycol]
        alpha = pd.Series(_alpha_from_q(dfc[qcol]), index=dfc.index)
        _scatter_by_scaffold(ax, dfc, y_col="Y", alpha=alpha)

        topk = _select_top_both(dfc, qcol=qcol, count_col="COUNT_FEAT", k=TOP_K_ANNOTATE)
        if len(topk):
            _annotate_labs(ax, topk, label_with_count=True)

        ax.set_ylabel(r"$-\log_{10}(q)$", fontsize=BASE_FONTSIZE)
        ax.set_title(f"{comp} — feature: GENE_count", fontsize=BASE_FONTSIZE+1)
        _axes_common_style(ax)
    axes[-1].set_xlabel("Genome position (Gb, scaffolds concatenated)", fontsize=BASE_FONTSIZE)
    fig.tight_layout()
    outA = outdir / f"manhattan_q_panels_GENE.png"
    fig.savefig(outA, dpi=220); plt.close(fig)

    # (A2) Optional GWAS overlay (red labels at significant windows overlapping ≥1 GWAS hit)
    outB = None
    if gwas_df is not None and len(gwas_df):
        gw_by_scaf = {s: sub["Pos"].astype(int).to_numpy() for s, sub in gwas_df.groupby("Scaffold")}
        gw_traits  = {s: sub.set_index("Pos")["traits"].astype(str) for s, sub in gwas_df.groupby("Scaffold")}

        fig2, axes2 = plt.subplots(3, 1, figsize=(14, PANEL_ROW_HEIGHT*3), sharex=True)
        for i, comp in enumerate(comps):
            ax = axes2[i]
            qcol = qmap[comp]
            ycol = f"neglog10_{comp}"
            dfc = base.copy()
            dfc["Y"] = dfc[ycol]
            alpha = pd.Series(_alpha_from_q(dfc[qcol]), index=dfc.index)
            _scatter_by_scaffold(ax, dfc, y_col="Y", alpha=alpha)

            hits = []
            for _, r in dfc[dfc[qcol] < Q_CUTOFF].iterrows():
                sc = r["CHROM"]
                if sc not in gw_by_scaf:
                    continue
                pos = gw_by_scaf[sc]
                ok  = (pos >= int(r["WIN_START"])) & (pos <= int(r["WIN_END"]))
                if np.any(ok):
                    first_pos = int(pos[ok][0])
                    trait     = gw_traits[sc].get(first_pos, "")
                    rr = r.copy()
                    rr["ann_text"] = f"{first_pos}:{trait}"
                    hits.append(rr)
            if hits:
                H = pd.DataFrame(hits)
                ax.scatter(H["X_Gb"], H["Y"], s=POINT_SIZE*1.4, c=COLOR_TOP_LABELS, zorder=5)
                for _, rr in H.iterrows():
                    ax.text(rr["X_Gb"], rr["Y"], rr["ann_text"], fontsize=BASE_FONTSIZE-1,
                            color=COLOR_TOP_LABELS, rotation=45, ha="left", va="bottom")

            ax.set_ylabel(r"$-\log_{10}(q)$", fontsize=BASE_FONTSIZE)
            ax.set_title(f"{comp} — feature: GENE_count (GWAS-overlap)", fontsize=BASE_FONTSIZE+1)
            _axes_common_style(ax)
        axes2[-1].set_xlabel("Genome position (Gb, scaffolds concatenated)", fontsize=BASE_FONTSIZE)
        fig2.tight_layout()
        outB = outdir / f"manhattan_q_panels_GENE_GWASannot.png"
        fig2.savefig(outB, dpi=220); plt.close(fig2)

    return outA, outB

# ---------- (B) GWAS-ONLY −log10(q) panels (no genes used) ----------
def manhattan_q_panels_GWAS_only(fst_df_with_coords, outdir, gwas_df):
    comps = ["3vs4","3vs5","4vs5"]
    qmap = {c: f"q_poi_{c}" for c in comps}
    base = fst_df_with_coords.copy()
    for c in comps:
        base[f"neglog10_{c}"] = -np.log10(base[qmap[c]].astype(float).clip(lower=1e-300, upper=1.0))

    # PRE-index GWAS
    gw_by_scaf = {}
    gw_traits  = {}
    if gwas_df is not None and len(gwas_df):
        gw_by_scaf = {s: sub["Pos"].astype(int).to_numpy() for s, sub in gwas_df.groupby("Scaffold")}
        gw_traits  = {s: sub.set_index("Pos")["traits"].astype(str) for s, sub in gwas_df.groupby("Scaffold")}

    fig, axes = plt.subplots(3, 1, figsize=(14, PANEL_ROW_HEIGHT*3), sharex=True)
    for i, comp in enumerate(comps):
        ax = axes[i]
        qcol = qmap[comp]
        ycol = f"neglog10_{comp}"
        dfc = base.copy()
        dfc["Y"] = dfc[ycol]
        alpha = pd.Series(np.where(dfc[qcol] < Q_CUTOFF, ALPHA_SIG, ALPHA_NONSIG), index=dfc.index)
        _scatter_by_scaffold(ax, dfc, y_col="Y", alpha=alpha)

        # overlay/annotate GWAS hit windows (significant only)
        if gw_by_scaf:
            hits = []
            for _, r in dfc[dfc[qcol] < Q_CUTOFF].iterrows():
                sc = r["CHROM"]
                if sc not in gw_by_scaf:
                    continue
                pos = gw_by_scaf[sc]
                ok  = (pos >= int(r["WIN_START"])) & (pos <= int(r["WIN_END"]))
                if np.any(ok):
                    first_pos = int(pos[ok][0])
                    trait     = gw_traits[sc].get(first_pos, "")
                    rr = r.copy()
                    rr["ann_text"] = f"{first_pos}:{trait}"
                    hits.append(rr)
            if hits:
                H = pd.DataFrame(hits)
                ax.scatter(H["X_Gb"], H["Y"], s=POINT_SIZE*1.4, c=COLOR_TOP_LABELS, zorder=5)
                for _, rr in H.iterrows():
                    ax.text(rr["X_Gb"], rr["Y"], rr["ann_text"], fontsize=BASE_FONTSIZE-1,
                            color=COLOR_TOP_LABELS, rotation=45, ha="left", va="bottom")

        ax.set_ylabel(r"$-\log_{10}(q)$", fontsize=BASE_FONTSIZE)
        ax.set_title(f"{comp} — GWAS-only overlay of significant windows", fontsize=BASE_FONTSIZE+1)
        _axes_common_style(ax)

    axes[-1].set_xlabel("Genome position (Gb, scaffolds concatenated)", fontsize=BASE_FONTSIZE)
    fig.tight_layout()
    out = outdir / "manhattan_q_panels_GWAS_ONLY.png"
    fig.savefig(out, dpi=220); plt.close(fig)
    return out

# ---------- (C) Gene COUNT PANELS (Y = raw counts) ----------
def manhattan_counts_panel_gene(df_aug, outdir):
    comps = ["3vs4","3vs5","4vs5"]
    qmap = {c: f"q_poi_{c}" for c in comps}

    base = df_aug[["CHROM","WIN_START","WIN_END","WIN_MID","GENE_count",
                   qmap["3vs4"], qmap["3vs5"], qmap["4vs5"]]].copy()
    base = base.rename(columns={"GENE_count":"Y"})
    base, _ = build_genome_coords(base)
    base["Y"] = pd.to_numeric(base["Y"], errors="coerce").fillna(0)

    fig, axes = plt.subplots(3, 1, figsize=(14, PANEL_ROW_HEIGHT*3), sharex=True)
    for i, comp in enumerate(comps):
        ax = axes[i]
        qcol = qmap[comp]
        alpha = pd.Series(_alpha_from_q(base[qcol]), index=base.index)
        _scatter_by_scaffold(ax, base, y_col="Y", alpha=alpha)
        ax.set_ylabel("GENE_count per window", fontsize=BASE_FONTSIZE)
        ax.set_title(f"{comp} — counts per window", fontsize=BASE_FONTSIZE+1)
        _axes_common_style(ax)
    axes[-1].set_xlabel("Genome position (Gb, scaffolds concatenated)", fontsize=BASE_FONTSIZE)
    fig.tight_layout()
    out = outdir / "manhattan_counts_GENE_count.png"
    fig.savefig(out, dpi=220); plt.close(fig)
    return out

# =========================
# MAIN
# =========================
def main():
    # 1) Load inputs
    fst   = read_fst_windows(FST_CSV)
    genes = read_gene_excel(GENE_XLSX)
    gwas  = read_gwas(GWAS_XLSX)

    # 2) Count genes per window
    gene_counts = count_genes_per_window(fst, genes)

    # 3) Build augmented FST windows table
    df_aug = fst.merge(gene_counts, on=["CHROM","WIN_START","WIN_END","WIN_LEN","WIN_MID"], how="left")
    num_cols = df_aug.select_dtypes(include=[np.number]).columns
    df_aug[num_cols] = df_aug[num_cols].fillna(0)

    # 4) Add GWAS overlap columns to TSV
    df_aug = add_gwas_overlap_columns(df_aug, gwas)

    # 5) Presence tests (significant vs not, presence = ≥1 gene)
    comps = ["3vs4","3vs5","4vs5"]
    comp_masks = {c: (df_aug[f"q_poi_{c}"] < Q_CUTOFF) for c in comps}
    rows = []
    for comp, is_sig in comp_masks.items():
        pres = presence_tests(df_aug, is_sig, count_col="GENE_count")
        rows.append({"comparison":comp,
                     "n_windows_sig":int(is_sig.sum()),
                     "n_windows_not":int((~is_sig).sum()),
                     "total_GENES_sig":int(df_aug.loc[is_sig,"GENE_count"].sum()),
                     "total_GENES_not":int(df_aug.loc[~is_sig,"GENE_count"].sum()),
                     **pres})
    pd.DataFrame(rows).to_csv(OUTDIR / "gene_presence_fisher_chisq.tsv", sep="\t", index=False)

    # 6) Manhattan families
    paths = []

    # (A) −log10(q) panels for the gene feature (with composite top-K labels)
    a, b = manhattan_q_panels_gene(df_aug, OUTDIR, gwas_df=gwas)
    paths += [p for p in (a,b) if p is not None]

    # (B) GWAS-only −log10(q) panels (no genes)
    fst_coords, _ = build_genome_coords(df_aug[["CHROM","WIN_START","WIN_END","WIN_MID",
                                                "q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]])
    paths.append(manhattan_q_panels_GWAS_only(fst_coords, OUTDIR, gwas_df=gwas))

    # (C) Gene COUNT panels
    paths.append(manhattan_counts_panel_gene(df_aug, OUTDIR))

    # 7) Save augmented TSV (with GWAS columns)
    out_aug = OUTDIR / "fst_windows_with_GENES.tsv"
    df_aug.to_csv(out_aug, sep="\t", index=False)

    # 8) Tiny PPTX listing example outputs
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "GENE vs FST — Manhattan suite (Gb axis, scaffold-ordered)"
    s.placeholders[1].text = (
        f"Q cutoff = {Q_CUTOFF}  |  Alpha: {ALPHA_SIG} (sig) / {ALPHA_NONSIG} (non-sig)\n"
        "A) −log10(q) for GENE_count with composite top-K labels (small q & high count)\n"
        "B) GWAS-only −log10(q) panels (no genes)\n"
        "C) GENE COUNT panels (raw counts per window)\n"
        f"Output dir: {OUTDIR.resolve()}"
    )
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Example outputs"
    tb = s2.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(3.6))
    tf = tb.text_frame; tf.clear()
    for p in paths[:10]:
        para = tf.paragraphs[0] if tf.text == "" else tf.add_paragraph()
        para.text = Path(p).name
    prs.save(str(OUTDIR / "GENE_vs_FST_manhattan_suite.pptx"))

    print("\nDone.")
    print(f"- Augmented TSV (GWAS cols): {out_aug.name}")
    print(f"- Presence tests TSV: gene_presence_fisher_chisq.tsv")
    print(f"- Figures in: {OUTDIR.resolve()}")
    if paths:
        print(f"- Examples: {', '.join([Path(p).name for p in paths[:5]])}")

if __name__ == "__main__":
    main()
