#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
TE vs FST windows — complete pipeline + new Manhattan plots:

  ✓ GLMs (density & counts), per-TE-type, Fisher/Chi² presence tests, PPTX
  ✓ New Manhattan panels (x in Gb, scaffold order numeric, alt colors):
      (A) −log10(q) per comparison, annot top-5 smallest q and top-5 largest counts
      (B) −log10(q) per comparison, annotate only windows that are q<0.05 AND
          overlap any significant GWAS hit (Position within window); label "pos:trait"

CUSTOMIZATION: see the CONFIG section below for colors, fonts, point sizes, etc.
"""

import os, re, math
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

from scipy.stats import norm, fisher_exact, chi2_contingency
import statsmodels.api as sm
from statsmodels.api import GLM
from statsmodels.genmod.families import links
from statsmodels.genmod.families.family import Poisson, NegativeBinomial
from statsmodels.stats.multitest import multipletests
from pptx import Presentation
from pptx.util import Inches

# =========================
# CONFIG (easy to tweak)
# =========================
FST_CSV = "fst_window_counts_tests_win10000000_ov2000000.csv"
TE_TSV  = os.path.expanduser(
    "~/Desktop/OSU_projects/conifers/LP/lodgepole_pine_assembly/col8_readable.TEanno.cds.scaf1ab.tsv"
)
GWAS_XLSX = "Significant_GWAS_list_Aria_v1.xlsx"   # put next to this script (or give absolute path)

Q_CUTOFF = 0.05

OUTDIR = Path("te_vs_fst_out")
OUTDIR.mkdir(parents=True, exist_ok=True)

# Figure style
BASE_FONTSIZE    = 11
PANEL_ROW_HEIGHT = 3.4     # inches per row for 3-row panels
POINT_SIZE       = 10
ALPHA_SIG        = 1.0     # alpha for q < Q_CUTOFF (bold)
ALPHA_NONSIG     = 0.18    # alpha for q ≥ Q_CUTOFF (faded)
COLOR_SCAF_A     = "#6BAED6"  # blue
COLOR_SCAF_B     = "#BDBDBD"  # gray
COLOR_TOP_LABELS = "red"
TOP_N_Q          = 5
TOP_N_COUNT      = 5
SPACER_BP        = 1_000_000  # gap between scaffolds in genome coords

np.random.seed(1)

# =========================
# IO & cleaning helpers
# =========================
def read_fst_windows(path: str) -> pd.DataFrame:
    df = pd.read_csv(path)
    need = ["CHROM","WIN_START","WIN_END","WIN_MID",
            "q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]
    miss = [c for c in need if c not in df.columns]
    if miss:
        raise ValueError(f"FST file missing required cols: {miss}")
    df["CHROM"]     = df["CHROM"].astype(str)
    df["WIN_START"] = pd.to_numeric(df["WIN_START"], errors="coerce").astype("Int64")
    df["WIN_END"]   = pd.to_numeric(df["WIN_END"],   errors="coerce").astype("Int64")
    df["WIN_MID"]   = pd.to_numeric(df["WIN_MID"],   errors="coerce").astype("Int64")
    df["WIN_LEN"]   = (df["WIN_END"] - df["WIN_START"] + 1).astype("Int64")
    # clamp q to (0,1] to avoid log10 issues if any zeros slipped in
    for c in ["q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]:
        df[c] = pd.to_numeric(df[c], errors="coerce").clip(lower=1e-300, upper=1.0).fillna(1.0)
    return df

def _has_header(path: str, first_col_name: str) -> bool:
    with open(path, "r") as fh:
        first = fh.readline().strip().split("\t")[0]
        return first == first_col_name

def read_te_table(path: str) -> pd.DataFrame:
    cols = ["seqid","source","sequence_ontology","start","end","score","strand"]
    te = pd.read_csv(path, sep="\t", names=cols,
                     header=0 if _has_header(path, cols[0]) else None)
    te["seqid"] = te["seqid"].astype(str)
    te["sequence_ontology"] = te["sequence_ontology"].astype(str)
    te["start"] = pd.to_numeric(te["start"], errors="coerce").astype("Int64")
    te["end"]   = pd.to_numeric(te["end"],   errors="coerce").astype("Int64")
    te = te.dropna(subset=["seqid","start","end"]).copy()
    te = te.sort_values(["seqid","start","end"]).reset_index(drop=True)
    return te

def read_gwas(path: str) -> pd.DataFrame:
    """
    Expected columns in Excel: Scaffold, Pos, MAF, p_value, neglogp, traits
    We'll use Scaffold (string), Pos (int), traits (string).
    """
    if not Path(path).exists():
        print(f"[WARN] GWAS file not found at {path}; GWAS-annotated panels will be skipped.")
        return pd.DataFrame(columns=["Scaffold","Pos","traits"])
    df = pd.read_excel(path)
    need = ["Scaffold","Pos","traits"]
    miss = [c for c in need if c not in df.columns]
    if miss:
        raise ValueError(f"GWAS Excel missing required cols: {miss}")
    df = df.copy()
    df["Scaffold"] = df["Scaffold"].astype(str)
    df["Pos"] = pd.to_numeric(df["Pos"], errors="coerce").astype("Int64")
    df["traits"] = df["traits"].astype(str)
    df = df.dropna(subset=["Scaffold","Pos"]).reset_index(drop=True)
    return df

# =========================
# Overlap counting
# =========================
def _count_overlaps_one_chrom(wins: pd.DataFrame, tes: pd.DataFrame, te_types: np.ndarray):
    w_st = wins["WIN_START"].to_numpy(dtype=np.int64)
    w_en = wins["WIN_END"].to_numpy(dtype=np.int64)
    nW   = len(wins)

    total_counts = np.zeros(nW, dtype=np.int64)
    type_to_idx = {t:i for i,t in enumerate(te_types)}
    per_type = np.zeros((nW, len(te_types)), dtype=np.int64)

    if tes.empty:
        return total_counts, per_type

    t_st = tes["start"].to_numpy(dtype=np.int64)
    t_en = tes["end"].to_numpy(dtype=np.int64)
    t_ty = tes["sequence_ontology"].astype(str).to_numpy()

    j = 0
    for i in range(nW):
        ws, we = w_st[i], w_en[i]
        while j < len(t_st) and t_en[j] < ws:
            j += 1
        k = j
        while k < len(t_st) and t_st[k] <= we:
            if not (t_en[k] < ws or t_st[k] > we):
                total_counts[i] += 1
                idx = type_to_idx.get(t_ty[k], None)
                if idx is not None:
                    per_type[i, idx] += 1
            k += 1
    return total_counts, per_type

def count_te_overlaps_per_window(fst_df: pd.DataFrame, te_df: pd.DataFrame):
    te_types = np.array(sorted(te_df["sequence_ontology"].astype(str).unique()))
    base = fst_df[["CHROM","WIN_START","WIN_END","WIN_LEN","WIN_MID"]].copy().reset_index(drop=True)

    total_counts_all = np.zeros(len(base), dtype=np.int64)
    per_type_all = np.zeros((len(base), len(te_types)), dtype=np.int64)

    for chrom, wins in fst_df.groupby("CHROM", sort=False):
        rows = wins.index.to_numpy()
        wins_sorted = wins.sort_values("WIN_START")
        order_back = np.argsort(wins_sorted.index.to_numpy())
        tes = te_df[te_df["seqid"] == chrom].sort_values("start")

        total_counts, per_type = _count_overlaps_one_chrom(wins_sorted, tes, te_types)
        total_counts_all[rows] = total_counts[order_back]
        per_type_all[rows, :]  = per_type[order_back, :]

    te_cols = {"TE_count": total_counts_all.astype(np.int64)}
    for j, t in enumerate(te_types):
        te_cols[f"TE_{t}"] = per_type_all[:, j].astype(np.int64)

    te_df_out = pd.DataFrame(te_cols)
    out = pd.concat([base, te_df_out], axis=1)

    win_len = out["WIN_LEN"].astype(float)
    out["TE_density"] = out["TE_count"] / win_len
    for t in te_types:
        out[f"TEdens_{t}"] = out[f"TE_{t}"] / win_len

    return out, te_types

# =========================
# GLM helpers (as before)
# =========================
def dispersion_ratio(y, mu, df_adj=1):
    y = np.asarray(y, float); mu = np.asarray(mu, float)
    ok = mu > 0
    if ok.sum() <= 1:
        return np.nan
    return float(np.sum(((y[ok]-mu[ok])**2)/mu[ok]) / max(ok.sum()-df_adj, 1))

def clamp_exp(x):  # safe exp for CI
    return float(np.exp(np.clip(x, -700, 700)))

def fit_glm_poisson(y, X, offset=None):
    fam = Poisson(link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    return model.fit()

def estimate_nb_alpha(y, mu):
    numer = np.sum((y - mu) ** 2 - mu)
    denom = np.sum(mu ** 2)
    if denom <= 0:
        return 1.0
    alpha = numer / denom
    return float(max(alpha, 1e-8))

def fit_glm_nb(y, X, offset=None, alpha=1.0):
    fam = NegativeBinomial(alpha=alpha, link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    return model.fit()

def glm_compare_counts(y, group, offset_log=None, prefer_nb_if_overdispersed=True):
    X = pd.DataFrame({"intercept": 1.0, "group": group.astype(float)})
    res_p = fit_glm_poisson(y, X, offset=offset_log)
    mu_p  = res_p.fittedvalues
    disp  = dispersion_ratio(y, mu_p)

    res_used = res_p
    family_used = "Poisson"
    if prefer_nb_if_overdispersed and np.isfinite(disp) and disp > 1.5:
        alpha = estimate_nb_alpha(y, mu_p)
        res_nb = fit_glm_nb(y, X, offset=offset_log, alpha=alpha)
        res_used = res_nb
        family_used = f"NegBin(alpha={alpha:.3g})"

    b  = res_used.params.get("group", np.nan)
    se = res_used.bse.get("group", np.nan)
    z  = (b / se) if se and np.isfinite(se) and se > 0 else np.nan
    p  = 2 * (1 - norm.cdf(abs(z))) if np.isfinite(z) else np.nan
    rr = math.exp(b) if np.isfinite(b) else np.nan
    ci_lo = clamp_exp(b - 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan
    ci_hi = clamp_exp(b + 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan

    return {"family": family_used, "coef": b, "se": se, "z": z, "p": p,
            "RR": rr, "RR_CI_lo": ci_lo, "RR_CI_hi": ci_hi,
            "dispersion": disp, "res": res_used}

# =========================
# Presence tests (Fisher/Chi²)
# =========================
def presence_tests(df, is_sig, count_col):
    a = int((df.loc[ is_sig, count_col] >= 1).sum())
    b = int((df.loc[ is_sig, count_col] == 0).sum())
    c = int((df.loc[~is_sig, count_col] >= 1).sum())
    d = int((df.loc[~is_sig, count_col] == 0).sum())
    table = np.array([[a,b],[c,d]], dtype=int)
    try: odds, p_f = fisher_exact(table, alternative="two-sided")
    except: odds, p_f = (np.nan, np.nan)
    try: chi2, p_c, dof, _ = chi2_contingency(table, correction=False)
    except: chi2, p_c, dof = (np.nan, np.nan, np.nan)
    return {"a_sig_ge1":a, "b_sig_eq0":b, "c_not_ge1":c, "d_not_eq0":d,
            "odds_ratio_fisher":odds, "p_fisher":p_f, "chi2":chi2, "p_chi2":p_c, "df":dof}

# =========================
# Genome coordinates (Gb) and scaffold order
# =========================
_scaf_pat = re.compile(r"scaffold[_\-]?(\d+)([ab])?$", re.IGNORECASE)
def scaffold_sort_key(s: str):
    """
    Order scaffolds in numeric order with 'a' before 'b':
      scaffold_1a < scaffold_1b < scaffold_2 < scaffold_3a ...
    Falls back to (len, string) if pattern not matched.
    """
    m = _scaf_pat.search(s)
    if not m:
        return (1e9, s)
    num = int(m.group(1))
    suf = m.group(2) or ""
    suf_rank = {"a":0, "b":1}.get(suf.lower(), 2)
    return (num, suf_rank)

def build_genome_coords(df):
    """
    Returns df2 with:
      X_bp : cumulative genome coordinate in bp (scaffold offsets + WIN_MID)
      X_Gb : X_bp / 1e9
      scaf_offsets : {scaffold: start_offset}
    """
    scafs = sorted(df["CHROM"].unique(), key=scaffold_sort_key)
    offsets = {}
    current = 0
    for s in scafs:
        offsets[s] = current
        s_len = int(df.loc[df["CHROM"]==s, "WIN_END"].max())
        current += s_len + SPACER_BP
    df2 = df.copy()
    df2["X_bp"] = df2.apply(lambda r: offsets[r["CHROM"]] + float(r["WIN_MID"]), axis=1)
    df2["X_Gb"] = df2["X_bp"] / 1e9
    return df2, offsets

# =========================
# Manhattan plotters (NEW)
# =========================
def _alpha_from_q(qseries):
    return np.where(qseries < Q_CUTOFF, ALPHA_SIG, ALPHA_NONSIG)

def _axes_common_style(ax):
    ax.grid(axis="y", alpha=0.25)
    ax.tick_params(labelsize=BASE_FONTSIZE-1)

def _scatter_by_scaffold(ax, df, color_a=COLOR_SCAF_A, color_b=COLOR_SCAF_B, y_col="Y", alpha=None):
    scafs = sorted(df["CHROM"].unique(), key=scaffold_sort_key)
    for i, s in enumerate(scafs):
        sub = df[df["CHROM"]==s]
        col = color_a if (i % 2 == 0) else color_b
        a = alpha[sub.index] if isinstance(alpha, pd.Series) else alpha
        ax.scatter(sub["X_Gb"], sub[y_col], c=col, s=POINT_SIZE, alpha=a, edgecolors="none", zorder=2)

def _annotate_labs(ax, rows, text_color=COLOR_TOP_LABELS):
    for _, r in rows.iterrows():
        lab = f"{r['CHROM']}:{int(r['WIN_START'])}-{int(r['WIN_END'])} [{int(r['COUNT_FEAT'])}]"
        ax.scatter(r["X_Gb"], r["Y"], s=POINT_SIZE*1.4, c=text_color, zorder=5)
        ax.text(r["X_Gb"], r["Y"], lab, fontsize=BASE_FONTSIZE-1, color=text_color,
                rotation=45, ha="left", va="bottom")

def manhattan_q_panels_by_feature(df_aug, feature_col, outdir, gwas_df=None):
    """
    Make TWO figures for a given feature (TE_count or TE_<class>):
      (A) −log10(q) panels, annotate union of top-5 by q and top-5 by counts.
      (B) −log10(q) panels, annotate windows that are q<cutoff AND overlap any GWAS hit.
    """
    comps = ["3vs4","3vs5","4vs5"]
    qmap = {c: f"q_poi_{c}" for c in comps}
    safe_feat = feature_col.replace("/", "_").replace(" ", "_")

    # Prepare base data
    base = df_aug[["CHROM","WIN_START","WIN_END","WIN_MID", feature_col,
                   qmap["3vs4"], qmap["3vs5"], qmap["4vs5"]]].copy()
    base = base.rename(columns={feature_col: "COUNT_FEAT"})
    base, offsets = build_genome_coords(base)  # adds X_bp, X_Gb
    # Pre-build −log10(q)
    for c in comps:
        base[f"neglog10_{c}"] = -np.log10(base[qmap[c]].astype(float).clip(lower=1e-300, upper=1.0))

    # (A) Top-q & top-count annotations
    fig, axes = plt.subplots(3, 1, figsize=(14, PANEL_ROW_HEIGHT*3), sharex=True)
    for i, comp in enumerate(comps):
        ax = axes[i]
        qcol = qmap[comp]
        ycol = f"neglog10_{comp}"
        dfc = base.copy()
        dfc["Y"] = dfc[ycol]
        alpha = pd.Series(_alpha_from_q(dfc[qcol]), index=dfc.index)
        _scatter_by_scaffold(ax, dfc, y_col="Y", alpha=alpha)

        # annotate union of top-5 by q and top-5 by count
        top_q = dfc.nsmallest(TOP_N_Q, qcol, keep="all")
        top_c = dfc.nlargest(TOP_N_COUNT, "COUNT_FEAT", keep="all")
        ann = pd.concat([top_q, top_c]).drop_duplicates()
        _annotate_labs(ax, ann)

        ax.set_ylabel(r"$-\log_{10}(q)$", fontsize=BASE_FONTSIZE)
        ax.set_title(f"{comp} — feature: {feature_col}", fontsize=BASE_FONTSIZE+1)
        _axes_common_style(ax)
    axes[-1].set_xlabel("Genome position (Gb, scaffolds concatenated)", fontsize=BASE_FONTSIZE)
    fig.tight_layout()
    outA = outdir / f"manhattan_q_panels_{safe_feat}.png"
    fig.savefig(outA, dpi=220); plt.close(fig)

    # (B) GWAS-overlap annotations (only if GWAS provided)
    outB = None
    if gwas_df is not None and len(gwas_df):
        # index GWAS by scaffold for quick lookup
        gw_by_scaf = {s: sub["Pos"].astype(int).to_numpy() for s, sub in gwas_df.groupby("Scaffold")}
        gw_traits  = {s: sub.set_index("Pos")["traits"].astype(str) for s, sub in gwas_df.groupby("Scaffold")}

        fig2, axes2 = plt.subplots(3, 1, figsize=(14, PANEL_ROW_HEIGHT*3), sharex=True)
        for i, comp in enumerate(comps):
            ax = axes2[i]
            qcol = qmap[comp]
            ycol = f"neglog10_{comp}"
            dfc = base.copy()
            dfc["Y"] = dfc[ycol]
            alpha = pd.Series(_alpha_from_q(dfc[qcol]), index=dfc.index)
            _scatter_by_scaffold(ax, dfc, y_col="Y", alpha=alpha)

            # find sig windows overlapping any GWAS pos
            hits = []
            for _, r in dfc[dfc[qcol] < Q_CUTOFF].iterrows():
                sc = r["CHROM"]
                if sc not in gw_by_scaf:
                    continue
                pos = gw_by_scaf[sc]
                ok = (pos >= int(r["WIN_START"])) & (pos <= int(r["WIN_END"]))
                if np.any(ok):
                    # annotate first few hits per window (avoid clutter)
                    first_pos = int(pos[ok][0])
                    trait = gw_traits[sc].get(first_pos, "")
                    rr = r.copy()
                    rr["ann_text"] = f"{first_pos}:{trait}"
                    hits.append(rr)
            if hits:
                H = pd.DataFrame(hits)
                # plot hits as red & annotated
                ax.scatter(H["X_Gb"], H["Y"], s=POINT_SIZE*1.4, c=COLOR_TOP_LABELS, zorder=5)
                for _, rr in H.iterrows():
                    ax.text(rr["X_Gb"], rr["Y"], rr["ann_text"], fontsize=BASE_FONTSIZE-1,
                            color=COLOR_TOP_LABELS, rotation=45, ha="left", va="bottom")

            ax.set_ylabel(r"$-\log_{10}(q)$", fontsize=BASE_FONTSIZE)
            ax.set_title(f"{comp} — feature: {feature_col} (GWAS-overlap annotations)", fontsize=BASE_FONTSIZE+1)
            _axes_common_style(ax)

        axes2[-1].set_xlabel("Genome position (Gb, scaffolds concatenated)", fontsize=BASE_FONTSIZE)
        fig2.tight_layout()
        outB = outdir / f"manhattan_q_panels_{safe_feat}_GWASannot.png"
        fig2.savefig(outB, dpi=220); plt.close(fig2)

    return outA, outB

# =========================
# Presence tests TSV (as in prior version)
# =========================
def run_presence_tests_and_write(df_aug, te_types, comp_masks, out_path):
    rows = []
    for comp, is_sig in comp_masks.items():
        # overall TE
        pres = presence_tests(df_aug, is_sig, "TE_count")
        rows.append({
            "comparison": comp,
            "feature": "TE_total",
            "n_windows_sig": int(is_sig.sum()),
            "n_windows_not": int((~is_sig).sum()),
            "total_TE_sig": int(df_aug.loc[is_sig,  "TE_count"].sum()),
            "total_TE_not": int(df_aug.loc[~is_sig, "TE_count"].sum()),
            **pres
        })
        # per TE class
        for t in te_types:
            col = f"TE_{t}"
            pre_t = presence_tests(df_aug, is_sig, col)
            rows.append({
                "comparison": comp,
                "feature": f"TE_{t}",
                "n_windows_sig": int(is_sig.sum()),
                "n_windows_not": int((~is_sig).sum()),
                "total_TE_sig": int(df_aug.loc[is_sig,  col].sum()),
                "total_TE_not": int(df_aug.loc[~is_sig, col].sum()),
                **pre_t
            })
    pd.DataFrame(rows).to_csv(out_path, sep="\t", index=False)

# =========================
# MAIN
# =========================
def main():
    # 1) Load inputs
    fst  = read_fst_windows(FST_CSV)
    te   = read_te_table(TE_TSV)
    gwas = read_gwas(GWAS_XLSX)

    # 2) Overlap → counts/densities
    te_counts, te_types = count_te_overlaps_per_window(fst, te)
    df_aug = fst.merge(te_counts, on=["CHROM","WIN_START","WIN_END","WIN_LEN","WIN_MID"], how="left")
    num_cols = df_aug.select_dtypes(include=[np.number]).columns
    df_aug[num_cols] = df_aug[num_cols].fillna(0)

    # 3) Significance masks per comparison
    comps = ["3vs4","3vs5","4vs5"]
    comp_masks = {c: (df_aug[f"q_poi_{c}"] < Q_CUTOFF) for c in comps}

    # 4) Presence tests TSV
    presence_tsv = OUTDIR / "te_presence_fisher_chisq.tsv"
    run_presence_tests_and_write(df_aug, te_types, comp_masks, presence_tsv)

    # 5) NEW Manhattan families
    #    (do once for overall TE_count, then for each TE class)
    paths = []
    a, b = manhattan_q_panels_by_feature(df_aug, "TE_count", OUTDIR, gwas_df=gwas)
    paths += [p for p in (a,b) if p is not None]
    for t in te_types:
        feat = f"TE_{t}"
        a, b = manhattan_q_panels_by_feature(df_aug, feat, OUTDIR, gwas_df=gwas)
        paths += [p for p in (a,b) if p is not None]

    # 6) Save augmented table for reference
    df_aug.to_csv(OUTDIR / "fst_windows_with_TEcounts.tsv", sep="\t", index=False)

    # 7) Small PPTX with entry points to the new figures
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "New Manhattan panels (−log10 q) with Gb axis"
    s.placeholders[1].text = (
        f"Q cutoff = {Q_CUTOFF}\n"
        "Panels: 3 rows (3vs4, 3vs5, 4vs5)\n"
        "Alpha: bold if q<cutoff, faded otherwise\n"
        "Annotations:\n"
        "  • Top-5 smallest q and top-5 largest counts per panel\n"
        "  • GWAS overlap (window q<cutoff AND has ≥1 GWAS hit) → red dot + 'pos:trait'\n"
        f"Presence tests TSV: {presence_tsv.name}\n"
        f"Output dir: {OUTDIR.resolve()}"
    )
    # Add one slide listing a few generated filenames
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Example outputs"
    tb = s2.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(3.6))
    tf = tb.text_frame; tf.clear()
    for p in paths[:10]:
        para = tf.paragraphs[0] if tf.text == "" else tf.add_paragraph()
        para.text = Path(p).name

    prs.save(str(OUTDIR / "TE_vs_FST_new_manhattan.pptx"))

    print("\nDone.")
    print(f"- Presence tests: {presence_tsv.name}")
    print(f"- New Manhattan figures written to: {OUTDIR.resolve()}")
    print(f"- Example files: {', '.join([Path(p).name for p in paths[:5]])}")

if __name__ == "__main__":
    main()

