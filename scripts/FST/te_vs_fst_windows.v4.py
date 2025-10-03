#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
TE overlap & density vs FST-enriched windows
Additions:
  • Fisher & Chi-square tests (presence-based 2x2) for overall TE and each TE class,
    per comparison (3vs4, 3vs5, 4vs5), with full inputs written to TSV.
  • Manhattan panels (3 rows: 3vs4, 3vs5, 4vs5) for TE counts and for each TE class count,
    alternating scaffold colors (gray/blue), alpha = 1 - q_poi_<comp>, top-10 labels.

This script supersedes the previous one; all earlier outputs are still produced.

INPUTS (same folder):
  fst_window_counts_tests_win10000000_ov2000000.csv
   — must include EXACT columns stated by the user including:
     CHROM, WIN_START, WIN_END, WIN_MID,
     q_poi_3vs4, q_poi_3vs5, q_poi_4vs5, ...
  TE TSV (absolute path):
     ~/Desktop/OSU_projects/conifers/LP/lodgepole_pine_assembly/col8_readable.TEanno.cds.scaf1ab.tsv
     columns: seqid, source, sequence_ontology, start, end, score, strand

Outputs in ./te_vs_fst_out/
"""

import os
import math
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

import statsmodels.api as sm
from statsmodels.api import GLM, families
from statsmodels.genmod.families import links
from statsmodels.stats.multitest import multipletests
from scipy.stats import norm, fisher_exact, chi2_contingency

from pptx import Presentation
from pptx.util import Inches

# ========== USER KNOBS (edit) ==========
FST_CSV = "fst_window_counts_tests_win10000000_ov2000000.csv"
TE_TSV  = os.path.expanduser(
    "~/Desktop/OSU_projects/conifers/LP/lodgepole_pine_assembly/col8_readable.TEanno.cds.scaf1ab.tsv"
)
Q_CUTOFF = 0.05
OUTDIR = Path("te_vs_fst_out")
OUTDIR.mkdir(parents=True, exist_ok=True)

# FIGURE STYLE (easy to tweak)
PANEL_ROW_HEIGHT = 3.0   # inches per row in Manhattan panels (increase for larger points)
BASE_FONTSIZE    = 11    # change axis/label font sizes
COLOR_SCAF_A     = "#6BAED6"  # blue for odd scaffolds
COLOR_SCAF_B     = "#BDBDBD"  # gray for even scaffolds
COLOR_TOP_LABELS = "red"
POINT_SIZE       = 8     # Manhattan point size
ALPHA_MIN        = 0.05  # floor for alpha

np.random.seed(1)

# =========================
# IO & cleaning
# =========================
def read_fst_windows(path: str) -> pd.DataFrame:
    df = pd.read_csv(path)
    need = ["CHROM","WIN_START","WIN_END","WIN_MID",
            "q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]
    missing = [c for c in need if c not in df.columns]
    if missing:
        raise ValueError(f"FST file missing required column(s): {missing}")
    df["CHROM"]     = df["CHROM"].astype(str)
    df["WIN_START"] = pd.to_numeric(df["WIN_START"], errors="coerce").astype("Int64")
    df["WIN_END"]   = pd.to_numeric(df["WIN_END"],   errors="coerce").astype("Int64")
    df["WIN_MID"]   = pd.to_numeric(df["WIN_MID"],   errors="coerce").astype("Int64")
    df["WIN_LEN"]   = (df["WIN_END"] - df["WIN_START"] + 1).astype("Int64")
    return df

def _has_header(path: str, first_col_name: str) -> bool:
    with open(path, "r") as fh:
        first = fh.readline().strip().split("\t")[0]
        return first == first_col_name

def read_te_table(path: str) -> pd.DataFrame:
    cols = ["seqid","source","sequence_ontology","start","end","score","strand"]
    te = pd.read_csv(path, sep="\t", names=cols,
                     header=0 if _has_header(path, cols[0]) else None)
    te["seqid"] = te["seqid"].astype(str)
    te["sequence_ontology"] = te["sequence_ontology"].astype(str)
    te["start"] = pd.to_numeric(te["start"], errors="coerce").astype("Int64")
    te["end"]   = pd.to_numeric(te["end"],   errors="coerce").astype("Int64")
    te = te.dropna(subset=["seqid","start","end"]).copy()
    te = te.sort_values(["seqid","start","end"]).reset_index(drop=True)
    return te

# =========================
# Overlap counting
# =========================
def _count_overlaps_one_chrom(wins: pd.DataFrame, tes: pd.DataFrame, te_types: np.ndarray):
    w_st = wins["WIN_START"].to_numpy(dtype=np.int64)
    w_en = wins["WIN_END"].to_numpy(dtype=np.int64)
    nW   = len(wins)

    total_counts = np.zeros(nW, dtype=np.int64)
    type_to_idx = {t:i for i,t in enumerate(te_types)}
    per_type = np.zeros((nW, len(te_types)), dtype=np.int64)

    if tes.empty:
        return total_counts, per_type

    t_st = tes["start"].to_numpy(dtype=np.int64)
    t_en = tes["end"].to_numpy(dtype=np.int64)
    t_ty = tes["sequence_ontology"].astype(str).to_numpy()

    j = 0
    for i in range(nW):
        ws, we = w_st[i], w_en[i]
        while j < len(t_st) and t_en[j] < ws:
            j += 1
        k = j
        while k < len(t_st) and t_st[k] <= we:
            if not (t_en[k] < ws or t_st[k] > we):
                total_counts[i] += 1
                idx = type_to_idx.get(t_ty[k], None)
                if idx is not None:
                    per_type[i, idx] += 1
            k += 1
    return total_counts, per_type

def count_te_overlaps_per_window(fst_df: pd.DataFrame, te_df: pd.DataFrame):
    te_types = np.array(sorted(te_df["sequence_ontology"].astype(str).unique()))
    base = fst_df[["CHROM","WIN_START","WIN_END","WIN_LEN","WIN_MID"]].copy().reset_index(drop=True)

    total_counts_all = np.zeros(len(base), dtype=np.int64)
    per_type_all = np.zeros((len(base), len(te_types)), dtype=np.int64)

    for chrom, wins in fst_df.groupby("CHROM", sort=False):
        rows = wins.index.to_numpy()
        wins_sorted = wins.sort_values("WIN_START")
        order_back = np.argsort(wins_sorted.index.to_numpy())
        tes = te_df[te_df["seqid"] == chrom].sort_values("start")

        total_counts, per_type = _count_overlaps_one_chrom(wins_sorted, tes, te_types)
        total_counts_all[rows] = total_counts[order_back]
        per_type_all[rows, :]  = per_type[order_back, :]

    te_cols = {"TE_count": total_counts_all.astype(np.int64)}
    for j, t in enumerate(te_types):
        te_cols[f"TE_{t}"] = per_type_all[:, j].astype(np.int64)

    te_df_out = pd.DataFrame(te_cols)
    out = pd.concat([base, te_df_out], axis=1)

    win_len = out["WIN_LEN"].astype(float)
    out["TE_density"] = out["TE_count"] / win_len
    for t in te_types:
        out[f"TEdens_{t}"] = out[f"TE_{t}"] / win_len

    return out, te_types

# =========================
# Modeling helpers
# =========================
def dispersion_ratio(y, mu, df_adj=1):
    y = np.asarray(y, float); mu = np.asarray(mu, float)
    ok = mu > 0
    if ok.sum() <= 1:
        return np.nan
    return float(np.sum(((y[ok]-mu[ok])**2)/mu[ok]) / max(ok.sum()-df_adj, 1))

def clamp_exp(x):
    return float(np.exp(np.clip(x, -700, 700)))

from statsmodels.genmod.families.family import Poisson, NegativeBinomial
def fit_glm_poisson(y, X, offset=None):
    fam = Poisson(link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    return model.fit()

def estimate_nb_alpha(y, mu):
    numer = np.sum((y - mu) ** 2 - mu)
    denom = np.sum(mu ** 2)
    if denom <= 0:
        return 1.0
    alpha = numer / denom
    return float(max(alpha, 1e-8))

def fit_glm_nb(y, X, offset=None, alpha=1.0):
    fam = NegativeBinomial(alpha=alpha, link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    return model.fit()

def glm_compare_counts(y, group, offset_log=None, prefer_nb_if_overdispersed=True):
    """
    GLM: count ~ group (0/1), optional offset = log(exposure)
    Returns dict(coef,se,z,p,RR,CI,dispersion,family,res).
    """
    X = pd.DataFrame({"intercept": 1.0, "group": group.astype(float)})

    res_p = fit_glm_poisson(y, X, offset=offset_log)
    mu_p  = res_p.fittedvalues
    disp  = dispersion_ratio(y, mu_p)

    res_used = res_p
    family_used = "Poisson"

    if prefer_nb_if_overdispersed and np.isfinite(disp) and disp > 1.5:
        alpha = estimate_nb_alpha(y, mu_p)
        res_nb = fit_glm_nb(y, X, offset=offset_log, alpha=alpha)
        res_used = res_nb
        family_used = f"NegBin(alpha={alpha:.3g})"

    b  = res_used.params.get("group", np.nan)
    se = res_used.bse.get("group", np.nan)
    z  = (b / se) if se and np.isfinite(se) and se > 0 else np.nan
    p  = 2 * (1 - norm.cdf(abs(z))) if np.isfinite(z) else np.nan
    rr = math.exp(b) if np.isfinite(b) else np.nan
    ci_lo = clamp_exp(b - 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan
    ci_hi = clamp_exp(b + 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan

    return {"family": family_used, "coef": b, "se": se, "z": z, "p": p,
            "RR": rr, "RR_CI_lo": ci_lo, "RR_CI_hi": ci_hi,
            "dispersion": disp, "res": res_used}

def global_rate_ratio_counts(y, mask):
    """Quick mean ratio of counts per window between groups."""
    m1 = np.mean(y[mask])  if mask.sum()>0 else np.nan
    m0 = np.mean(y[~mask]) if (~mask).sum()>0 else np.nan
    rr = (m1/m0) if (m0 not in (0, np.nan) and np.isfinite(m0) and np.isfinite(m1)) else np.nan
    return rr, m1, m0

# =========================
# Fisher & Chi-square (presence-based 2x2)
# =========================
def presence_tests(df, is_sig, count_col):
    """
    Build a 2x2 table on window presence (>=1 vs 0) and do Fisher & Chi-square.
    Returns dict with all inputs and results.
    """
    a = int((df.loc[ is_sig, count_col] >= 1).sum())  # sig windows with ≥1
    b = int((df.loc[ is_sig, count_col] == 0).sum())  # sig windows with 0
    c = int((df.loc[~is_sig, count_col] >= 1).sum())  # not windows with ≥1
    d = int((df.loc[~is_sig, count_col] == 0).sum())  # not windows with 0

    table = np.array([[a, b],
                      [c, d]], dtype=int)

    # Fisher exact (two-sided)
    try:
        odds, p_fisher = fisher_exact(table, alternative='two-sided')
    except Exception:
        odds, p_fisher = (np.nan, np.nan)

    # Chi-square test of independence
    try:
        chi2, p_chi2, dof, exp = chi2_contingency(table, correction=False)
    except Exception:
        chi2, p_chi2, dof = (np.nan, np.nan, np.nan)

    return {
        "a_sig_ge1": a, "b_sig_eq0": b,
        "c_not_ge1": c, "d_not_eq0": d,
        "odds_ratio_fisher": odds, "p_fisher": p_fisher,
        "chi2": chi2, "p_chi2": p_chi2, "df": dof
    }

# =========================
# Manhattan (counts) helpers
# =========================
def build_concat_coords(df):
    """
    Build a genome-wide x coordinate by concatenating scaffolds in alphanumeric order.
    Returns df with X and a tick table for plotting.
    """
    scafs = sorted(df["CHROM"].unique(), key=lambda s: (len(s), s))
    offsets = {}
    current = 0
    ticks_pos = []
    ticks_lab = []
    for i, s in enumerate(scafs):
        s_len = int(df.loc[df["CHROM"]==s, "WIN_END"].max())
        offsets[s] = current
        mid = current + s_len/2.0
        ticks_pos.append(mid)
        ticks_lab.append(s)
        current += s_len + 1_000_000  # spacer; change if desired
    X = df.apply(lambda r: offsets[r["CHROM"]] + float(r["WIN_MID"]), axis=1)
    tick_df = pd.DataFrame({"pos": ticks_pos, "lab": ticks_lab})
    return X, tick_df, offsets

def annotate_top(ax, sub, comp, qcol, how_many=10):
    top = sub.nsmallest(how_many, qcol)
    for _, r in top.iterrows():
        label = f"{r['CHROM']}:{int(r['WIN_START'])}-{int(r['WIN_END'])}"
        ax.scatter(r["X"], r["Y"], s=POINT_SIZE*1.4, c=COLOR_TOP_LABELS, zorder=5)
        ax.text(r["X"], r["Y"], label, fontsize=BASE_FONTSIZE-1, color=COLOR_TOP_LABELS,
                rotation=45, ha="left", va="bottom")

def manhattan_counts_panel(df_aug, te_types, outdir, feature="TE_count"):
    """
    Make a 3-row Manhattan panel for `feature` which is either:
      - "TE_count"  (overall)
      - "TE_<TEtype>"  (a specific class)
    Y = counts per window for that feature.
    Color alternates per scaffold; alpha = 1 - q_poi_<comp>.
    """
    comps = ["3vs4","3vs5","4vs5"]
    qcols = {c: f"q_poi_{c}" for c in comps}

    # Build genome-wide X for the whole table once
    df_local = df_aug[["CHROM","WIN_START","WIN_END","WIN_MID", feature, *qcols.values()]].copy()
    df_local = df_local.rename(columns={feature: "Y"})
    df_local["Y"] = pd.to_numeric(df_local["Y"], errors="coerce").fillna(0)

    X, ticks, offsets = build_concat_coords(df_local)
    df_local["X"] = X

    # figure
    fig, axes = plt.subplots(3, 1, figsize=(14, PANEL_ROW_HEIGHT*3), sharex=True)
    for i, comp in enumerate(comps):
        ax = axes[i]
        qcol = qcols[comp]
        df_local["alpha"] = 1.0 - df_local[qcol].clip(lower=0, upper=1).astype(float)
        df_local["alpha"] = df_local["alpha"].clip(lower=ALPHA_MIN, upper=1.0)

        # alternate scaffold colors
        scafs = sorted(df_local["CHROM"].unique(), key=lambda s: (len(s), s))
        for si, s in enumerate(scafs):
            sub = df_local[df_local["CHROM"]==s].copy()
            color = COLOR_SCAF_A if (si % 2 == 0) else COLOR_SCAF_B
            ax.scatter(sub["X"], sub["Y"], s=POINT_SIZE,
                       c=color, alpha=sub["alpha"], edgecolors='none')

        # annotate top-10 windows by q (red)
        annotate_top(ax, df_local, comp, qcol, how_many=10)

        ax.set_ylabel(f"{feature} per window", fontsize=BASE_FONTSIZE)
        ax.set_title(f"{comp} — alpha = 1 − {qcol}", fontsize=BASE_FONTSIZE+1)
        ax.grid(axis="y", alpha=0.25)

    axes[-1].set_xticks(ticks["pos"].to_list(), ticks["lab"].to_list(), rotation=60)
    axes[-1].set_xlabel("Genome (scaffolds concatenated)", fontsize=BASE_FONTSIZE)
    for ax in axes:
        ax.tick_params(labelsize=BASE_FONTSIZE-1)

    fig.tight_layout()
    safe_feat = feature.replace("/", "_").replace(" ", "_")
    out = outdir / f"manhattan_counts_{safe_feat}.png"
    fig.savefig(out, dpi=220)
    plt.close(fig)
    return out

# =========================
# Report helper
# =========================
def add_slide(prs, title, bullets=None, image=None):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if bullets:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.6))
        tf = tb.text_frame; tf.clear()
        for i,b in enumerate(bullets):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.text = b; p.level = 0
    if image and Path(image).exists():
        slide.shapes.add_picture(str(image), Inches(0.5), Inches(2.7), height=Inches(3.8))

# =========================
# MAIN
# =========================
def main():
    # Load data
    fst = read_fst_windows(FST_CSV)
    te  = read_te_table(TE_TSV)

    # Overlaps
    te_counts, te_types = count_te_overlaps_per_window(fst, te)
    df_aug = fst.merge(te_counts, on=["CHROM","WIN_START","WIN_END","WIN_LEN","WIN_MID"], how="left")
    num_cols = df_aug.select_dtypes(include=[np.number]).columns
    df_aug[num_cols] = df_aug[num_cols].fillna(0)

    # === DENSITY & COUNT GLMs (unchanged from prior version, omitted for brevity) ===
    # (We re-run the dual-track analyses to keep all previous outputs)
    comps = ["3vs4","3vs5","4vs5"]
    comp_masks_density = {}
    per_comp_overall_density = {}
    per_comp_types_density = {}
    comp_masks_count = {}
    per_comp_overall_count = {}
    per_comp_types_count = {}

    for comp in comps:
        qcol = f"q_poi_{comp}"
        is_sig = (df_aug[qcol] < Q_CUTOFF)

        # Density track
        y_den = df_aug["TE_count"].astype(int).to_numpy()
        grp   = is_sig.astype(int).to_numpy()
        off   = np.log(df_aug["WIN_LEN"].astype(float).to_numpy())
        res_den = glm_compare_counts(y_den, grp, offset_log=off)
        # per-type (density)
        rows_den = []
        for t in te_types:
            y_t = df_aug[f"TE_{t}"].astype(int).to_numpy()
            rows_den.append({
                "TE_type": t, **glm_compare_counts(y_t, grp, offset_log=off)
            })
        df_types_den = pd.DataFrame(rows_den)
        df_types_den["q"] = multipletests(df_types_den["p"].fillna(1.0), method="fdr_bh")[1]
        df_types_den = df_types_den.sort_values("q")

        comp_masks_density[comp] = is_sig
        per_comp_overall_density[comp] = {
            "model_used": res_den["family"], "RR_overall": res_den["RR"],
            "RR_CI_lo": res_den["RR_CI_lo"], "RR_CI_hi": res_den["RR_CI_hi"],
            "p_model": res_den["p"], "dispersion_est": res_den["dispersion"],
            "n_sig": int(is_sig.sum()), "n_not": int((~is_sig).sum())
        }
        per_comp_types_density[comp] = df_types_den
        df_types_den.to_csv(OUTDIR / f"te_density_by_type_{comp}.tsv", sep="\t", index=False)

        # Count track (no offset)
        y_cnt = df_aug["TE_count"].astype(int).to_numpy()
        res_cnt = glm_compare_counts(y_cnt, grp, offset_log=None)
        rows_cnt = []
        for t in te_types:
            y_t = df_aug[f"TE_{t}"].astype(int).to_numpy()
            rows_cnt.append({
                "TE_type": t, **glm_compare_counts(y_t, grp, offset_log=None)
            })
        df_types_cnt = pd.DataFrame(rows_cnt)
        df_types_cnt["q"] = multipletests(df_types_cnt["p"].fillna(1.0), method="fdr_bh")[1]
        df_types_cnt = df_types_cnt.sort_values("q")

        comp_masks_count[comp] = is_sig
        rr_cnt, m1, m0 = global_rate_ratio_counts(y_cnt, is_sig.values)
        per_comp_overall_count[comp] = {
            "model_used": res_cnt["family"], "RR_overall": res_cnt["RR"],
            "RR_CI_lo": res_cnt["RR_CI_lo"], "RR_CI_hi": res_cnt["RR_CI_hi"],
            "p_model": res_cnt["p"], "dispersion_est": res_cnt["dispersion"],
            "n_sig": int(is_sig.sum()), "n_not": int((~is_sig).sum()),
            "mean_TE_sig": float(m1), "mean_TE_not": float(m0), "RR_means_sig_over_not": float(rr_cnt)
        }
        per_comp_types_count[comp] = df_types_cnt
        df_types_cnt.to_csv(OUTDIR / f"te_count_by_type_{comp}.tsv", sep="\t", index=False)

    # Write augmented windows
    df_aug.to_csv(OUTDIR / "fst_windows_with_TEcounts.tsv", sep="\t", index=False)

    # =========================
    # Fisher & Chi-square (presence-based)
    # =========================
    fisher_rows = []
    for comp in comps:
        is_sig = comp_masks_density[comp]  # same mask for all tests for this comp

        # Overall TE presence
        pres = presence_tests(df_aug, is_sig, "TE_count")
        fisher_rows.append({
            "comparison": comp,
            "feature": "TE_total",
            "n_windows_sig": int(is_sig.sum()),
            "n_windows_not": int((~is_sig).sum()),
            "total_TE_sig": int(df_aug.loc[is_sig, "TE_count"].sum()),
            "total_TE_not": int(df_aug.loc[~is_sig, "TE_count"].sum()),
            **pres
        })

        # Each TE class presence
        for t in te_types:
            col = f"TE_{t}"
            pres_t = presence_tests(df_aug, is_sig, col)
            fisher_rows.append({
                "comparison": comp,
                "feature": f"TE_{t}",
                "n_windows_sig": int(is_sig.sum()),
                "n_windows_not": int((~is_sig).sum()),
                "total_TE_sig": int(df_aug.loc[is_sig, col].sum()),
                "total_TE_not": int(df_aug.loc[~is_sig, col].sum()),
                **pres_t
            })

    fisher_df = pd.DataFrame(fisher_rows)
    fisher_out = OUTDIR / "te_presence_fisher_chisq.tsv"
    fisher_df.to_csv(fisher_out, sep="\t", index=False)

    # =========================
    # Manhattan panels (counts)
    # =========================
    # Overall TE counts
    manh_overall = manhattan_counts_panel(df_aug, te_types, OUTDIR, feature="TE_count")
    # Each TE class
    manh_teclass_paths = []
    for t in te_types:
        manh_teclass_paths.append(
            manhattan_counts_panel(df_aug, te_types, OUTDIR, feature=f"TE_{t}")
        )

    # =========================
    # PPTX
    # =========================
    prs = Presentation()
    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TE vs FST windows — density & count tracks + presence tests"
    slide.placeholders[1].text = (
        f"Significance by Poisson q< {Q_CUTOFF} per comparison (3vs4/3vs5/4vs5)\n"
        "Density: GLM with log(offset = window length)\n"
        "Counts:  GLM without offset (raw counts per window)\n"
        "Presence tests: Fisher & Chi-square on windows with ≥1 TE vs 0 TE"
    )
    # Add one slide per comparison summarizing presence tests (overall TE)
    for comp in comps:
        row = fisher_df[(fisher_df["comparison"]==comp) & (fisher_df["feature"]=="TE_total")].iloc[0]
        bullets = [
            f"{comp}: windows sig/not = {int(row['n_windows_sig']):,} / {int(row['n_windows_not']):,}",
            f"Total TE events: sig={int(row['total_TE_sig']):,}, not={int(row['total_TE_not']):,}",
            f"Presence 2x2: [[{row['a_sig_ge1']},{row['b_sig_eq0']}],[{row['c_not_ge1']},{row['d_not_eq0']}]]",
            f"Fisher OR={row['odds_ratio_fisher']:.3f}, p={row['p_fisher']:.2e}",
            f"Chi-square χ²={row['chi2']:.2f} (df={int(row['df'])}), p={row['p_chi2']:.2e}"
        ]
        add_slide(prs, f"Presence tests (overall TE) — {comp}", bullets=bullets)

    # Manhattan overview slides
    add_slide(prs, "Manhattan (counts) — overall TE", image=str(manh_overall))
    # (Optional) add a slide that lists where per-class panels were written
    add_slide(prs, "Per-class Manhattan panels (counts)", bullets=[
        "Generated one 3-row panel per TE class (counts per window).",
        f"Files (examples): {Path(manh_teclass_paths[0]).name}, {Path(manh_teclass_paths[-1]).name}",
        f"Output folder: {OUTDIR.resolve()}"
    ])

    pptx_path = OUTDIR / "TE_vs_FST_report.pptx"
    prs.save(str(pptx_path))

    print(f"\nDone. Outputs in: {OUTDIR.resolve()}")
    print(f"- Augmented windows: fst_windows_with_TEcounts.tsv")
    print(f"- Presence tests TSV: {fisher_out.name}")
    print(f"- Manhattan (overall): {Path(manh_overall).name}")
    print(f"- {len(manh_teclass_paths)} per-class Manhattan panels saved.")
    print(f"- Report: {pptx_path.name}")

if __name__ == "__main__":
    main()
