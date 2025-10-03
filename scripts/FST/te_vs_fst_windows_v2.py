#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
TE overlap & density vs FST-enriched windows
— Separate results for 3vs4, 3vs5, 4vs5 + panel figures

Inputs (same folder as this script):
  fst_window_counts_tests_win10000000_ov2000000.csv
  (MUST include these exact columns, per user):
  CHROM,WIN_START,WIN_END,C0_3vs4,C0.25_3vs4,C0_3vs5,C0.25_3vs5,C0_4vs5,C0.25_4vs5,
  R_3vs4,R_3vs5,R_4vs5,p_glob_3vs4,lambda_3vs4,disp_bin_3vs4,disp_poi_3vs4,
  rho_bb_3vs4,theta_nb_3vs4,p_bin_3vs4,q_bin_3vs4,p_poi_3vs4,q_poi_3vs4,
  p_bb_3vs4,q_bb_3vs4,p_nb_3vs4,q_nb_3vs4,p_glob_3vs5,lambda_3vs5,disp_bin_3vs5,
  disp_poi_3vs5,rho_bb_3vs5,theta_nb_3vs5,p_bin_3vs5,q_bin_3vs5,p_poi_3vs5,
  q_poi_3vs5,p_bb_3vs5,q_bb_3vs5,p_nb_3vs5,q_nb_3vs5,p_glob_4vs5,lambda_4vs5,
  disp_bin_4vs5,disp_poi_4vs5,rho_bb_4vs5,theta_nb_4vs5,p_bin_4vs5,q_bin_4vs5,
  p_poi_4vs5,q_poi_4vs5,p_bb_4vs5,q_bb_4vs5,p_nb_4vs5,q_nb_4vs5,WIN_MID

TE TSV (outside this folder):
  ~/Desktop/OSU_projects/conifers/LP/lodgepole_pine_assembly/col8_readable.TEanno.cds.scaf1ab.tsv
  Columns (tab-separated): seqid, source, sequence_ontology, start, end, score, strand

Outputs (in ./te_vs_fst_out/):
  - fst_windows_with_TEcounts.tsv  (augmented table: TE counts & densities)
  - te_density_group_summary.tsv   (group-level stats & model summaries for all 3 comps)
  - te_density_by_type_<comp>.tsv  (per-TE-type model results with FDR per comparison)
  - figures:
      te_density_violin_box_PANEL.png   (3-row panel)
      te_type_rate_ratio_forest_<comp>.png  (per comp)
      te_type_rate_ratio_forest_PANEL.png   (compact 3-row forest)
  - PowerPoint: TE_vs_FST_report.pptx
"""

import os
import math
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

import statsmodels.api as sm
from statsmodels.api import GLM, families
from statsmodels.genmod.families import links
from statsmodels.stats.multitest import multipletests
from scipy.stats import norm
from pptx import Presentation
from pptx.util import Inches

# =========================
# User knobs (edit here)
# =========================
FST_CSV = "fst_window_counts_tests_win10000000_ov2000000.csv"
TE_TSV  = os.path.expanduser(
    "~/Desktop/OSU_projects/conifers/LP/lodgepole_pine_assembly/col8_readable.TEanno.cds.scaf1ab.tsv"
)

# Poisson q-value cutoff to define “significant FST windows” per comparison
Q_CUTOFF = 0.05

# Output folder
OUTDIR = Path("te_vs_fst_out")
OUTDIR.mkdir(parents=True, exist_ok=True)

np.random.seed(1)

# =========================
# IO & cleaning
# =========================
def read_fst_windows(path: str) -> pd.DataFrame:
    df = pd.read_csv(path)
    need = ["CHROM","WIN_START","WIN_END","WIN_MID",
            "q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]
    missing = [c for c in need if c not in df.columns]
    if missing:
        raise ValueError(f"FST file missing required column(s): {missing}")

    df["CHROM"]     = df["CHROM"].astype(str)
    df["WIN_START"] = pd.to_numeric(df["WIN_START"], errors="coerce").astype("Int64")
    df["WIN_END"]   = pd.to_numeric(df["WIN_END"],   errors="coerce").astype("Int64")
    df["WIN_MID"]   = pd.to_numeric(df["WIN_MID"],   errors="coerce").astype("Int64")
    df["WIN_LEN"]   = (df["WIN_END"] - df["WIN_START"] + 1).astype("Int64")
    return df

def _has_header(path: str, first_col_name: str) -> bool:
    with open(path, "r") as fh:
        first = fh.readline().strip().split("\t")[0]
        return first == first_col_name

def read_te_table(path: str) -> pd.DataFrame:
    cols = ["seqid","source","sequence_ontology","start","end","score","strand"]
    te = pd.read_csv(path, sep="\t", names=cols,
                     header=0 if _has_header(path, cols[0]) else None)
    te["seqid"] = te["seqid"].astype(str)
    te["sequence_ontology"] = te["sequence_ontology"].astype(str)
    te["start"] = pd.to_numeric(te["start"], errors="coerce").astype("Int64")
    te["end"]   = pd.to_numeric(te["end"],   errors="coerce").astype("Int64")
    te = te.dropna(subset=["seqid","start","end"]).copy()
    te = te.sort_values(["seqid","start","end"]).reset_index(drop=True)
    return te

# =========================
# Overlap counting
# =========================
def _count_overlaps_one_chrom(wins: pd.DataFrame, tes: pd.DataFrame, te_types: np.ndarray):
    w_st = wins["WIN_START"].to_numpy(dtype=np.int64)
    w_en = wins["WIN_END"].to_numpy(dtype=np.int64)
    nW   = len(wins)

    total_counts = np.zeros(nW, dtype=np.int64)
    type_to_idx = {t:i for i,t in enumerate(te_types)}
    per_type = np.zeros((nW, len(te_types)), dtype=np.int64)

    if tes.empty:
        return total_counts, per_type

    t_st = tes["start"].to_numpy(dtype=np.int64)
    t_en = tes["end"].to_numpy(dtype=np.int64)
    t_ty = tes["sequence_ontology"].astype(str).to_numpy()

    j = 0
    for i in range(nW):
        ws, we = w_st[i], w_en[i]
        while j < len(t_st) and t_en[j] < ws:
            j += 1
        k = j
        while k < len(t_st) and t_st[k] <= we:
            if not (t_en[k] < ws or t_st[k] > we):
                total_counts[i] += 1
                idx = type_to_idx.get(t_ty[k], None)
                if idx is not None:
                    per_type[i, idx] += 1
            k += 1
    return total_counts, per_type

def count_te_overlaps_per_window(fst_df: pd.DataFrame, te_df: pd.DataFrame):
    te_types = np.array(sorted(te_df["sequence_ontology"].astype(str).unique()))
    base = fst_df[["CHROM","WIN_START","WIN_END","WIN_LEN"]].copy().reset_index(drop=True)

    total_counts_all = np.zeros(len(base), dtype=np.int64)
    per_type_all = np.zeros((len(base), len(te_types)), dtype=np.int64)

    for chrom, wins in fst_df.groupby("CHROM", sort=False):
        rows = wins.index.to_numpy()
        wins_sorted = wins.sort_values("WIN_START")
        order_back = np.argsort(wins_sorted.index.to_numpy())
        tes = te_df[te_df["seqid"] == chrom].sort_values("start")

        total_counts, per_type = _count_overlaps_one_chrom(wins_sorted, tes, te_types)
        total_counts_all[rows] = total_counts[order_back]
        per_type_all[rows, :]  = per_type[order_back, :]

    te_cols = {"TE_count": total_counts_all.astype(np.int64)}
    for j, t in enumerate(te_types):
        te_cols[f"TE_{t}"] = per_type_all[:, j].astype(np.int64)

    te_df_out = pd.DataFrame(te_cols)
    out = pd.concat([base, te_df_out], axis=1)

    win_len = out["WIN_LEN"].astype(float)
    out["TE_density"] = out["TE_count"] / win_len
    for t in te_types:
        out[f"TEdens_{t}"] = out[f"TE_{t}"] / win_len

    return out, te_types

# =========================
# Modeling helpers
# =========================
def dispersion_ratio(y, mu, df_adj=1):
    y = np.asarray(y, float)
    mu = np.asarray(mu, float)
    ok = mu > 0
    if ok.sum() <= 1:
        return np.nan
    return float(np.sum(((y[ok]-mu[ok])**2)/mu[ok]) / max(ok.sum()-df_adj, 1))

def clamp_exp(x):
    return float(np.exp(np.clip(x, -700, 700)))

def fit_glm_poisson(y, X, offset):
    fam = families.Poisson(link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    res = model.fit()
    return res

def estimate_nb_alpha(y, mu):
    numer = np.sum((y - mu) ** 2 - mu)
    denom = np.sum(mu ** 2)
    if denom <= 0:
        return 1.0
    alpha = numer / denom
    return float(max(alpha, 1e-8))

def fit_glm_nb(y, X, offset, alpha):
    fam = families.NegativeBinomial(alpha=alpha, link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    res = model.fit()
    return res

def glm_compare_counts(y, group, offset_log, prefer_nb_if_overdispersed=True):
    X = pd.DataFrame({"intercept": 1.0, "group": group.astype(float)})

    # Poisson
    res_p = fit_glm_poisson(y, X, offset_log)
    mu_p  = res_p.fittedvalues
    disp  = dispersion_ratio(y, mu_p)

    res_used = res_p
    family_used = "Poisson"

    if prefer_nb_if_overdispersed and np.isfinite(disp) and disp > 1.5:
        alpha = estimate_nb_alpha(y, mu_p)
        res_nb = fit_glm_nb(y, X, offset_log, alpha)
        res_used = res_nb
        family_used = f"NegBin(alpha={alpha:.3g})"

    b = res_used.params.get("group", np.nan)
    se = res_used.bse.get("group", np.nan)
    z  = (b / se) if se and np.isfinite(se) and se > 0 else np.nan
    p  = 2 * (1 - norm.cdf(abs(z))) if np.isfinite(z) else np.nan
    rr = math.exp(b) if np.isfinite(b) else np.nan
    ci_lo = clamp_exp(b - 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan
    ci_hi = clamp_exp(b + 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan

    return {"family": family_used, "coef": b, "se": se, "z": z, "p": p,
            "RR": rr, "RR_CI_lo": ci_lo, "RR_CI_hi": ci_hi,
            "dispersion": disp, "res": res_used}

def global_rate_ratio_test(k_sig, e_sig, k_nsig, e_nsig):
    lam1 = k_sig / e_sig if e_sig > 0 else np.nan
    lam0 = k_nsig / e_nsig if e_nsig > 0 else np.nan
    rr = lam1/lam0 if (lam0 and lam1) else np.nan
    if k_sig > 0 and k_nsig > 0 and np.isfinite(rr):
        se = math.sqrt(1.0/k_sig + 1.0/k_nsig)
        z  = math.log(rr)/se
        p  = 2*(1 - norm.cdf(abs(z)))
        ci_lo = clamp_exp(math.log(rr) - 1.96*se)
        ci_hi = clamp_exp(math.log(rr) + 1.96*se)
    else:
        p = np.nan; ci_lo = np.nan; ci_hi = np.nan
    return rr, ci_lo, ci_hi, p

# =========================
# Plots (panels)
# =========================
def violin_box_panel(df_aug: pd.DataFrame, comp_masks: dict, outdir: Path):
    """
    3-row panel: each row is a comparison (3vs4, 3vs5, 4vs5).
    """
    fig, axes = plt.subplots(3, 1, figsize=(7.5, 9), sharey=True)
    for i, comp in enumerate(["3vs4","3vs5","4vs5"]):
        mask = comp_masks[comp]
        data = [df_aug.loc[~mask, "TE_density"], df_aug.loc[mask, "TE_density"]]
        ax = axes[i]
        ax.violinplot(data, positions=[0,1], showmeans=True, widths=0.9)
        ax.boxplot(data, positions=[0,1], widths=0.25, showcaps=True, showfliers=False)
        ax.set_xticks([0,1], ["not significant", "significant"])
        ax.set_ylabel("TE density (TEs / bp)")
        ax.set_title(f"TE density by FST significance — {comp} (q_poi < 0.05)")
        ax.grid(alpha=0.2, axis="y")
    axes[-1].set_xlabel("Window group")
    fig.tight_layout()
    out = outdir / "te_density_violin_box_PANEL.png"
    fig.savefig(out, dpi=220); plt.close(fig)
    return out

def forest_per_comp(df_types: pd.DataFrame, comp: str, outdir: Path):
    d = df_types.sort_values("RR", ascending=False).copy()
    y = np.arange(len(d))[::-1]
    fig, ax = plt.subplots(figsize=(7.5, 0.35*max(6, len(d))))
    ax.errorbar(d["RR"], y,
                xerr=[d["RR"]-d["RR_CI_lo"], d["RR_CI_hi"]-d["RR"]],
                fmt='o', capsize=3)
    ax.axvline(1.0, color="#999999", lw=1, ls="--")
    ax.set_yticks(y, d["TE_type"])
    ax.set_xscale("log")
    ax.set_xlabel("Rate ratio (significant / not)   [log scale]")
    ax.set_title(f"Per-TE-type enrichment in significant windows — {comp}")
    ax.grid(alpha=0.2, axis="x")
    out = outdir / f"te_type_rate_ratio_forest_{comp}.png"
    fig.tight_layout(); fig.savefig(out, dpi=220); plt.close(fig)
    return out

def compact_forest_panel(per_comp_types: dict, outdir: Path, max_types=12):
    """
    Compact 3-row panel; each row shows top-N TE types by significance (q) for that comparison.
    """
    fig, axes = plt.subplots(3, 1, figsize=(9, 9), sharex=False)
    for i, comp in enumerate(["3vs4","3vs5","4vs5"]):
        d = per_comp_types[comp].sort_values("q").head(max_types).copy()
        d = d.sort_values("RR", ascending=False)
        y = np.arange(len(d))[::-1]
        ax = axes[i]
        ax.errorbar(d["RR"], y, xerr=[d["RR"]-d["RR_CI_lo"], d["RR_CI_hi"]-d["RR"]],
                    fmt='o', capsize=3)
        ax.axvline(1.0, color="#999999", lw=1, ls="--")
        ax.set_yticks(y, d["TE_type"])
        ax.set_xscale("log")
        ax.set_xlabel("Rate ratio (sig/not)   [log scale]")
        ax.set_title(f"Top TE types (by FDR) — {comp}")
        ax.grid(alpha=0.2, axis="x")
    fig.tight_layout()
    out = outdir / "te_type_rate_ratio_forest_PANEL.png"
    fig.savefig(out, dpi=220); plt.close(fig)
    return out

# =========================
# Report
# =========================
def add_slide(prs, title, bullets=None, image=None):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if bullets:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.6))
        tf = tb.text_frame; tf.clear()
        for i,b in enumerate(bullets):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.text = b; p.level = 0
    if image and Path(image).exists():
        slide.shapes.add_picture(str(image), Inches(0.5), Inches(2.7), height=Inches(3.8))

# =========================
# Main
# =========================
def main():
    # Load data
    fst = read_fst_windows(FST_CSV)
    te  = read_te_table(TE_TSV)

    # TE overlaps & densities
    te_counts, te_types = count_te_overlaps_per_window(fst, te)
    df_aug = fst.merge(te_counts, on=["CHROM","WIN_START","WIN_END","WIN_LEN"], how="left")
    num_cols = df_aug.select_dtypes(include=[np.number]).columns
    df_aug[num_cols] = df_aug[num_cols].fillna(0)

    # Per-comparison analyses
    comps = ["3vs4","3vs5","4vs5"]
    comp_masks = {}           # {comp: boolean Series is_sig_comp}
    per_comp_overall = {}     # {comp: dict of overall stats}
    per_comp_types = {}       # {comp: DataFrame of per-type RR + FDR}
    per_comp_forest_png = {}  # {comp: path}

    for comp in comps:
        qcol = f"q_poi_{comp}"
        if qcol not in df_aug.columns:
            raise ValueError(f"Expected Poisson q-value column missing: {qcol}")

        # mask
        is_sig = (df_aug[qcol] < Q_CUTOFF)
        comp_masks[comp] = is_sig

        # overall model: TE_count ~ is_sig + offset(log(WIN_LEN))
        y = df_aug["TE_count"].astype(int).to_numpy()
        group = is_sig.astype(int).to_numpy()
        offset = np.log(df_aug["WIN_LEN"].astype(float).to_numpy())

        res_used = glm_compare_counts(y, group, offset)
        model_used = res_used["family"]

        # global rate-ratio
        k_sig   = int(df_aug.loc[is_sig,  "TE_count"].sum())
        e_sig   = float(df_aug.loc[is_sig,  "WIN_LEN"].sum())
        k_nsig  = int(df_aug.loc[~is_sig, "TE_count"].sum())
        e_nsig  = float(df_aug.loc[~is_sig, "WIN_LEN"].sum())
        rr_glob, ci_lo_glob, ci_hi_glob, p_glob = global_rate_ratio_test(k_sig, e_sig, k_nsig, e_nsig)

        # per-type models (BH-FDR per comparison)
        rows = []
        for t in te_types:
            y_t = df_aug[f"TE_{t}"].astype(int).to_numpy()
            res_t = glm_compare_counts(y_t, group, offset)
            rows.append({
                "TE_type": t,
                "family": res_t["family"],
                "RR": res_t["RR"], "RR_CI_lo": res_t["RR_CI_lo"], "RR_CI_hi": res_t["RR_CI_hi"],
                "p": res_t["p"], "dispersion": res_t["dispersion"]
            })
        df_types_comp = pd.DataFrame(rows)
        df_types_comp["q"] = multipletests(df_types_comp["p"].fillna(1.0), method="fdr_bh")[1]
        df_types_comp = df_types_comp.sort_values("q")

        # store per-comp results
        per_comp_overall[comp] = {
            "model_used": model_used,
            "RR_overall": res_used["RR"], "RR_CI_lo": res_used["RR_CI_lo"], "RR_CI_hi": res_used["RR_CI_hi"],
            "p_model": res_used["p"], "dispersion_est": res_used["dispersion"],
            "global_RR": rr_glob, "global_CI_lo": ci_lo_glob, "global_CI_hi": ci_hi_glob, "p_global": p_glob,
            "n_sig": int(is_sig.sum()), "n_not": int((~is_sig).sum())
        }
        per_comp_types[comp] = df_types_comp

        # write per-comp types table
        out_types = OUTDIR / f"te_density_by_type_{comp}.tsv"
        df_types_comp.to_csv(out_types, sep="\t", index=False)

        # per-comp forest plot
        per_comp_forest_png[comp] = forest_per_comp(df_types_comp, comp, OUTDIR)

    # Write the augmented windows (once)
    out_aug = OUTDIR / "fst_windows_with_TEcounts.tsv"
    df_aug.to_csv(out_aug, sep="\t", index=False)

    # Build a combined group summary table for all comparisons
    rows = []
    for comp in comps:
        is_sig = comp_masks[comp]
        sub_sig  = df_aug[is_sig]
        sub_nsig = df_aug[~is_sig]
        ov = per_comp_overall[comp]
        rows.append({
            "comparison": comp,
            "n_windows": len(df_aug),
            "n_sig": ov["n_sig"], "n_not": ov["n_not"],
            "mean_density_sig": float((sub_sig["TE_count"]/sub_sig["WIN_LEN"]).mean()) if len(sub_sig)>0 else np.nan,
            "mean_density_not": float((sub_nsig["TE_count"]/sub_nsig["WIN_LEN"]).mean()) if len(sub_nsig)>0 else np.nan,
            "model_used": ov["model_used"],
            "RR_overall": ov["RR_overall"], "RR_CI_lo": ov["RR_CI_lo"], "RR_CI_hi": ov["RR_CI_hi"], "p_model": ov["p_model"],
            "dispersion_est": ov["dispersion_est"],
            "global_RR": ov["global_RR"], "global_CI_lo": ov["global_CI_lo"], "global_CI_hi": ov["global_CI_hi"], "p_global": ov["p_global"]
        })
    df_summary_all = pd.DataFrame(rows)
    out_summ = OUTDIR / "te_density_group_summary.tsv"
    df_summary_all.to_csv(out_summ, sep="\t", index=False)

    # Panel plots
    fig_violin_panel = violin_box_panel(df_aug, comp_masks, OUTDIR)
    fig_forest_panel = compact_forest_panel(per_comp_types, OUTDIR, max_types=12)

    # PPTX
    prs = Presentation()
    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TE density vs FST-enriched windows (per comparison)"
    slide.placeholders[1].text = (
        f"Grouping by Poisson q< {Q_CUTOFF} for each of: 3vs4, 3vs5, 4vs5\n"
        "Counts modeled by GLM with log(offset=window length)\n"
        "Auto-switch to Negative Binomial under overdispersion\n"
        "Per-TE-type effects FDR-controlled within each comparison"
    )

    # One slide per comparison (bullets)
    for comp in comps:
        ov = per_comp_overall[comp]
        bullets = [
            f"{comp}: windows sig/not = {ov['n_sig']:,} / {ov['n_not']:,}",
            f"Overall RR (sig/not) = {ov['RR_overall']:.3f} "
            f"[{ov['RR_CI_lo']:.3f}, {ov['RR_CI_hi']:.3f}], p={ov['p_model']:.2e}",
            f"Model used: {ov['model_used']} (Poisson dispersion≈{ov['dispersion_est']:.2f})",
            f"Global RR = {ov['global_RR']:.3f} "
            f"[{ov['global_CI_lo']:.3f}, {ov['global_CI_hi']:.3f}], p={ov['p_global']:.2e}",
            ("Top TE types by FDR (q<0.05): " +
             ", ".join(per_comp_types[comp].loc[per_comp_types[comp]["q"]<0.05, "TE_type"]
                       .head(10).astype(str)))
                if (per_comp_types[comp]["q"]<0.05).any()
                else "No TE types pass FDR<0.05"
        ]
        add_slide(prs, f"Summary — {comp}", bullets=bullets, image=str(per_comp_forest_png[comp]))

    # Panel figures slides
    add_slide(prs, "TE density violin+box — PANEL", bullets=None, image=str(fig_violin_panel))
    add_slide(prs, "Per-TE-type RR (top by FDR) — PANEL", bullets=None, image=str(fig_forest_panel))

    pptx_path = OUTDIR / "TE_vs_FST_report.pptx"
    prs.save(str(pptx_path))

    print(f"\nDone. Outputs in: {OUTDIR.resolve()}")
    print(f"- Augmented windows: {out_aug.name}")
    print(f"- Group summary:     {out_summ.name}")
    for comp in comps:
        print(f"- Per-type ({comp}): te_density_by_type_{comp}.tsv, figure: {per_comp_forest_png[comp].name}")
    print(f"- Panel figures:     {Path(fig_violin_panel).name}, {Path(fig_forest_panel).name}")
    print(f"- Report:            {pptx_path.name}")

if __name__ == "__main__":
    main()

