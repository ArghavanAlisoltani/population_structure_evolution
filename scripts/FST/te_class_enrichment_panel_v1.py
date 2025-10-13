#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
TE class enrichment tied to significant FST windows
(uses your exact v9 headers; adds Gypsy↔Copia, Helitron↔Copia, Helitron↔Gypsy)

Inputs (TSV; tab-separated) must contain (from your header):
  CHROM, WIN_START, WIN_END, WIN_MID, WIN_LEN,
  q_poi_3vs4, q_poi_3vs5, q_poi_4vs5,
  TE_Gypsy_LTR_retrotransposon,   TEdens_Gypsy_LTR_retrotransposon,
  TE_Copia_LTR_retrotransposon,   TEdens_Copia_LTR_retrotransposposon,
  TE_helitron,                    TEdens_helitron
(and many more columns we don't touch)

Outputs:
  - augmented.tsv (adds is_sig_* flags + sanity)
  - fisher_presence.tsv (Gypsy/Copia/Helitron presence in sig vs non; per pair)
  - glm_interaction.tsv (Poisson & NB; class x significance interaction; per pair)
  - glm_within_sig.tsv (class-vs-class within sig-only; Poisson & NB; per pair)
  - class_props.tsv (proportions for plotting)
  - figures/*.png
  - slides/TE_enrichment_panel.pptx  (optional)
"""

import os
import math
import argparse
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

from scipy.stats import fisher_exact
import statsmodels.api as sm
from statsmodels.genmod.families import Poisson, NegativeBinomial
from statsmodels.genmod.families.links import log as LogLink

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAVE_PPTX = True
except Exception:
    HAVE_PPTX = False

PAIRS = ["3vs4","3vs5","4vs5"]
CLASSES = {
    "Gypsy":    {"count":"TE_Gypsy_LTR_retrotransposon",  "dens":"TEdens_Gypsy_LTR_retrotransposon"},
    "Copia":    {"count":"TE_Copia_LTR_retrotransposon",  "dens":"TEdens_Copia_LTR_retrotransposon"},
    "Helitron": {"count":"TE_helitron",                   "dens":"TEdens_helitron"},
}

CMP_SETS = [
    ("Gypsy","Copia"),
    ("Helitron","Copia"),
    ("Helitron","Gypsy"),
]

def parse_args():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="infile", required=True, help="v9 TSV: fst_windows_with_TEcounts.tsv")
    ap.add_argument("--outdir", default="te_class_panel_out", help="Output directory")
    ap.add_argument("--qcut", type=float, default=0.05, help="Poisson q-value cutoff for FST significance")
    ap.add_argument("--make_plots", action="store_true", help="Emit PNGs")
    ap.add_argument("--make_slides", action="store_true", help="Make PPTX (python-pptx required)")
    return ap.parse_args()

def ensure_columns(df):
    base = ["CHROM","WIN_START","WIN_END","WIN_MID","WIN_LEN",
            "q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]
    for k,v in CLASSES.items():
        base.append(v["count"])
    missing = [c for c in base if c not in df.columns]
    if missing:
        raise ValueError(f"Missing required columns: {missing}")

def add_sig_flags(df, qcut):
    df = df.copy()
    for pair in PAIRS:
        df[f"is_sig_{pair}"] = (pd.to_numeric(df[f"q_poi_{pair}"], errors="coerce") < qcut).astype(int)
    return df

def fisher_presence(sig_mask, has_feat):
    a = int((sig_mask & has_feat).sum())
    b = int((sig_mask & ~has_feat).sum())
    c = int((~sig_mask & has_feat).sum())
    d = int((~sig_mask & ~has_feat).sum())
    if (a+b==0) or (c+d==0):
        odds, p = (np.nan, np.nan)
    else:
        odds, p = fisher_exact([[a,b],[c,d]], alternative="two-sided")
    return a,b,c,d,odds,p

def glm_interaction_counts(df, pair, clsA, clsB, family="poisson", qcut=0.05):
    """Count model with offset: count ~ te_type(A vs B) + is_sig + interaction; offset=log(WIN_LEN)"""
    d = df.copy()
    length = pd.to_numeric(d["WIN_LEN"], errors="coerce").fillna((d["WIN_END"]-d["WIN_START"]+1)).clip(lower=1).astype(float)
    sig = (pd.to_numeric(d[f"q_poi_{pair}"], errors="coerce") < qcut).astype(int)

    yA = pd.to_numeric(d[CLASSES[clsA]["count"]], errors="coerce").fillna(0).astype(float).values
    yB = pd.to_numeric(d[CLASSES[clsB]["count"]], errors="coerce").fillna(0).astype(float).values

    y = np.concatenate([yA, yB]).astype(float)
    te_is_A = np.concatenate([np.ones(len(d),dtype=int), np.zeros(len(d),dtype=int)])
    is_sig  = np.concatenate([sig.values, sig.values]).astype(int)
    off     = np.log(np.concatenate([length.values, length.values]).clip(1.0))

    X = pd.DataFrame({
        "intercept": 1.0,
        "te_is_A": te_is_A,            # B is reference
        "is_sig":  is_sig,
        "int_te_sig": te_is_A * is_sig
    })

    fam = Poisson(LogLink()) if family=="poisson" else NegativeBinomial(LogLink())
    fit = sm.GLM(y, X, family=fam, offset=off).fit()

    # dispersion
    mu = fit.fittedvalues
    var = mu if family=="poisson" else fit.family.variance(mu)
    ok = var > 0
    dfree = max(1, ok.sum() - X.shape[1])
    dispersion = float(np.sum(((y - mu)[ok]**2)/var[ok]) / dfree)

    # interaction (does A−B difference change between sig vs non?)
    b = fit.params["int_te_sig"]; se = fit.bse["int_te_sig"]
    if se > 0 and np.isfinite(b):
        z = b/se
        p = 2*(1 - 0.5*(1+math.erf(abs(z)/math.sqrt(2))))
    else:
        z = np.nan; p = np.nan
    rr_int = float(np.exp(b)) if np.isfinite(b) else np.nan

    # main effects (optional)
    b_AB = fit.params["te_is_A"]; se_AB = fit.bse["te_is_A"]
    p_AB = 2*(1 - 0.5*(1+math.erf(abs(b_AB/se_AB)/math.sqrt(2)))) if se_AB>0 else np.nan
    b_sig = fit.params["is_sig"]; se_sig = fit.bse["is_sig"]
    p_sig = 2*(1 - 0.5*(1+math.erf(abs(b_sig/se_sig)/math.sqrt(2)))) if se_sig>0 else np.nan

    return {
        "dispersion": dispersion,
        "beta_A_vs_B": float(b_AB), "p_A_vs_B": float(p_AB),
        "beta_is_sig": float(b_sig), "p_is_sig": float(p_sig),
        "beta_int": float(b), "se_int": float(se), "z_int": float(z),
        "p_int": float(p), "RR_int": rr_int
    }

def glm_within(df, pair, clsA, clsB, family="poisson", qcut=0.05, which="sig"):
    """Class-vs-class **within** sig-only (which='sig') or non-sig (which='non') windows."""
    sel = (pd.to_numeric(df[f"q_poi_{pair}"], errors="coerce") < qcut) if which=="sig" \
          else ~(pd.to_numeric(df[f"q_poi_{pair}"], errors="coerce") < qcut)
    d = df.loc[sel].copy()
    length = pd.to_numeric(d["WIN_LEN"], errors="coerce").fillna((d["WIN_END"]-d["WIN_START"]+1)).clip(lower=1).astype(float)

    yA = pd.to_numeric(d[CLASSES[clsA]["count"]], errors="coerce").fillna(0).astype(float).values
    yB = pd.to_numeric(d[CLASSES[clsB]["count"]], errors="coerce").fillna(0).astype(float).values

    y = np.concatenate([yA, yB]).astype(float)
    te_is_A = np.concatenate([np.ones(len(d),dtype=int), np.zeros(len(d),dtype=int)])
    off     = np.log(np.concatenate([length.values, length.values]).clip(1.0))
    X = pd.DataFrame({"intercept":1.0, "te_is_A": te_is_A})

    fam = Poisson(LogLink()) if family=="poisson" else NegativeBinomial(LogLink())
    fit = sm.GLM(y, X, family=fam, offset=off).fit()

    mu = fit.fittedvalues
    var = mu if family=="poisson" else fit.family.variance(mu)
    ok = var > 0
    dfree = max(1, ok.sum() - X.shape[1])
    dispersion = float(np.sum(((y - mu)[ok]**2)/var[ok]) / dfree)

    b = fit.params["te_is_A"]; se = fit.bse["te_is_A"]
    if se > 0 and np.isfinite(b):
        z = b/se
        p = 2*(1 - 0.5*(1+math.erf(abs(z)/math.sqrt(2))))
    else:
        z = np.nan; p = np.nan
    rr = float(np.exp(b)) if np.isfinite(b) else np.nan

    return {
        "n_windows": int(len(d)),
        "dispersion": dispersion,
        "beta_A_vs_B": float(b), "se": float(se), "z": float(z), "p": float(p), "RR": rr,
        "which": which
    }

def summarize_props(df, qcut):
    """Proportion present (any>0) per class × pair × sig/non (for plots)."""
    rows = []
    for pair in PAIRS:
        sig = (pd.to_numeric(df[f"q_poi_{pair}"], errors="coerce") < qcut)
        for cls, meta in CLASSES.items():
            present = (pd.to_numeric(df[meta["count"]], errors="coerce").fillna(0) > 0)
            a = int((sig & present).sum()); b = int((sig & ~present).sum())
            c = int((~sig & present).sum()); d = int((~sig & ~present).sum())
            n_sig = a + b; n_non = c + d
            prop_sig = a / n_sig if n_sig>0 else np.nan
            prop_non = c / n_non if n_non>0 else np.nan
            rows.append({"pair":pair,"class":cls,"n_sig":n_sig,"n_non":n_non,
                         "prop_present_sig":prop_sig,"prop_present_non":prop_non,
                         "a":a,"b":b,"c":c,"d":d})
    return pd.DataFrame(rows)

def plot_props(df_props, outdir):
    for pair in PAIRS:
        sub = df_props[df_props["pair"]==pair].copy()
        labels = sub["class"].tolist()
        sigv = sub["prop_present_sig"].values
        nonv = sub["prop_present_non"].values
        x = np.arange(len(labels))
        w = 0.35
        fig, ax = plt.subplots(figsize=(8,4))
        ax.bar(x - w/2, nonv, width=w, label="Not-sig")
        ax.bar(x + w/2, sigv, width=w, label="Sig")
        ax.set_xticks(x); ax.set_xticklabels(labels)
        ax.set_ylim(0,1)
        ax.set_ylabel("Proportion of windows with ≥1 TE")
        ax.set_title(f"Presence proportions by class — {pair}")
        ax.legend(frameon=False)
        fig.tight_layout()
        fig.savefig(os.path.join(outdir, f"prop_presence_{pair}.png"), dpi=300)
        plt.close(fig)

def plot_violin_counts(df, pair, outdir, qcut):
    sig = (pd.to_numeric(df[f"q_poi_{pair}"], errors="coerce") < qcut)
    fig, axes = plt.subplots(1, 3, figsize=(12,3.5), sharey=True)
    for ax, cls in zip(axes, ["Gypsy","Copia","Helitron"]):
        cnt = pd.to_numeric(df[CLASSES[cls]["count"]], errors="coerce").fillna(0)
        data = [cnt[~sig].values, cnt[sig].values]
        ax.violinplot(data, showmedians=True)
        ax.set_title(cls)
        ax.set_xticks([1,2]); ax.set_xticklabels(["Not-sig","Sig"])
    axes[0].set_ylabel("TE count per window")
    fig.suptitle(f"Counts by significance — {pair}")
    fig.tight_layout(rect=[0,0,1,0.95])
    fig.savefig(os.path.join(outdir, f"counts_violin_{pair}.png"), dpi=300)
    plt.close(fig)

def make_slides(pngs, tables, out_pptx):
    if not HAVE_PPTX:
        print("[!] python-pptx not installed; skipping slides.")
        return
    prs = Presentation()
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = "TE class enrichment vs FST significance"
    title.placeholders[1].text = "Gypsy↔Copia, Helitron↔Copia, Helitron↔Gypsy\nPoisson & NB with offsets; Fisher presence"

    for label, path in pngs:
        if not os.path.exists(path): continue
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = label
        slide.shapes.add_picture(path, Inches(0.5), Inches(1.5), width=Inches(9.0))

    for label, path in tables:
        if not os.path.exists(path): continue
        df = pd.read_csv(path, sep="\t")
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = label
        # write a small top-rows text box
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(0.8))
        tb.text_frame.text = f"{os.path.basename(path)} (first 10 rows):"
        # dump first 10 rows to a simple preformatted text box
        pre = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9.0), Inches(4.5))
        pre.text_frame.text = df.head(10).to_string(index=False)

    prs.save(out_pptx)

if __name__ == "__main__":
    args = parse_args()
    INFILE  = args.infile
    OUTDIR  = args.outdir
    Q_CUT   = args.qcut
    MAKEP   = args.make_plots
    MAKES   = args.make_slides

    os.makedirs(OUTDIR, exist_ok=True)
    figdir = os.path.join(OUTDIR, "figures")
    os.makedirs(figdir, exist_ok=True)
    slidedir = os.path.join(OUTDIR, "slides")
    os.makedirs(slidedir, exist_ok=True)

    df = pd.read_csv(INFILE, sep="\t")
    ensure_columns(df)
    df = add_sig_flags(df, Q_CUT)

    # Save augmented table
    out_aug = os.path.join(OUTDIR, "augmented.tsv")
    df.to_csv(out_aug, sep="\t", index=False)

    # Fisher presence per class and pair
    fisher_rows = []
    diag_rows   = []  # raw 2x2
    for pair in PAIRS:
        sig = (pd.to_numeric(df[f"q_poi_{pair}"], errors="coerce") < Q_CUT)
        for cls, meta in CLASSES.items():
            present = (pd.to_numeric(df[meta["count"]], errors="coerce").fillna(0) > 0)
            a,b,c,d,odds,p = fisher_presence(sig, present)
            fisher_rows.append({"pair":pair,"class":cls,"a_sig_present":a,"b_sig_absent":b,
                                "c_non_present":c,"d_non_absent":d,"odds":odds,"p":p})
            diag_rows.append({"pair":pair,"class":cls,"sig_present":a,"sig_absent":b,
                              "non_present":c,"non_absent":d})
    fisher_df = pd.DataFrame(fisher_rows)
    fisher_df["q_bh_perpair"] = (fisher_df.groupby("pair")["p"]
                                 .transform(lambda p: (p*p.notna().sum()/p.rank(method="min")).clip(upper=1.0)))
    fisher_df.to_csv(os.path.join(OUTDIR,"fisher_presence.tsv"), sep="\t", index=False)
    pd.DataFrame(diag_rows).to_csv(os.path.join(OUTDIR,"fisher_presence_raw_counts.tsv"), sep="\t", index=False)

    # Class-vs-class interaction tests (Poisson & NB)
    glm_rows = []
    for pair in PAIRS:
        for (A,B) in CMP_SETS:
            for fam in ["poisson","negbin"]:
                res = glm_interaction_counts(df, pair, A, B, family=fam, qcut=Q_CUT)
                res.update({"pair":pair,"classA":A,"classB":B,"family":fam})
                glm_rows.append(res)
    glm_df = pd.DataFrame(glm_rows)
    # FDR across 3 pairs per comparison/family
    glm_df["comp"] = glm_df["classA"] + "_vs_" + glm_df["classB"] + "_" + glm_df["family"]
    glm_df["q_int_bh"] = (glm_df.groupby("comp")["p_int"]
                          .transform(lambda p: (p*p.notna().sum()/p.rank(method="min")).clip(upper=1.0)))
    glm_df.to_csv(os.path.join(OUTDIR,"glm_interaction.tsv"), sep="\t", index=False)

    # Within sig-only and non-sig-only (direct A vs B)
    wit_rows = []
    for pair in PAIRS:
        for (A,B) in CMP_SETS:
            for fam in ["poisson","negbin"]:
                for which in ["sig","non"]:
                    res = glm_within(df, pair, A, B, family=fam, qcut=Q_CUT, which=which)
                    res.update({"pair":pair,"classA":A,"classB":B,"family":fam})
                    wit_rows.append(res)
    wit_df = pd.DataFrame(wit_rows)
    wit_df["comp"] = wit_df["classA"] + "_vs_" + wit_df["classB"] + "_" + wit_df["family"] + "_" + wit_df["which"]
    wit_df["q_bh"] = (wit_df.groupby("comp")["p"]
                      .transform(lambda p: (p*p.notna().sum()/p.rank(method="min")).clip(upper=1.0)))
    wit_df.to_csv(os.path.join(OUTDIR,"glm_within_sig.tsv"), sep="\t", index=False)

    # Proportion bars + violins
    props = summarize_props(df, Q_CUT)
    props.to_csv(os.path.join(OUTDIR,"class_props.tsv"), sep="\t", index=False)
    if MAKEP:
        plot_props(props, figdir)
        for pair in PAIRS:
            plot_violin_counts(df, pair, figdir, Q_CUT)

    # Slides (optional)
    if MAKES:
        pngs = []
        for pair in PAIRS:
            pngs.append((f"Presence proportions — {pair}", os.path.join(figdir, f"prop_presence_{pair}.png")))
            pngs.append((f"Counts violin — {pair}",       os.path.join(figdir, f"counts_violin_{pair}.png")))
        tables = [
            ("Fisher presence (top) and raw counts (bottom)",
             os.path.join(OUTDIR,"fisher_presence.tsv")),
            ("GLM interaction", os.path.join(OUTDIR,"glm_interaction.tsv")),
            ("Within sig/non GLM", os.path.join(OUTDIR,"glm_within_sig.tsv")),
        ]
        make_slides(pngs, tables, os.path.join(slidedir,"TE_enrichment_panel.pptx"))

    print(f"[✓] Wrote to: {OUTDIR}")

