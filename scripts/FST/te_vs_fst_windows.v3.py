#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
TE overlap & density vs FST-enriched windows
— Separate results for 3vs4, 3vs5, 4vs5 + panel figures
— Adds *count-based* (no offset) comparisons alongside density-based analyses

Inputs (same folder as this script):
  fst_window_counts_tests_win10000000_ov2000000.csv
  (MUST include these exact columns, per user):
  CHROM,WIN_START,WIN_END,C0_3vs4,C0.25_3vs4,C0_3vs5,C0.25_3vs5,C0_4vs5,C0.25_4vs5,
  R_3vs4,R_3vs5,R_4vs5,p_glob_3vs4,lambda_3vs4,disp_bin_3vs4,disp_poi_3vs4,
  rho_bb_3vs4,theta_nb_3vs4,p_bin_3vs4,q_bin_3vs4,p_poi_3vs4,q_poi_3vs4,
  p_bb_3vs4,q_bb_3vs4,p_nb_3vs4,q_nb_3vs4,p_glob_3vs5,lambda_3vs5,disp_bin_3vs5,
  disp_poi_3vs5,rho_bb_3vs5,theta_nb_3vs5,p_bin_3vs5,q_bin_3vs5,p_poi_3vs5,
  q_poi_3vs5,p_bb_3vs5,q_bb_3vs5,p_nb_3vs5,q_nb_3vs5,p_glob_4vs5,lambda_4vs5,
  disp_bin_4vs5,disp_poi_4vs5,rho_bb_4vs5,theta_nb_4vs5,p_bin_4vs5,q_bin_4vs5,
  p_poi_4vs5,q_poi_4vs5,p_bb_4vs5,q_bb_4vs5,p_nb_4vs5,q_nb_4vs5,WIN_MID

TE TSV (outside this folder):
  ~/Desktop/OSU_projects/conifers/LP/lodgepole_pine_assembly/col8_readable.TEanno.cds.scaf1ab.tsv
  Columns (tab-separated): seqid, source, sequence_ontology, start, end, score, strand

Outputs (in ./te_vs_fst_out/):
  - fst_windows_with_TEcounts.tsv
  - te_density_group_summary.tsv            (density track, all comps)
  - te_density_by_type_<comp>.tsv           (density track, per TE type)
  - te_count_group_summary.tsv              (count track, all comps)
  - te_count_by_type_<comp>.tsv             (count track, per TE type)
  - figures:
      te_density_violin_box_PANEL.png
      te_type_rate_ratio_forest_<comp>.png
      te_type_rate_ratio_forest_PANEL.png
      te_count_violin_box_PANEL.png
      te_type_count_rate_ratio_forest_<comp>.png
      te_type_count_rate_ratio_forest_PANEL.png
  - PowerPoint: TE_vs_FST_report.pptx
"""
import os
import math
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

import statsmodels.api as sm
from statsmodels.api import GLM, families
from statsmodels.genmod.families import links
from statsmodels.stats.multitest import multipletests
from scipy.stats import norm
from pptx import Presentation
from pptx.util import Inches

# =========================
# User knobs (edit here)
# =========================
FST_CSV = "fst_window_counts_tests_win10000000_ov2000000.csv"
TE_TSV  = os.path.expanduser(
    "~/Desktop/OSU_projects/conifers/LP/lodgepole_pine_assembly/col8_readable.TEanno.cds.scaf1ab.tsv"
)

# Poisson q-value cutoff to define “significant FST windows” per comparison
Q_CUTOFF = 0.05

# Output folder
OUTDIR = Path("te_vs_fst_out")
OUTDIR.mkdir(parents=True, exist_ok=True)

np.random.seed(1)

# =========================
# IO & cleaning
# =========================
def read_fst_windows(path: str) -> pd.DataFrame:
    df = pd.read_csv(path)
    need = ["CHROM","WIN_START","WIN_END","WIN_MID",
            "q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]
    missing = [c for c in need if c not in df.columns]
    if missing:
        raise ValueError(f"FST file missing required column(s): {missing}")

    df["CHROM"]     = df["CHROM"].astype(str)
    df["WIN_START"] = pd.to_numeric(df["WIN_START"], errors="coerce").astype("Int64")
    df["WIN_END"]   = pd.to_numeric(df["WIN_END"],   errors="coerce").astype("Int64")
    df["WIN_MID"]   = pd.to_numeric(df["WIN_MID"],   errors="coerce").astype("Int64")
    df["WIN_LEN"]   = (df["WIN_END"] - df["WIN_START"] + 1).astype("Int64")
    return df

def _has_header(path: str, first_col_name: str) -> bool:
    with open(path, "r") as fh:
        first = fh.readline().strip().split("\t")[0]
        return first == first_col_name

def read_te_table(path: str) -> pd.DataFrame:
    cols = ["seqid","source","sequence_ontology","start","end","score","strand"]
    te = pd.read_csv(path, sep="\t", names=cols,
                     header=0 if _has_header(path, cols[0]) else None)
    te["seqid"] = te["seqid"].astype(str)
    te["sequence_ontology"] = te["sequence_ontology"].astype(str)
    te["start"] = pd.to_numeric(te["start"], errors="coerce").astype("Int64")
    te["end"]   = pd.to_numeric(te["end"],   errors="coerce").astype("Int64")
    te = te.dropna(subset=["seqid","start","end"]).copy()
    te = te.sort_values(["seqid","start","end"]).reset_index(drop=True)
    return te

# =========================
# Overlap counting
# =========================
def _count_overlaps_one_chrom(wins: pd.DataFrame, tes: pd.DataFrame, te_types: np.ndarray):
    w_st = wins["WIN_START"].to_numpy(dtype=np.int64)
    w_en = wins["WIN_END"].to_numpy(dtype=np.int64)
    nW   = len(wins)

    total_counts = np.zeros(nW, dtype=np.int64)
    type_to_idx = {t:i for i,t in enumerate(te_types)}
    per_type = np.zeros((nW, len(te_types)), dtype=np.int64)

    if tes.empty:
        return total_counts, per_type

    t_st = tes["start"].to_numpy(dtype=np.int64)
    t_en = tes["end"].to_numpy(dtype=np.int64)
    t_ty = tes["sequence_ontology"].astype(str).to_numpy()

    j = 0
    for i in range(nW):
        ws, we = w_st[i], w_en[i]
        while j < len(t_st) and t_en[j] < ws:
            j += 1
        k = j
        while k < len(t_st) and t_st[k] <= we:
            if not (t_en[k] < ws or t_st[k] > we):
                total_counts[i] += 1
                idx = type_to_idx.get(t_ty[k], None)
                if idx is not None:
                    per_type[i, idx] += 1
            k += 1
    return total_counts, per_type

def count_te_overlaps_per_window(fst_df: pd.DataFrame, te_df: pd.DataFrame):
    te_types = np.array(sorted(te_df["sequence_ontology"].astype(str).unique()))
    base = fst_df[["CHROM","WIN_START","WIN_END","WIN_LEN"]].copy().reset_index(drop=True)

    total_counts_all = np.zeros(len(base), dtype=np.int64)
    per_type_all = np.zeros((len(base), len(te_types)), dtype=np.int64)

    for chrom, wins in fst_df.groupby("CHROM", sort=False):
        rows = wins.index.to_numpy()
        wins_sorted = wins.sort_values("WIN_START")
        order_back = np.argsort(wins_sorted.index.to_numpy())
        tes = te_df[te_df["seqid"] == chrom].sort_values("start")

        total_counts, per_type = _count_overlaps_one_chrom(wins_sorted, tes, te_types)
        total_counts_all[rows] = total_counts[order_back]
        per_type_all[rows, :]  = per_type[order_back, :]

    te_cols = {"TE_count": total_counts_all.astype(np.int64)}
    for j, t in enumerate(te_types):
        te_cols[f"TE_{t}"] = per_type_all[:, j].astype(np.int64)

    te_df_out = pd.DataFrame(te_cols)
    out = pd.concat([base, te_df_out], axis=1)

    win_len = out["WIN_LEN"].astype(float)
    out["TE_density"] = out["TE_count"] / win_len
    for t in te_types:
        out[f"TEdens_{t}"] = out[f"TE_{t}"] / win_len

    return out, te_types

# =========================
# Modeling helpers
# =========================
def dispersion_ratio(y, mu, df_adj=1):
    y = np.asarray(y, float)
    mu = np.asarray(mu, float)
    ok = mu > 0
    if ok.sum() <= 1:
        return np.nan
    return float(np.sum(((y[ok]-mu[ok])**2)/mu[ok]) / max(ok.sum()-df_adj, 1))

def clamp_exp(x):
    return float(np.exp(np.clip(x, -700, 700)))

def fit_glm_poisson(y, X, offset=None):
    fam = families.Poisson(link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    res = model.fit()
    return res

def estimate_nb_alpha(y, mu):
    numer = np.sum((y - mu) ** 2 - mu)
    denom = np.sum(mu ** 2)
    if denom <= 0:
        return 1.0
    alpha = numer / denom
    return float(max(alpha, 1e-8))

def fit_glm_nb(y, X, offset=None, alpha=1.0):
    fam = families.NegativeBinomial(alpha=alpha, link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    res = model.fit()
    return res

def glm_compare_counts(y, group, offset_log=None, prefer_nb_if_overdispersed=True):
    """
    GLM: count ~ group (0/1), optional offset = log(exposure)
    Returns dict(coef,se,z,p,RR,CI,dispersion,family,res).
    """
    X = pd.DataFrame({"intercept": 1.0, "group": group.astype(float)})

    # Poisson
    res_p = fit_glm_poisson(y, X, offset=offset_log)
    mu_p  = res_p.fittedvalues
    disp  = dispersion_ratio(y, mu_p)

    res_used = res_p
    family_used = "Poisson"

    if prefer_nb_if_overdispersed and np.isfinite(disp) and disp > 1.5:
        alpha = estimate_nb_alpha(y, mu_p)
        res_nb = fit_glm_nb(y, X, offset=offset_log, alpha=alpha)
        res_used = res_nb
        family_used = f"NegBin(alpha={alpha:.3g})"

    b  = res_used.params.get("group", np.nan)
    se = res_used.bse.get("group", np.nan)
    z  = (b / se) if se and np.isfinite(se) and se > 0 else np.nan
    p  = 2 * (1 - norm.cdf(abs(z))) if np.isfinite(z) else np.nan
    rr = math.exp(b) if np.isfinite(b) else np.nan
    ci_lo = clamp_exp(b - 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan
    ci_hi = clamp_exp(b + 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan

    return {"family": family_used, "coef": b, "se": se, "z": z, "p": p,
            "RR": rr, "RR_CI_lo": ci_lo, "RR_CI_hi": ci_hi,
            "dispersion": disp, "res": res_used}

def global_rate_ratio_test(k_sig, e_sig, k_nsig, e_nsig):
    lam1 = k_sig / e_sig if e_sig > 0 else np.nan
    lam0 = k_nsig / e_nsig if e_nsig > 0 else np.nan
    rr = lam1/lam0 if (lam0 and lam1) else np.nan
    if k_sig > 0 and k_nsig > 0 and np.isfinite(rr):
        se = math.sqrt(1.0/k_sig + 1.0/k_nsig)
        z  = math.log(rr)/se
        p  = 2*(1 - norm.cdf(abs(z)))
        ci_lo = clamp_exp(math.log(rr) - 1.96*se)
        ci_hi = clamp_exp(math.log(rr) + 1.96*se)
    else:
        p = np.nan; ci_lo = np.nan; ci_hi = np.nan
    return rr, ci_lo, ci_hi, p

# =========================
# Plots (panels)
# =========================
def violin_box_panel(df_aug: pd.DataFrame, comp_masks: dict, outdir: Path, value_col="TE_density", tag="density"):
    """
    3-row panel: each row is a comparison (3vs4, 3vs5, 4vs5).
    value_col: "TE_density" (density track) or "TE_count" (count track)
    tag: used in output filename
    """
    fig, axes = plt.subplots(3, 1, figsize=(7.5, 9), sharey=(value_col=="TE_density"))
    for i, comp in enumerate(["3vs4","3vs5","4vs5"]):
        mask = comp_masks[comp]
        data = [df_aug.loc[~mask, value_col], df_aug.loc[mask, value_col]]
        ax = axes[i]
        ax.violinplot(data, positions=[0,1], showmeans=True, widths=0.9)
        ax.boxplot(data, positions=[0,1], widths=0.25, showcaps=True, showfliers=False)
        ax.set_xticks([0,1], ["not significant", "significant"])
        ax.set_ylabel(f"{'TEs / bp' if value_col=='TE_density' else 'TEs per window'}")
        title_mid = "TE density" if value_col=="TE_density" else "TE count"
        ax.set_title(f"{title_mid} by FST significance — {comp} (q_poi < 0.05)")
        ax.grid(alpha=0.2, axis="y")
    axes[-1].set_xlabel("Window group")
    fig.tight_layout()
    out = outdir / f"te_{tag}_violin_box_PANEL.png"
    fig.savefig(out, dpi=220); plt.close(fig)
    return out

def forest_per_comp(df_types: pd.DataFrame, comp: str, outdir: Path, tag="density"):
    d = df_types.sort_values("RR", ascending=False).copy()
    y = np.arange(len(d))[::-1]
    fig, ax = plt.subplots(figsize=(7.5, 0.35*max(6, len(d))))
    ax.errorbar(d["RR"], y,
                xerr=[d["RR"]-d["RR_CI_lo"], d["RR_CI_hi"]-d["RR"]],
                fmt='o', capsize=3)
    ax.axvline(1.0, color="#999999", lw=1, ls="--")
    ax.set_yticks(y, d["TE_type"])
    ax.set_xscale("log")
    ax.set_xlabel("Rate ratio (significant / not)   [log scale]")
    ax.set_title(f"Per-TE-type enrichment in significant windows — {comp} ({tag})")
    ax.grid(alpha=0.2, axis="x")
    out = outdir / f"te_type_{tag}_rate_ratio_forest_{comp}.png"
    fig.tight_layout(); fig.savefig(out, dpi=220); plt.close(fig)
    return out

def compact_forest_panel(per_comp_types: dict, outdir: Path, tag="density", max_types=12):
    """
    Compact 3-row panel; each row shows top-N TE types by significance (q) for that comparison.
    """
    fig, axes = plt.subplots(3, 1, figsize=(9, 9), sharex=False)
    for i, comp in enumerate(["3vs4","3vs5","4vs5"]):
        d = per_comp_types[comp].sort_values("q").head(max_types).copy()
        d = d.sort_values("RR", ascending=False)
        y = np.arange(len(d))[::-1]
        ax = axes[i]
        ax.errorbar(d["RR"], y, xerr=[d["RR"]-d["RR_CI_lo"], d["RR_CI_hi"]-d["RR"]],
                    fmt='o', capsize=3)
        ax.axvline(1.0, color="#999999", lw=1, ls="--")
        ax.set_yticks(y, d["TE_type"])
        ax.set_xscale("log")
        ax.set_xlabel("Rate ratio (sig/not)   [log scale]")
        ax.set_title(f"Top TE types (by FDR) — {comp} ({tag})")
        ax.grid(alpha=0.2, axis="x")
    fig.tight_layout()
    out = outdir / f"te_type_{tag}_rate_ratio_forest_PANEL.png"
    fig.savefig(out, dpi=220); plt.close(fig)
    return out

# =========================
# Report
# =========================
def add_slide(prs, title, bullets=None, image=None):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if bullets:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.6))
        tf = tb.text_frame; tf.clear()
        for i,b in enumerate(bullets):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.text = b; p.level = 0
    if image and Path(image).exists():
        slide.shapes.add_picture(str(image), Inches(0.5), Inches(2.7), height=Inches(3.8))

# =========================
# Main
# =========================
def main():
    # Load data
    fst = read_fst_windows(FST_CSV)
    te  = read_te_table(TE_TSV)

    # TE overlaps & densities
    te_counts, te_types = count_te_overlaps_per_window(fst, te)
    df_aug = fst.merge(te_counts, on=["CHROM","WIN_START","WIN_END","WIN_LEN"], how="left")
    num_cols = df_aug.select_dtypes(include=[np.number]).columns
    df_aug[num_cols] = df_aug[num_cols].fillna(0)

    # Per-comparison analyses
    comps = ["3vs4","3vs5","4vs5"]

    # --- Density track containers ---
    comp_masks_density = {}
    per_comp_overall_density = {}
    per_comp_types_density = {}
    per_comp_forest_density = {}

    # --- Count track containers ---
    comp_masks_count = {}
    per_comp_overall_count = {}
    per_comp_types_count = {}
    per_comp_forest_count = {}

    for comp in comps:
        qcol = f"q_poi_{comp}"
        if qcol not in df_aug.columns:
            raise ValueError(f"Expected Poisson q-value column missing: {qcol}")

        # Mask: significant windows for THIS comparison
        is_sig = (df_aug[qcol] < Q_CUTOFF)

        # =========================
        # DENSITY TRACK (offset = log WIN_LEN)
        # =========================
        y_den = df_aug["TE_count"].astype(int).to_numpy()
        grp   = is_sig.astype(int).to_numpy()
        off   = np.log(df_aug["WIN_LEN"].astype(float).to_numpy())

        res_used_den = glm_compare_counts(y_den, grp, offset_log=off)
        model_used_den = res_used_den["family"]

        # global RR (rates per bp)
        k_sig   = int(df_aug.loc[is_sig,  "TE_count"].sum())
        e_sig   = float(df_aug.loc[is_sig,  "WIN_LEN"].sum())
        k_nsig  = int(df_aug.loc[~is_sig, "TE_count"].sum())
        e_nsig  = float(df_aug.loc[~is_sig, "WIN_LEN"].sum())
        rr_glob_den, ci_lo_den, ci_hi_den, p_glob_den = global_rate_ratio_test(k_sig, e_sig, k_nsig, e_nsig)

        # per-TE-type (density)
        rows_den = []
        for t in te_types:
            y_t = df_aug[f"TE_{t}"].astype(int).to_numpy()
            res_t = glm_compare_counts(y_t, grp, offset_log=off)
            rows_den.append({
                "TE_type": t,
                "family": res_t["family"],
                "RR": res_t["RR"], "RR_CI_lo": res_t["RR_CI_lo"], "RR_CI_hi": res_t["RR_CI_hi"],
                "p": res_t["p"], "dispersion": res_t["dispersion"]
            })
        df_types_den = pd.DataFrame(rows_den)
        df_types_den["q"] = multipletests(df_types_den["p"].fillna(1.0), method="fdr_bh")[1]
        df_types_den = df_types_den.sort_values("q")

        comp_masks_density[comp] = is_sig
        per_comp_overall_density[comp] = {
            "model_used": model_used_den,
            "RR_overall": res_used_den["RR"], "RR_CI_lo": res_used_den["RR_CI_lo"],
            "RR_CI_hi": res_used_den["RR_CI_hi"], "p_model": res_used_den["p"],
            "dispersion_est": res_used_den["dispersion"],
            "global_RR": rr_glob_den, "global_CI_lo": ci_lo_den, "global_CI_hi": ci_hi_den, "p_global": p_glob_den,
            "n_sig": int(is_sig.sum()), "n_not": int((~is_sig).sum())
        }
        per_comp_types_density[comp] = df_types_den
        # write per-comp table (density)
        OUTDIR.mkdir(parents=True, exist_ok=True)
        df_types_den.to_csv(OUTDIR / f"te_density_by_type_{comp}.tsv", sep="\t", index=False)
        # forest (density)
        per_comp_forest_density[comp] = forest_per_comp(df_types_den, comp, OUTDIR, tag="density")

        # =========================
        # COUNT TRACK (NO offset)
        # =========================
        y_cnt = df_aug["TE_count"].astype(int).to_numpy()
        res_used_cnt = glm_compare_counts(y_cnt, grp, offset_log=None)  # <-- no offset
        model_used_cnt = res_used_cnt["family"]

        # global “RR” using totals and same exposures for convenience (still reported)
        # Here ‘exposure’ is the number of windows (for a quick summary).
        # The GLM without offset already compares means directly; this RR is auxiliary.
        n_win_sig  = int(is_sig.sum())
        n_win_nsig = int((~is_sig).sum())
        mean_sig  = (df_aug.loc[is_sig,  "TE_count"].mean() if n_win_sig  > 0 else np.nan)
        mean_nsig = (df_aug.loc[~is_sig, "TE_count"].mean() if n_win_nsig > 0 else np.nan)
        rr_glob_cnt = (mean_sig / mean_nsig) if (mean_sig is not np.nan and mean_nsig not in (np.nan, 0)) else np.nan

        # per-TE-type (counts)
        rows_cnt = []
        for t in te_types:
            y_t = df_aug[f"TE_{t}"].astype(int).to_numpy()
            res_t = glm_compare_counts(y_t, grp, offset_log=None)
            rows_cnt.append({
                "TE_type": t,
                "family": res_t["family"],
                "RR": res_t["RR"], "RR_CI_lo": res_t["RR_CI_lo"], "RR_CI_hi": res_t["RR_CI_hi"],
                "p": res_t["p"], "dispersion": res_t["dispersion"]
            })
        df_types_cnt = pd.DataFrame(rows_cnt)
        df_types_cnt["q"] = multipletests(df_types_cnt["p"].fillna(1.0), method="fdr_bh")[1]
        df_types_cnt = df_types_cnt.sort_values("q")

        comp_masks_count[comp] = is_sig
        per_comp_overall_count[comp] = {
            "model_used": model_used_cnt,
            "RR_overall": res_used_cnt["RR"], "RR_CI_lo": res_used_cnt["RR_CI_lo"],
            "RR_CI_hi": res_used_cnt["RR_CI_hi"], "p_model": res_used_cnt["p"],
            "dispersion_est": res_used_cnt["dispersion"],
            "mean_TE_sig": mean_sig, "mean_TE_not": mean_nsig, "RR_means_sig_over_not": rr_glob_cnt,
            "n_sig": n_win_sig, "n_not": n_win_nsig
        }
        per_comp_types_count[comp] = df_types_cnt
        # write per-comp table (counts)
        df_types_cnt.to_csv(OUTDIR / f"te_count_by_type_{comp}.tsv", sep="\t", index=False)
        # forest (counts)
        per_comp_forest_count[comp] = forest_per_comp(df_types_cnt, comp, OUTDIR, tag="count")

    # Write the augmented windows (once)
    out_aug = OUTDIR / "fst_windows_with_TEcounts.tsv"
    df_aug.to_csv(out_aug, sep="\t", index=False)

    # Build combined summary tables for all comparisons
    # Density track
    rows_den = []
    for comp in comps:
        is_sig = comp_masks_density[comp]
        sub_sig  = df_aug[is_sig]
        sub_nsig = df_aug[~is_sig]
        ov = per_comp_overall_density[comp]
        rows_den.append({
            "track": "density",
            "comparison": comp,
            "n_windows": len(df_aug),
            "n_sig": ov["n_sig"], "n_not": ov["n_not"],
            "mean_density_sig": float((sub_sig["TE_count"]/sub_sig["WIN_LEN"]).mean()) if len(sub_sig)>0 else np.nan,
            "mean_density_not": float((sub_nsig["TE_count"]/sub_nsig["WIN_LEN"]).mean()) if len(sub_nsig)>0 else np.nan,
            "model_used": ov["model_used"],
            "RR_overall": ov["RR_overall"], "RR_CI_lo": ov["RR_CI_lo"], "RR_CI_hi": ov["RR_CI_hi"], "p_model": ov["p_model"],
            "dispersion_est": ov["dispersion_est"],
            "global_RR": ov["global_RR"], "global_CI_lo": ov["global_CI_lo"], "global_CI_hi": ov["global_CI_hi"], "p_global": ov["p_global"]
        })
    df_summary_den = pd.DataFrame(rows_den)
    df_summary_den.to_csv(OUTDIR / "te_density_group_summary.tsv", sep="\t", index=False)

    # Count track
    rows_cnt = []
    for comp in comps:
        is_sig = comp_masks_count[comp]
        sub_sig  = df_aug[is_sig]
        sub_nsig = df_aug[~is_sig]
        ov = per_comp_overall_count[comp]
        rows_cnt.append({
            "track": "count",
            "comparison": comp,
            "n_windows": len(df_aug),
            "n_sig": ov["n_sig"], "n_not": ov["n_not"],
            "mean_TE_sig": float(sub_sig["TE_count"].mean()) if len(sub_sig)>0 else np.nan,
            "mean_TE_not": float(sub_nsig["TE_count"].mean()) if len(sub_nsig)>0 else np.nan,
            "model_used": ov["model_used"],
            "RR_overall": ov["RR_overall"], "RR_CI_lo": ov["RR_CI_lo"], "RR_CI_hi": ov["RR_CI_hi"], "p_model": ov["p_model"],
            "dispersion_est": ov["dispersion_est"],
            "RR_means_sig_over_not": ov["RR_means_sig_over_not"]
        })
    df_summary_cnt = pd.DataFrame(rows_cnt)
    df_summary_cnt.to_csv(OUTDIR / "te_count_group_summary.tsv", sep="\t", index=False)

    # Panel plots
    fig_violin_panel_den = violin_box_panel(df_aug, comp_masks_density, OUTDIR, value_col="TE_density", tag="density")
    fig_violin_panel_cnt = violin_box_panel(df_aug, comp_masks_count,   OUTDIR, value_col="TE_count",   tag="count")
    fig_forest_panel_den = compact_forest_panel(per_comp_types_density, OUTDIR, tag="density", max_types=12)
    fig_forest_panel_cnt = compact_forest_panel(per_comp_types_count,   OUTDIR, tag="count",   max_types=12)

    # PPTX
    prs = Presentation()
    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TE vs FST windows — density & count tracks (per comparison)"
    slide.placeholders[1].text = (
        f"Significance: q_poi < {Q_CUTOFF} for each of: 3vs4, 3vs5, 4vs5\n"
        "Density track: GLM with log(offset = window length)\n"
        "Count track:   GLM with no offset (raw counts per window)\n"
        "Auto-switch to NegBin if overdispersed; FDR per comparison & per TE type"
    )

    # Per-comparison bullets (both tracks)
    for comp in comps:
        ovd = per_comp_overall_density[comp]
        ovc = per_comp_overall_count[comp]
        bullets = [
            f"{comp}: significant / not = {ovd['n_sig']:,} / {ovd['n_not']:,}",
            f"[Density] RR(sig/not) = {ovd['RR_overall']:.3f} "
            f"[{ovd['RR_CI_lo']:.3f}, {ovd['RR_CI_hi']:.3f}], p={ovd['p_model']:.2e}; model={ovd['model_used']}, disp≈{ovd['dispersion_est']:.2f}",
            f"[Count]   RR(exp coef) = {ovc['RR_overall']:.3f} "
            f"[{ovc['RR_CI_lo']:.3f}, {ovc['RR_CI_hi']:.3f}], p={ovc['p_model']:.2e}; model={ovc['model_used']}, disp≈{ovc['dispersion_est']:.2f}",
            f"[Count]   Means: sig={ovc['mean_TE_sig']:.3f}, not={ovc['mean_TE_not']:.3f}  (ratio≈{ovc['RR_means_sig_over_not']:.3f})",
        ]
        add_slide(prs, f"Summary — {comp}", bullets=bullets, image=str(forest_per_comp(per_comp_types_density[comp], comp, OUTDIR, tag="density")))

    # Panel slides
    add_slide(prs, "TE density violin+box — PANEL", bullets=None, image=str(fig_violin_panel_den))
    add_slide(prs, "TE count violin+box — PANEL",   bullets=None, image=str(fig_violin_panel_cnt))
    add_slide(prs, "Per-TE-type RR (density) — PANEL", bullets=None, image=str(fig_forest_panel_den))
    add_slide(prs, "Per-TE-type RR (count)   — PANEL", bullets=None, image=str(fig_forest_panel_cnt))

    pptx_path = OUTDIR / "TE_vs_FST_report.pptx"
    prs.save(str(pptx_path))

    print(f"\nDone. Outputs in: {OUTDIR.resolve()}")
    print(f"- Augmented windows: {out_aug.name}")
    print(f"- Density summaries: te_density_group_summary.tsv; per-type: te_density_by_type_<comp>.tsv")
    print(f"- Count summaries:   te_count_group_summary.tsv;   per-type: te_count_by_type_<comp>.tsv")
    print(f"- Panel figures:     {Path(fig_violin_panel_den).name}, {Path(fig_violin_panel_cnt).name},")
    print(f"                     {Path(fig_forest_panel_den).name}, {Path(fig_forest_panel_cnt).name}")
    print(f"- Report:            {pptx_path.name}")

if __name__ == "__main__":
    main()
