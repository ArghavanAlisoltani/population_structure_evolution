#pip install pandas numpy scipy statsmodels matplotlib python-pptx
#!/usr/bin/env python3
"""
FST window enrichment with Binomial, Poisson, Beta–Binomial, Negative Binomial;
BH-FDR; benchmarking multiple window sizes/overlaps; plots; PowerPoint report.

Input columns (exact): CHROM, POS, PROC3_vs_PROC4, PROC3_vs_PROC5, PROC4_vs_PROC5
"""

import re
import itertools
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

from scipy.stats import binom, poisson, betabinom, nbinom
from statsmodels.stats.multitest import multipletests
from pptx import Presentation
from pptx.util import Inches

# =========================
# User knobs
# =========================
IN_PATH = "24_scaff_fst_for_window_richness.tsv"

# Benchmarks
WIN_SIZES = [10_000, 100_000, 1_000_000, 10_000_000]
OVERLAPS  = [2_000,   20_000,  200_000,   2_000_000]   # (Option A guard will skip invalid pairs)

# Thresholds
THR_LO = 0.0     # counts C0_* = number of sites with FST > THR_LO
THR_HI = 0.25    # counts C0.25_* = number of sites with FST >= THR_HI
RATIO_WHEN_C0_IS0 = 0.0

# Plotting / outputs
SCAFFOLD_OF_INTEREST = "scaffold_1a"   # optional: used for per-combo ratio plots
OUTDIR = Path("py_fst_windows_out")
OUTDIR.mkdir(parents=True, exist_ok=True)

# ==================================
# IO & cleaning
# ==================================
def load_and_clean(in_path: str) -> pd.DataFrame:
    need = ["CHROM", "POS", "PROC3_vs_PROC4", "PROC3_vs_PROC5", "PROC4_vs_PROC5"]
    df = pd.read_csv(in_path, sep="\t")
    missing = [c for c in need if c not in df.columns]
    if missing:
        raise ValueError(f"Missing required columns: {missing}")
    df = df[need].copy()
    for c in ["PROC3_vs_PROC4", "PROC3_vs_PROC5", "PROC4_vs_PROC5"]:
        x = pd.to_numeric(df[c], errors="coerce")
        x = x.replace([np.inf, -np.inf], np.nan).fillna(0.0)
        x = x.mask(x < 0, 0.0)
        df[c] = x.astype(float)
    df["POS"] = pd.to_numeric(df["POS"], errors="coerce").fillna(0).astype(np.int64)
    df["CHROM"] = df["CHROM"].astype(str)
    df = df.sort_values(["CHROM", "POS"]).reset_index(drop=True)
    return df

# ==================================
# Windowing & counting
# ==================================
def windows_for_scaffold(pos_min: int, pos_max: int, win_size: int, overlap: int) -> np.ndarray:
    step = win_size - overlap
    if step <= 0:
        raise ValueError("overlap must be smaller than window size (step <= 0).")
    starts = np.arange(pos_min, pos_max + 1, step, dtype=np.int64)
    ends   = np.minimum(starts + win_size - 1, pos_max)
    return np.vstack([starts, ends]).T

def count_in_windows_for_scaffold(pos: np.ndarray, vals: dict, win_arr: np.ndarray) -> dict:
    # Two-pointer scan
    N = pos.size
    M = win_arr.shape[0]
    out = {}
    for lab in vals.keys():
        out[f"C0_{lab}"] = np.zeros(M, dtype=np.int64)
        out[f"C0.25_{lab}"] = np.zeros(M, dtype=np.int64)

    i_left = 0
    i_right = 0
    for j, (w_start, w_end) in enumerate(win_arr):
        while i_left < N and pos[i_left] < w_start:
            i_left += 1
        if i_right < i_left:
            i_right = i_left
        while i_right < N and pos[i_right] <= w_end:
            i_right += 1
        if i_left < i_right:
            sl = slice(i_left, i_right)
            for lab, arr in vals.items():
                arr_slice = arr[sl]
                out[f"C0_{lab}"][j]    = int(np.sum(arr_slice >  THR_LO))
                out[f"C0.25_{lab}"][j] = int(np.sum(arr_slice >= THR_HI))
    return out

# ==================================
# Tests (Bin, Poi, BetaBin, NegBin) + FDR + dispersion
# ==================================
def bh_fdr(p):
    p = np.asarray(p, float)
    mask = np.isfinite(p)
    q = np.full_like(p, np.nan, dtype=float)
    if mask.sum() > 0:
        q[mask] = multipletests(p[mask], method="fdr_bh")[1]
    return q

def tests_for_comp(dfw: pd.DataFrame, lab: str) -> pd.DataFrame:
    """
    lab in {"3vs4", "3vs5", "4vs5"}.
    Adds columns:
      p_glob_*, lambda_*, disp_bin_*, disp_poi_*,
      p_bin_*, q_bin_*, p_poi_*, q_poi_*,
      p_bb_*,  q_bb_*,  p_nb_*,  q_nb_*,
      theta_nb_*, rho_bb_*, p_hat_bb_* (rho only)
    """
    G = dfw[f"C0_{lab}"].to_numpy(np.int64)
    K = dfw[f"C0.25_{lab}"].to_numpy(np.int64)

    totG = G.sum()
    totK = K.sum()
    p_glob = (totK / totG) if totG > 0 else np.nan
    lam = p_glob

    # --- Binomial (upper tail)
    if np.isnan(p_glob):
        p_bin = np.full_like(G, np.nan, dtype=float)
    else:
        p_bin = np.ones_like(G, dtype=float)
        ok = G > 0
        p_bin[ok] = 1.0 - binom.cdf(np.maximum(0, K[ok]) - 1, G[ok], p_glob)
        p_bin[~ok & (K == 0)] = 1.0
        p_bin[~ok & (K >  0)] = 0.0

    # --- Poisson with exposure (upper tail)
    if np.isnan(lam):
        mu = np.full_like(G, np.nan, dtype=float)
        p_poi = np.full_like(G, np.nan, dtype=float)
    else:
        mu = lam * G
        p_poi = np.ones_like(G, dtype=float)
        ok = mu > 0
        p_poi[ok] = 1.0 - poisson.cdf(np.maximum(0, K[ok]) - 1, mu[ok])
        p_poi[~ok & (K == 0)] = 1.0
        p_poi[~ok & (K >  0)] = 0.0

    # --- Dispersion diagnostics (global)
    if not (np.isnan(p_glob) or p_glob <= 0 or p_glob >= 1):
        V = G * p_glob * (1.0 - p_glob)
        ok = V > 0
        dfree = max(1, int(ok.sum()) - 1)
        disp_bin = float(np.sum(((K[ok] - G[ok]*p_glob)**2) / V[ok]) / dfree)
    else:
        disp_bin = np.nan

    if np.any(np.isfinite(mu)) and np.any(mu > 0):
        ok = mu > 0
        dfree = max(1, int(ok.sum()) - 1)
        disp_poi = float(np.sum(((K[ok] - mu[ok])**2) / mu[ok]) / dfree)
    else:
        disp_poi = np.nan

    # --- Beta–Binomial (method-of-moments for rho)
    n_bar = G[G > 0].mean() if (G > 0).any() else np.nan
    if np.isfinite(disp_bin) and np.isfinite(n_bar) and n_bar > 1.0:
        rho_est = max(0.0, min(0.9999, (disp_bin - 1.0) / (n_bar - 1.0)))
    else:
        rho_est = 0.0
    if rho_est <= 1e-10 or np.isnan(p_glob) or p_glob <= 0 or p_glob >= 1:
        p_bb = p_bin.copy()
    else:
        t = (1.0 / rho_est) - 1.0
        a_hat = p_glob * t
        b_hat = (1.0 - p_glob) * t
        p_bb = np.ones_like(G, dtype=float)
        ok = G > 0
        p_bb[ok] = betabinom.sf(np.maximum(0, K[ok]) - 1, G[ok], a_hat, b_hat)
        p_bb[~ok & (K == 0)] = 1.0
        p_bb[~ok & (K >  0)] = 0.0

    # --- Negative Binomial (method-of-moments for theta)
    if np.any(np.isfinite(mu)):
        resid = (K - mu)
        denom = np.sum((resid**2) - mu)
        numer = np.sum(mu**2)
        if denom > 1e-12 and numer > 0:
            theta_est = max(1e-6, numer / denom)
        else:
            theta_est = 1e9  # ~Poisson
    else:
        theta_est = np.nan

    p_nb = np.full_like(G, np.nan, dtype=float)
    if np.isfinite(theta_est) and np.any(np.isfinite(mu)):
        ok = mu >= 0
        r = theta_est
        p_param = np.where(ok, r / (r + mu), np.nan)
        p_nb = np.ones_like(G, dtype=float)
        has_trials = G > 0
        p_nb[has_trials] = nbinom.sf(np.maximum(0, K[has_trials]) - 1, r, p_param[has_trials])
        p_nb[~has_trials & (K == 0)] = 1.0
        p_nb[~has_trials & (K >  0)] = 0.0

    out = pd.DataFrame({
        f"p_glob_{lab}": np.repeat(p_glob, len(dfw)),
        f"lambda_{lab}": np.repeat(lam, len(dfw)),
        f"disp_bin_{lab}": np.repeat(disp_bin, len(dfw)),
        f"disp_poi_{lab}": np.repeat(disp_poi, len(dfw)),
        f"rho_bb_{lab}": np.repeat(rho_est, len(dfw)),
        f"theta_nb_{lab}": np.repeat(theta_est, len(dfw)),
        f"p_bin_{lab}": p_bin,  f"q_bin_{lab}": bh_fdr(p_bin),
        f"p_poi_{lab}": p_poi,  f"q_poi_{lab}": bh_fdr(p_poi),
        f"p_bb_{lab}":  p_bb,   f"q_bb_{lab}":  bh_fdr(p_bb),
        f"p_nb_{lab}":  p_nb,   f"q_nb_{lab}":  bh_fdr(p_nb),
    })
    return out

# ==================================
# Analysis per (win, overlap)
# ==================================
def analyze_one_combo(df: pd.DataFrame, win_size: int, overlap: int) -> pd.DataFrame:
    results = []
    for chrom, dfc in df.groupby("CHROM", sort=False):
        pos = dfc["POS"].to_numpy(np.int64)
        vals = {
            "3vs4": dfc["PROC3_vs_PROC4"].to_numpy(float),
            "3vs5": dfc["PROC3_vs_PROC5"].to_numpy(float),
            "4vs5": dfc["PROC4_vs_PROC5"].to_numpy(float),
        }
        if pos.size == 0:
            continue
        try:
            win_arr = windows_for_scaffold(int(pos.min()), int(pos.max()), win_size, overlap)
        except ValueError:
            print(f"[SKIP] {chrom}: invalid step for win={win_size:,}, ov={overlap:,}")
            continue
        if win_arr.size == 0:
            continue

        counts = count_in_windows_for_scaffold(pos, vals, win_arr)
        dfw = pd.DataFrame({"CHROM": chrom,
                            "WIN_START": win_arr[:, 0],
                            "WIN_END":   win_arr[:, 1]})
        for k, v in counts.items():
            dfw[k] = v

        # Ratios
        for lab in ["3vs4", "3vs5", "4vs5"]:
            c0, c025 = dfw[f"C0_{lab}"].to_numpy(), dfw[f"C0.25_{lab}"].to_numpy()
            with np.errstate(divide="ignore", invalid="ignore"):
                dfw[f"R_{lab}"] = np.where(c0 > 0, c025 / c0, RATIO_WHEN_C0_IS0)

        # Tests
        t34 = tests_for_comp(dfw, "3vs4")
        t35 = tests_for_comp(dfw, "3vs5")
        t45 = tests_for_comp(dfw, "4vs5")
        dfw = pd.concat([dfw, t34, t35, t45], axis=1)
        dfw["WIN_MID"] = ((dfw["WIN_START"] + dfw["WIN_END"]) // 2).astype(np.int64)
        results.append(dfw)

    if not results:
        return pd.DataFrame()
    return pd.concat(results, ignore_index=True)

# ==================================
# Plot helpers (ratio + Manhattan with 3 panels)
# ==================================
def save_ratio_plot(dfw: pd.DataFrame, scaffold: str, win_size: int, overlap: int, outdir: Path) -> Path | None:
    sub = dfw[dfw["CHROM"] == scaffold].copy()
    if sub.empty:
        return None
    plot_df = pd.melt(sub, id_vars=["WIN_MID"],
                      value_vars=["R_3vs4","R_3vs5","R_4vs5"],
                      var_name="comparison", value_name="ratio").sort_values(["comparison","WIN_MID"])
    fig, axes = plt.subplots(3, 1, figsize=(10, 6), sharex=True)
    for i, comp in enumerate(["R_3vs4","R_3vs5","R_4vs5"]):
        dd = plot_df[plot_df["comparison"] == comp]
        axes[i].plot(dd["WIN_MID"], dd["ratio"], lw=0.8)
        axes[i].scatter(dd["WIN_MID"], dd["ratio"], s=6)
        axes[i].set_ylim(0, 1)
        axes[i].set_ylabel("ratio")
        axes[i].set_title(comp)
        axes[i].grid(alpha=0.2)
    axes[-1].set_xlabel("Genomic position (bp)")
    fig.suptitle(f"High-FST fraction per window: {scaffold}\n(win={win_size:,}, overlap={overlap:,})", y=0.98)
    fig.tight_layout(rect=[0,0,1,0.95])
    out_png = outdir / f"ratios_{scaffold}_win{win_size}_ov{overlap}.png"
    fig.savefig(out_png, dpi=200); plt.close(fig)
    return out_png

def scaffold_sort_key(chrom: str):
    # Extract first integer in CHROM for numeric ordering; fall back to name
    m = re.search(r"(\d+)", chrom)
    return (int(m.group(1)) if m else 10**12, chrom)

def make_manhattan_three(dfw: pd.DataFrame, test_kind: str, out_png: Path, title_prefix: str,
                         qcols_map: dict, top_n: int = 10):
    """
    Make a single Manhattan figure with 3 rows (3vs4, 3vs5, 4vs5) for a given test kind.
    qcols_map: {"3vs4": "q_poi_3vs4", "3vs5": "...", "4vs5": "..."} or the NB equivalents
    Colors: sig (q<0.05)=black, non-sig=gray; top10=red with labels.
    """
    if dfw.empty:
        return None

    # Order scaffolds numerically and build cumulative offsets shared across panels
    chroms = sorted(dfw["CHROM"].unique(), key=scaffold_sort_key)
    offsets = {}
    cursor = 0
    xticks = []
    xtlabs = []
    # Use per-chrom max WIN_MID for spacing on the x-axis
    for ch in chroms:
        sub = dfw[dfw["CHROM"] == ch]
        if sub.empty: 
            continue
        max_pos = int(sub["WIN_MID"].max())
        offsets[ch] = cursor
        xticks.append(cursor + max_pos // 2)
        xtlabs.append(ch)
        cursor += max_pos + 1

    comps = ["3vs4", "3vs5", "4vs5"]
    fig, axes = plt.subplots(3, 1, figsize=(12, 8), sharex=True)

    for idx, comp in enumerate(comps):
        qcol = qcols_map[comp]
        if qcol not in dfw.columns:
            continue
        dd = dfw[["CHROM","WIN_START","WIN_END","WIN_MID", qcol]].copy()
        dd["offset"] = dd["CHROM"].map(offsets).fillna(0).astype(np.int64)
        dd["cumpos"] = dd["offset"] + dd["WIN_MID"]
        q = dd[qcol].astype(float)
        ml10 = -np.log10(np.clip(q, 1e-300, 1.0))
        dd["ml10"] = ml10
        dd["sig"] = (q < 0.05)

        # Top N
        top_idx = np.argsort(-ml10)[:top_n]
        dd_top = dd.iloc[top_idx]

        ax = axes[idx]
        # gray (not sig), black (sig)
        ax.scatter(dd.loc[~dd["sig"], "cumpos"], dd.loc[~dd["sig"], "ml10"],
                   s=6, c="#bdbdbd", linewidths=0)
        ax.scatter(dd.loc[dd["sig"], "cumpos"], dd.loc[dd["sig"], "ml10"],
                   s=8, c="black", linewidths=0)
        ax.scatter(dd_top["cumpos"], dd_top["ml10"], s=22, c="red", linewidths=0, zorder=5)
        for _, r in dd_top.iterrows():
            ax.text(r["cumpos"], r["ml10"] + 0.1,
                    f"{r['CHROM']}:{int(r['WIN_START'])}-{int(r['WIN_END'])}",
                    fontsize=7, color="red", rotation=45, ha="left", va="bottom")
        ax.set_ylabel(r"$-\log_{10}(q)$")
        ax.set_title(comp)
        ax.grid(alpha=0.2, axis="y")

    axes[-1].set_xlabel("Genome (ordered scaffolds)")
    plt.xticks(xticks, xtlabs, rotation=90, fontsize=7)
    fig.suptitle(title_prefix, y=0.98)
    fig.tight_layout(rect=[0,0,1,0.95])
    fig.savefig(out_png, dpi=220)
    plt.close(fig)
    return out_png

# ==================================
# Report slides
# ==================================
def add_slide(prs, title, bullets=None, image=None):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if bullets:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(3.7))
        tf = tb.text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b; p.level = 0
    if image and Path(image).exists():
        slide.shapes.add_picture(str(image), Inches(0.5), Inches(2.8), height=Inches(3.8))

def make_comparison_plots(summary_df: pd.DataFrame, outdir: Path) -> list[Path]:
    imgs = []
    # test kinds to plot across window sizes (lines per overlap)
    for testkind in ["bin", "poi", "bb", "nb"]:
        for comp in ["3vs4", "3vs5", "4vs5"]:
            col = f"nsig_{testkind}_{comp}"
            if col not in summary_df.columns:
                continue
            pivot = summary_df.pivot_table(index="win_size", columns="overlap",
                                           values=col, aggfunc="sum").sort_index()
            fig, ax = plt.subplots(figsize=(8, 5))
            for ov in pivot.columns:
                ax.plot(pivot.index, pivot[ov], marker="o", label=f"overlap={ov:,}")
            ax.set_xscale("log")
            ax.set_xlabel("Window size (bp, log scale)")
            ax.set_ylabel(f"# windows with FDR<0.05 ({testkind.upper()}, {comp})")
            ax.set_title(f"Significant windows vs. window size — {testkind.upper()} — {comp}")
            ax.grid(alpha=0.3)
            ax.legend(fontsize=8)
            fig.tight_layout()
            out_png = outdir / f"compare_nsig_{testkind}_{comp}.png"
            fig.savefig(out_png, dpi=200); plt.close(fig)
            imgs.append(out_png)
    return imgs

def make_report(all_summaries: list, pptx_path: Path, comparison_imgs: list[Path]):
    prs = Presentation()
    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "FST Window Enrichment Report"
    slide.placeholders[1].text = (
        "Binomial, Poisson, Beta–Binomial, Negative Binomial with BH-FDR\n"
        "Benchmarking across window sizes and overlaps\n"
        "Includes dispersion diagnostics and Manhattan plots (3 panels per figure)"
    )
    # Per-combo slides
    for s in all_summaries:
        title = f"win={s['win_size']:,}, overlap={s['overlap']:,}"
        bullets = [
            f"Windows: {s['n_windows']:,} across {s['n_scaffolds']} scaffolds",
            f"Dispersion (BIN): 3v4={s['disp_bin_3vs4']:.2f}, 3v5={s['disp_bin_3vs5']:.2f}, 4v5={s['disp_bin_4vs5']:.2f}",
            f"Dispersion (POI): 3v4={s['disp_poi_3vs4']:.2f}, 3v5={s['disp_poi_3vs5']:.2f}, 4v5={s['disp_poi_4vs5']:.2f}",
            f"rho (BB):         3v4={s['rho_bb_3vs4']:.3f}, 3v5={s['rho_bb_3vs5']:.3f}, 4v5={s['rho_bb_4vs5']:.3f}",
            f"theta (NB):       3v4={s['theta_nb_3vs4']:.2f}, 3v5={s['theta_nb_3vs5']:.2f}, 4v5={s['theta_nb_4vs5']:.2f}",
            f"Sig windows FDR<0.05 BIN/POI/BB/NB — 3v4: {s['nsig_bin_3vs4']:,}/{s['nsig_poi_3vs4']:,}/{s['nsig_bb_3vs4']:,}/{s['nsig_nb_3vs4']:,}",
            f"                                                   3v5: {s['nsig_bin_3vs5']:,}/{s['nsig_poi_3vs5']:,}/{s['nsig_bb_3vs5']:,}/{s['nsig_nb_3vs5']:,}",
            f"                                                   4v5: {s['nsig_bin_4vs5']:,}/{s['nsig_poi_4vs5']:,}/{s['nsig_bb_4vs5']:,}/{s['nsig_nb_4vs5']:,}",
            f"Shared significant (POI): {s['shared_poi_n']} (file: {s['shared_poi_csv']})",
            f"Shared significant (NB):  {s['shared_nb_n']}  (file: {s['shared_nb_csv']})",
            f"Output: {s['csv_name']}",
            "Manhattan plots (POI & NB, 3 panels) saved in combo folder."
        ]
        add_slide(prs, title, bullets=bullets, image=s.get("ratio_plot_path"))
    # Comparison slides
    for img in comparison_imgs:
        add_slide(prs, f"Comparison: {img.stem}", image=img)
    prs.save(str(pptx_path))

# ==================================
# Main
# ==================================
def main():
    df = load_and_clean(IN_PATH)

    all_summaries = []

    # run all combos with Option A guard
    for win_size, overlap in itertools.product(WIN_SIZES, OVERLAPS):
        if overlap >= win_size:
            print(f"[SKIP] win={win_size:,} overlap={overlap:,} (overlap must be < window size)")
            continue

        combo_dir = OUTDIR / f"win{win_size}_ov{overlap}"
        combo_dir.mkdir(parents=True, exist_ok=True)

        dfw = analyze_one_combo(df, win_size, overlap)
        if dfw.empty:
            print(f"[WARN] No windows for win={win_size:,}, ov={overlap:,}")
            continue

        # Save full table
        csv_path = combo_dir / f"fst_window_counts_tests_win{win_size}_ov{overlap}.csv"
        dfw.to_csv(csv_path, index=False)

        # Per-comparison numeric summaries
        def comp_nums(prefix: str):
            d_bin = float(dfw[f"disp_bin_{prefix}"].iloc[0])
            d_poi = float(dfw[f"disp_poi_{prefix}"].iloc[0])
            rho   = float(dfw[f"rho_bb_{prefix}"].iloc[0])
            theta = float(dfw[f"theta_nb_{prefix}"].iloc[0])
            ns_bin = int((dfw[f"q_bin_{prefix}"] < 0.05).sum())
            ns_poi = int((dfw[f"q_poi_{prefix}"] < 0.05).sum())
            ns_bb  = int((dfw[f"q_bb_{prefix}"]  < 0.05).sum())
            ns_nb  = int((dfw[f"q_nb_{prefix}"]  < 0.05).sum())
            return d_bin, d_poi, rho, theta, ns_bin, ns_poi, ns_bb, ns_nb

        c34 = comp_nums("3vs4")
        c35 = comp_nums("3vs5")
        c45 = comp_nums("4vs5")

        # ratio plot for one scaffold
        ratio_plot_path = save_ratio_plot(dfw, SCAFFOLD_OF_INTEREST, win_size, overlap, combo_dir)

        # Manhattan plots: one figure per test kind (POI + NB), each with 3 rows (3vs4, 3vs5, 4vs5)
        if "WIN_MID" not in dfw.columns:
            dfw["WIN_MID"] = ((dfw["WIN_START"] + dfw["WIN_END"]) // 2).astype(np.int64)

        poi_png = combo_dir / f"manhattan_POI_3panels_win{win_size}_ov{overlap}.png"
        nb_png  = combo_dir / f"manhattan_NB_3panels_win{win_size}_ov{overlap}.png"
        make_manhattan_three(
            dfw, "poi", poi_png,
            title_prefix=f"Poisson FDR (q) — 3 panels — win={win_size:,}, ov={overlap:,}",
            qcols_map={"3vs4": "q_poi_3vs4", "3vs5": "q_poi_3vs5", "4vs5": "q_poi_4vs5"},
            top_n=10
        )
        make_manhattan_three(
            dfw, "nb", nb_png,
            title_prefix=f"NegBin FDR (q) — 3 panels — win={win_size:,}, ov={overlap:,}",
            qcols_map={"3vs4": "q_nb_3vs4", "3vs5": "q_nb_3vs5", "4vs5": "q_nb_4vs5"},
            top_n=10
        )

        # Shared significant windows across all three comparisons (q<0.05 in each)
        # Do this for Poisson and NB separately
        poi_mask_all = (dfw["q_poi_3vs4"] < 0.05) & (dfw["q_poi_3vs5"] < 0.05) & (dfw["q_poi_4vs5"] < 0.05)
        nb_mask_all  = (dfw["q_nb_3vs4"]  < 0.05) & (dfw["q_nb_3vs5"]  < 0.05) & (dfw["q_nb_4vs5"]  < 0.05)

        shared_cols = ["CHROM","WIN_START","WIN_END","WIN_MID",
                       "C0_3vs4","C0.25_3vs4","C0_3vs5","C0.25_3vs5","C0_4vs5","C0.25_4vs5",
                       "q_poi_3vs4","q_poi_3vs5","q_poi_4vs5","q_nb_3vs4","q_nb_3vs5","q_nb_4vs5"]
        shared_cols = [c for c in shared_cols if c in dfw.columns]

        shared_poi = dfw.loc[poi_mask_all, shared_cols].copy()
        shared_nb  = dfw.loc[nb_mask_all,  shared_cols].copy()

        shared_poi_csv = combo_dir / f"shared_significant_POI_win{win_size}_ov{overlap}.csv"
        shared_nb_csv  = combo_dir / f"shared_significant_NB_win{win_size}_ov{overlap}.csv"
        shared_poi.to_csv(shared_poi_csv, index=False)
        shared_nb.to_csv(shared_nb_csv, index=False)

        # summary row
        summary = {
            "win_size": win_size, "overlap": overlap,
            "n_windows": len(dfw), "n_scaffolds": dfw["CHROM"].nunique(),
            "disp_bin_3vs4": c34[0], "disp_poi_3vs4": c34[1], "rho_bb_3vs4": c34[2], "theta_nb_3vs4": c34[3],
            "disp_bin_3vs5": c35[0], "disp_poi_3vs5": c35[1], "rho_bb_3vs5": c35[2], "theta_nb_3vs5": c35[3],
            "disp_bin_4vs5": c45[0], "disp_poi_4vs5": c45[1], "rho_bb_4vs5": c45[2], "theta_nb_4vs5": c45[3],
            "nsig_bin_3vs4": c34[4], "nsig_poi_3vs4": c34[5], "nsig_bb_3vs4": c34[6], "nsig_nb_3vs4": c34[7],
            "nsig_bin_3vs5": c35[4], "nsig_poi_3vs5": c35[5], "nsig_bb_3vs5": c35[6], "nsig_nb_3vs5": c35[7],
            "nsig_bin_4vs5": c45[4], "nsig_poi_4vs5": c45[5], "nsig_bb_4vs5": c45[6], "nsig_nb_4vs5": c45[7],
            "csv_name": csv_path.name,
            "ratio_plot_path": ratio_plot_path,
            # shared results
            "shared_poi_n": int(shared_poi.shape[0]),
            "shared_nb_n":  int(shared_nb.shape[0]),
            "shared_poi_csv": shared_poi_csv.name,
            "shared_nb_csv":  shared_nb_csv.name
        }
        all_summaries.append(summary)

        # write per-combo README
        with open(combo_dir / "summary.txt", "w") as fh:
            fh.write(f"win={win_size:,}, overlap={overlap:,}\n")
            fh.write(f"Windows: {summary['n_windows']:,} across {summary['n_scaffolds']} scaffolds\n")
            fh.write("Dispersion BIN (3v4/3v5/4v5): "
                     f"{summary['disp_bin_3vs4']:.2f}/{summary['disp_bin_3vs5']:.2f}/{summary['disp_bin_4vs5']:.2f}\n")
            fh.write("Dispersion POI (3v4/3v5/4v5): "
                     f"{summary['disp_poi_3vs4']:.2f}/{summary['disp_poi_3vs5']:.2f}/{summary['disp_poi_4vs5']:.2f}\n")
            fh.write("rho BB (3v4/3v5/4v5): "
                     f"{summary['rho_bb_3vs4']:.3f}/{summary['rho_bb_3vs5']:.3f}/{summary['rho_bb_4vs5']:.3f}\n")
            fh.write("theta NB (3v4/3v5/4v5): "
                     f"{summary['theta_nb_3vs4']:.2f}/{summary['theta_nb_3vs5']:.2f}/{summary['theta_nb_4vs5']:.2f}\n")
            fh.write("Sig windows FDR<0.05 BIN (3v4/3v5/4v5): "
                     f"{summary['nsig_bin_3vs4']}/{summary['nsig_bin_3vs5']}/{summary['nsig_bin_4vs5']}\n")
            fh.write("Sig windows FDR<0.05 POI (3v4/3v5/4v5): "
                     f"{summary['nsig_poi_3vs4']}/{summary['nsig_poi_3vs5']}/{summary['nsig_poi_4vs5']}\n")
            fh.write("Sig windows FDR<0.05  BB (3v4/3v5/4v5): "
                     f"{summary['nsig_bb_3vs4']}/{summary['nsig_bb_3vs5']}/{summary['nsig_bb_4vs5']}\n")
            fh.write("Sig windows FDR<0.05  NB (3v4/3v5/4v5): "
                     f"{summary['nsig_nb_3vs4']}/{summary['nsig_nb_3vs5']}/{summary['nsig_nb_4vs5']}\n")
            fh.write(f"Shared significant (POI): {summary['shared_poi_n']} — file: {summary['shared_poi_csv']}\n")
            fh.write(f"Shared significant (NB):  {summary['shared_nb_n']} — file: {summary['shared_nb_csv']}\n")
            fh.write(f"CSV: {csv_path.name}\n")

    # Build comparison plots across window sizes
    if all_summaries:
        summary_df = pd.DataFrame(all_summaries)
        # Add numeric summary fields for comparison slides
        def add_nsig_cols(df, testkind):
            for comp in ["3vs4","3vs5","4vs5"]:
                col = f"nsig_{testkind}_{comp}"
                if col not in df.columns:
                    # already in summary for testkinds
                    pass
        cmp_imgs = make_comparison_plots(summary_df, OUTDIR)
    else:
        summary_df = pd.DataFrame()
        cmp_imgs = []

    # Build PPTX
    pptx_path = OUTDIR / "FST_window_report.pptx"
    make_report(all_summaries, pptx_path, cmp_imgs)

    print(f"\nAll done.\nOutputs: {OUTDIR.resolve()}\nReport: {pptx_path.name}")
    if not summary_df.empty:
        print("\nSummary (head):")
        print(summary_df.head())

if __name__ == "__main__":
    main()

