#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
TE vs FST windows — v9 (superset of v8)

Adds:
  • Poisson GLM + log-exposure with dispersion diagnostics (χ²/df), figures & TSV
  • Negative Binomial GLM (offset kept) with RR, 95% CI, p, AIC
  • Wilcoxon test on TE density with rank-biserial effect and BH-FDR
  • Presence/absence Fisher & Chi-square (kept, enriched)
  • Optional Zero-Inflated Poisson/NB models (AIC/LL comparison)
  • Optional block permutation test (difference in mean TE density)
  • New figures: dispersion barplots, density violins/boxes, count histograms
  • Detailed PPTX including new analyses
Retains all v8 functionality (TE overlaps, per-class counts, GWAS overlay, Manhattan figures, presence tables).

CONFIG section below controls behavior and appearance.
"""

import os, re, math, random
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

from scipy.stats import norm, fisher_exact, chi2_contingency, mannwhitneyu
from statsmodels.api import GLM
import statsmodels.api as sm
from statsmodels.genmod.families import links
from statsmodels.genmod.families.family import Poisson, NegativeBinomial
from statsmodels.stats.multitest import multipletests
from pptx import Presentation
from pptx.util import Inches

# Optional zero-inflated models (enabled via CONFIG flag)
try:
    from statsmodels.discrete.count_model import ZeroInflatedPoisson, ZeroInflatedNegativeBinomialP
    HAS_ZI = True
except Exception:
    HAS_ZI = False

# =========================
# CONFIG
# =========================
FST_CSV = "fst_window_counts_tests_win1000000_ov200000.csv"
TE_TSV  = os.path.expanduser(
    "~/Desktop/OSU_projects/conifers/LP/lodgepole_pine_assembly/col8_readable.TEanno.cds.scaf1ab.tsv"
)
GWAS_XLSX = "Significant_GWAS_list_Aria_v1.xlsx"   # optional

Q_CUTOFF = 0.05

# Advanced options
RUN_ZERO_INFLATED = False   # set True to fit ZI Poisson/NB (can be slow)
RUN_PERMUTATION   = False   # set True for block permutation test on mean TE density

# Permutation settings
N_PERM      = 1000
BLOCK_SIZE  = 5_000_000     # 5 Mb blocks for label shuffling within scaffold

OUTDIR = Path("te_vs_fst_out_v9")
OUTDIR.mkdir(parents=True, exist_ok=True)

# Figure styling
BASE_FONTSIZE    = 11
PANEL_ROW_HEIGHT = 3.6
POINT_SIZE       = 10
ALPHA_SIG        = 1.0
ALPHA_NONSIG     = 0.18
COLOR_SCAF_A     = "#6BAED6"
COLOR_SCAF_B     = "#BDBDBD"
COLOR_TOP_LABELS = "red"
TOP_K_ANNOTATE   = 5
SPACER_BP        = 1_000_000

COLOR_SIG        = "#1b9e77"
COLOR_NONSIG     = "#d95f02"

np.random.seed(1)
random.seed(1)

# =========================
# Utility & IO
# =========================
_scaf_pat = re.compile(r"scaffold[_\-]?(\d+)([ab])?$", re.IGNORECASE)

def scaffold_sort_key(s: str):
    """Return a sortable tuple for scaffold IDs like 'scaffold_1a'."""

    m = _scaf_pat.search(s)
    if not m:
        return (1e9, s)
    num = int(m.group(1))
    suf = m.group(2) or ""
    suf_rank = {"a":0, "b":1}.get(suf.lower(), 2)
    return (num, suf_rank)

def read_fst_windows(path: str) -> pd.DataFrame:
    """Read the primary FST window CSV and enforce the required schema."""

    df = pd.read_csv(path)
    need = ["CHROM","WIN_START","WIN_END","WIN_MID",
            "q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]
    miss = [c for c in need if c not in df.columns]
    if miss:
        raise ValueError(f"FST file missing required cols: {miss}")
    df["CHROM"]     = df["CHROM"].astype(str)
    df["WIN_START"] = pd.to_numeric(df["WIN_START"], errors="coerce").astype("Int64")
    df["WIN_END"]   = pd.to_numeric(df["WIN_END"],   errors="coerce").astype("Int64")
    df["WIN_MID"]   = pd.to_numeric(df["WIN_MID"],   errors="coerce").astype("Int64")
    df["WIN_LEN"]   = (df["WIN_END"] - df["WIN_START"] + 1).astype("Int64")
    for c in ["q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]:
        df[c] = pd.to_numeric(df[c], errors="coerce").clip(lower=1e-300, upper=1.0).fillna(1.0)
    return df

def _has_header(path: str, first_col_name: str) -> bool:
    """Return ``True`` if the first field of ``path`` matches ``first_col_name``."""

    with open(path, "r") as fh:
        first = fh.readline().strip().split("\t")[0]
        return first == first_col_name

def read_te_table(path: str) -> pd.DataFrame:
    """Read the TE annotation table and normalise column types/order."""

    cols = ["seqid","source","sequence_ontology","start","end","score","strand"]
    te = pd.read_csv(path, sep="\t", names=cols,
                     header=0 if _has_header(path, cols[0]) else None)
    te["seqid"] = te["seqid"].astype(str)
    te["sequence_ontology"] = te["sequence_ontology"].astype(str)
    te["start"] = pd.to_numeric(te["start"], errors="coerce").astype("Int64")
    te["end"]   = pd.to_numeric(te["end"],   errors="coerce").astype("Int64")
    te = te.dropna(subset=["seqid","start","end"]).copy()
    te = te.sort_values(["seqid","start","end"]).reset_index(drop=True)
    return te

def read_gwas(path: str) -> pd.DataFrame:
    """Read the optional GWAS Excel file used for Manhattan annotations."""

    if not Path(path).exists():
        print(f"[WARN] GWAS file not found at {path}; GWAS-annotated panels will be skipped.")
        return pd.DataFrame(columns=["Scaffold","Pos","traits"])
    df = pd.read_excel(path, engine="openpyxl")
    need = ["Scaffold","Pos","traits"]
    if not all(c in df.columns for c in need):
        raise ValueError(f"GWAS Excel missing required cols: {need}")
    df["Scaffold"] = df["Scaffold"].astype(str)
    df["Pos"]      = pd.to_numeric(df["Pos"], errors="coerce").astype("Int64")
    df["traits"]   = df["traits"].astype(str)
    df = df.dropna(subset=["Scaffold","Pos"]).reset_index(drop=True)
    return df

# =========================
# TE overlaps per window
# =========================
def _count_overlaps_one_chrom(wins: pd.DataFrame, tes: pd.DataFrame, te_types: np.ndarray):
    """Count TE overlaps for windows on a single chromosome using a two-pointer scan."""

    w_st = wins["WIN_START"].to_numpy(dtype=np.int64)
    w_en = wins["WIN_END"].to_numpy(dtype=np.int64)
    nW   = len(wins)

    total_counts = np.zeros(nW, dtype=np.int64)
    type_to_idx = {t:i for i,t in enumerate(te_types)}
    per_type = np.zeros((nW, len(te_types)), dtype=np.int64)

    if tes.empty:
        return total_counts, per_type

    t_st = tes["start"].to_numpy(dtype=np.int64)
    t_en = tes["end"].to_numpy(dtype=np.int64)
    t_ty = tes["sequence_ontology"].astype(str).to_numpy()

    j = 0
    for i in range(nW):
        ws, we = w_st[i], w_en[i]
        while j < len(t_st) and t_en[j] < ws:
            j += 1
        k = j
        while k < len(t_st) and t_st[k] <= we:
            if not (t_en[k] < ws or t_st[k] > we):
                total_counts[i] += 1
                idx = type_to_idx.get(t_ty[k], None)
                if idx is not None:
                    per_type[i, idx] += 1
            k += 1
    return total_counts, per_type

def count_te_overlaps_per_window(fst_df: pd.DataFrame, te_df: pd.DataFrame):
    """Augment the FST windows with total and per-class TE overlap counts."""

    te_types = np.array(sorted(te_df["sequence_ontology"].astype(str).unique()))
    base = fst_df[["CHROM","WIN_START","WIN_END","WIN_LEN","WIN_MID"]].copy().reset_index(drop=True)

    total_counts_all = np.zeros(len(base), dtype=np.int64)
    per_type_all = np.zeros((len(base), len(te_types)), dtype=np.int64)

    for chrom, wins in fst_df.groupby("CHROM", sort=False):
        rows = wins.index.to_numpy()
        wins_sorted = wins.sort_values("WIN_START")
        order_back = np.argsort(wins_sorted.index.to_numpy())
        tes = te_df[te_df["seqid"] == chrom].sort_values("start")

        total_counts, per_type = _count_overlaps_one_chrom(wins_sorted, tes, te_types)
        total_counts_all[rows] = total_counts[order_back]
        per_type_all[rows, :]  = per_type[order_back, :]

    te_cols = {"TE_count": total_counts_all.astype(np.int64)}
    for j, t in enumerate(te_types):
        te_cols[f"TE_{t}"] = per_type_all[:, j].astype(np.int64)

    te_df_out = pd.DataFrame(te_cols)
    out = pd.concat([base, te_df_out], axis=1)

    win_len = out["WIN_LEN"].astype(float)
    out["TE_density"] = out["TE_count"] / win_len
    for t in te_types:
        out[f"TEdens_{t}"] = out[f"TE_{t}"] / win_len

    return out, te_types

# =========================
# Genome axis (Gb) helpers
# =========================
def build_genome_coords(df):
    """Create genome-wide x coordinates for Manhattan plots."""

    scafs = sorted(df["CHROM"].unique(), key=scaffold_sort_key)
    offsets = {}
    current = 0
    for s in scafs:
        offsets[s] = current
        s_len = int(df.loc[df["CHROM"]==s, "WIN_END"].max())
        current += s_len + SPACER_BP
    df2 = df.copy()
    df2["X_bp"] = df2.apply(lambda r: offsets[r["CHROM"]] + float(r["WIN_MID"]), axis=1)
    df2["X_Gb"] = df2["X_bp"] / 1e9
    return df2, offsets

# =========================
# GWAS overlap
# =========================
def add_gwas_overlap_columns(df_aug: pd.DataFrame, gwas_df: pd.DataFrame) -> pd.DataFrame:
    """Append per-window GWAS overlap indicators and trait summaries."""

    if gwas_df is None or len(gwas_df) == 0:
        df_aug = df_aug.copy()
        df_aug["GWAS_n_hits"] = 0
        df_aug["GWAS_positions"] = ""
        df_aug["GWAS_traits"] = ""
        df_aug["GWAS_has_hit"] = 0
        return df_aug

    gw_by_scaf = {s: sub[["Pos","traits"]].astype({"Pos":int}).sort_values("Pos").reset_index(drop=True)
                  for s, sub in gwas_df.groupby("Scaffold")}
    n = len(df_aug)
    g_n = np.zeros(n, dtype=int)
    g_pos = []
    g_trt = []
    g_flag = np.zeros(n, dtype=int)

    for idx, r in df_aug.iterrows():
        sc = r["CHROM"]
        if sc not in gw_by_scaf:
            g_pos.append(""); g_trt.append(""); continue
        sub = gw_by_scaf[sc]
        ok  = (sub["Pos"] >= int(r["WIN_START"])) & (sub["Pos"] <= int(r["WIN_END"]))
        hits = sub.loc[ok]
        g_n[idx] = len(hits)
        g_flag[idx] = 1 if len(hits) > 0 else 0
        if len(hits):
            pos_list = hits["Pos"].astype(int).tolist()
            trt_list = hits["traits"].astype(str).tolist()
            g_pos.append(",".join(map(str, pos_list[:10])) + ("..." if len(pos_list)>10 else ""))
            uniq_tr = []
            seen = set()
            for t in trt_list:
                if t not in seen:
                    uniq_tr.append(t); seen.add(t)
            g_trt.append(";".join(uniq_tr[:10]) + ("..." if len(uniq_tr)>10 else ""))
        else:
            g_pos.append(""); g_trt.append("")
    out = df_aug.copy()
    out["GWAS_n_hits"] = g_n
    out["GWAS_positions"] = g_pos
    out["GWAS_traits"] = g_trt
    out["GWAS_has_hit"] = g_flag
    return out

# =========================
# GLM helpers
# =========================
def dispersion_ratio(y, mu, df_adj=1):
    """Return the Pearson Chi² dispersion ratio for GLM residuals."""

    y = np.asarray(y, float); mu = np.asarray(mu, float)
    ok = mu > 0
    if ok.sum() <= 1:
        return np.nan
    return float(np.sum(((y[ok]-mu[ok])**2)/mu[ok]) / max(ok.sum()-df_adj, 1))

def clamp_exp(x):  # safe exp for CI
    """Exponentiate while avoiding floating point overflow."""

    return float(np.exp(np.clip(x, -700, 700)))

def fit_glm_poisson(y, X, offset=None):
    """Fit a Poisson GLM with optional log-offset and return the result."""

    fam = Poisson(link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    return model.fit()

def estimate_nb_alpha(y, mu):
    """Method-of-moments estimate for the NB dispersion parameter alpha."""

    numer = np.sum((y - mu) ** 2 - mu)
    denom = np.sum(mu ** 2)
    if denom <= 0:
        return 1.0
    alpha = numer / denom
    return float(max(alpha, 1e-8))

def fit_glm_nb(y, X, offset=None, alpha=1.0):
    """Fit a Negative Binomial GLM with a fixed alpha."""

    fam = NegativeBinomial(alpha=alpha, link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    return model.fit()

def glm_compare_counts(y, group, offset_log=None, prefer_nb_if_overdispersed=True):
    """Compare group counts using Poisson/NB GLMs and return tidy summaries."""

    X = pd.DataFrame({"intercept": 1.0, "group": group.astype(float)})
    res_p = fit_glm_poisson(y, X, offset=offset_log)
    mu_p  = res_p.fittedvalues
    disp  = dispersion_ratio(y, mu_p)

    res_used = res_p
    family_used = "Poisson"
    if prefer_nb_if_overdispersed and np.isfinite(disp) and disp > 1.5:
        alpha = estimate_nb_alpha(y, mu_p)
        res_nb = fit_glm_nb(y, X, offset=offset_log, alpha=alpha)
        res_used = res_nb
        family_used = f"NegBin(alpha={alpha:.3g})"

    b  = res_used.params.get("group", np.nan)
    se = res_used.bse.get("group", np.nan)
    z  = (b / se) if se and np.isfinite(se) and se > 0 else np.nan
    p  = 2 * (1 - norm.cdf(abs(z))) if np.isfinite(z) else np.nan
    rr = math.exp(b) if np.isfinite(b) else np.nan
    ci_lo = clamp_exp(b - 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan
    ci_hi = clamp_exp(b + 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan

    return {"family": family_used, "coef": b, "se": se, "z": z, "p": p,
            "RR": rr, "RR_CI_lo": ci_lo, "RR_CI_hi": ci_hi,
            "dispersion": disp, "res": res_used, "res_p": res_p}

# =========================
# Presence tests (overall & per-class)
# =========================
def presence_tests(df, is_sig, count_col):
    """Fisher exact tests for TE presence vs FST significance."""

    a = int((df.loc[ is_sig, count_col] >= 1).sum())
    b = int((df.loc[ is_sig, count_col] == 0).sum())
    c = int((df.loc[~is_sig, count_col] >= 1).sum())
    d = int((df.loc[~is_sig, count_col] == 0).sum())
    table = np.array([[a,b],[c,d]], dtype=int)
    try: odds, p_f = fisher_exact(table, alternative="two-sided")
    except: odds, p_f = (np.nan, np.nan)
    try: chi2, p_c, dof, _ = chi2_contingency(table, correction=False)
    except: chi2, p_c, dof = (np.nan, np.nan, np.nan)
    return {"a_sig_ge1":a, "b_sig_eq0":b, "c_not_ge1":c, "d_not_eq0":d,
            "odds_ratio_fisher":odds, "p_fisher":p_f, "chi2":chi2, "p_chi2":p_c, "df":dof}

# =========================
# Nonparametric & permutation
# =========================
def wilcoxon_density(df, is_sig):
    """Wilcoxon rank-sum statistics for TE densities split by significance."""

    x_sig = df.loc[ is_sig, "TE_density"].astype(float).to_numpy()
    x_not = df.loc[~is_sig, "TE_density"].astype(float).to_numpy()
    if len(x_sig)==0 or len(x_not)==0:
        return dict(p=np.nan, U=np.nan, n_sig=len(x_sig), n_not=len(x_not),
                    med_sig=np.nan, med_not=np.nan, rbes=np.nan)
    U, p = mannwhitneyu(x_sig, x_not, alternative="two-sided")
    n1, n2 = len(x_sig), len(x_not)
    # rank-biserial effect size
    rbes = 1 - 2*U/(n1*n2)
    return dict(p=float(p), U=float(U), n_sig=n1, n_not=n2,
                med_sig=float(np.median(x_sig)), med_not=float(np.median(x_not)),
                rbes=float(rbes))

def block_permutation_mean_diff(df, is_sig, n_perm=1000, block_size=5_000_000):
    """Block permutation test for mean density differences respecting spatial structure."""

    # Permutation test on mean TE_density, shuffling labels within scaffold blocks.
    # Build blocks
    blocks = []
    for sc, sub in df.groupby("CHROM"):
        sub = sub.sort_values("WIN_START").copy()
        starts = sub["WIN_START"].to_numpy()
        ends   = sub["WIN_END"].to_numpy()
        idx = sub.index.to_numpy()
        b_start = starts[0]
        b_idx = [idx[0]]
        for i in range(1, len(sub)):
            if (ends[i-1] - b_start + 1) >= block_size:
                blocks.append(b_idx)
                b_start = starts[i]
                b_idx = [idx[i]]
            else:
                b_idx.append(idx[i])
        if b_idx:
            blocks.append(b_idx)
    # observed diff
    obs = df.loc[is_sig, "TE_density"].mean() - df.loc[~is_sig, "TE_density"].mean()
    # permute
    greater = 0
    for _ in range(n_perm):
        perm_mask = is_sig.copy()
        for b in blocks:
            submask = perm_mask.loc[b].to_numpy()
            np.random.shuffle(submask)
            perm_mask.loc[b] = submask
        val = df.loc[perm_mask, "TE_density"].mean() - df.loc[~perm_mask, "TE_density"].mean()
        if abs(val) >= abs(obs): greater += 1
    p = (greater + 1) / (n_perm + 1)
    return dict(p_perm=float(p), mean_diff=float(obs), n_perm=int(n_perm), block_size=int(block_size))

# =========================
# Plot helpers
# =========================
def _alpha_from_q(qseries):
    """Convert q-values to alpha thresholds while clamping extremes."""

    return np.where(qseries < Q_CUTOFF, ALPHA_SIG, ALPHA_NONSIG)

def _axes_common_style(ax):
    """Apply shared styling tweaks used by Manhattan panels."""

    ax.grid(axis="y", alpha=0.25)
    ax.tick_params(labelsize=BASE_FONTSIZE-1)

def _scatter_by_scaffold(ax, df, color_a=COLOR_SCAF_A, color_b=COLOR_SCAF_B, y_col="Y", alpha=None):
    """Scatter helper alternating scaffold colours for readability."""

    scafs = sorted(df["CHROM"].unique(), key=scaffold_sort_key)
    for i, s in enumerate(scafs):
        sub = df[df["CHROM"]==s]
        col = color_a if (i % 2 == 0) else color_b
        a = alpha[sub.index] if isinstance(alpha, pd.Series) else alpha
        ax.scatter(sub["X_Gb"], sub[y_col], c=col, s=POINT_SIZE, alpha=a, edgecolors="none", zorder=2)

def manhattan_q_panels_by_feature(df_aug, feature_col, outdir, gwas_df=None):
    """Render multi-row Manhattan plots for a TE feature across comparisons."""

    comps = ["3vs4","3vs5","4vs5"]
    qmap = {c: f"q_poi_{c}" for c in comps}
    safe_feat = feature_col.replace("/", "_").replace(" ", "_")

    base = df_aug[["CHROM","WIN_START","WIN_END","WIN_MID", feature_col,
                   qmap["3vs4"], qmap["3vs5"], qmap["4vs5"]]].copy()
    base = base.rename(columns={feature_col: "COUNT_FEAT"})
    base, _ = build_genome_coords(base)
    for c in comps:
        base[f"neglog10_{c}"] = -np.log10(base[qmap[c]].astype(float).clip(lower=1e-300, upper=1.0))

    fig, axes = plt.subplots(3, 1, figsize=(14, PANEL_ROW_HEIGHT*3), sharex=True)
    for i, comp in enumerate(comps):
        ax = axes[i]
        qcol = qmap[comp]
        ycol = f"neglog10_{comp}"
        dfc = base.copy(); dfc["Y"] = dfc[ycol]
        alpha = pd.Series(_alpha_from_q(dfc[qcol]), index=dfc.index)
        _scatter_by_scaffold(ax, dfc, y_col="Y", alpha=alpha)

        # annotate top-K by combined rank (q small & count high)
        sig = dfc[dfc[qcol] < Q_CUTOFF].copy()
        if len(sig):
            sig["r_q"] = sig[qcol].rank(method="min", ascending=True)
            sig["r_c"] = (-sig["COUNT_FEAT"]).rank(method="min", ascending=True)
            sig["score"] = sig["r_q"] + sig["r_c"]
            topk = sig.nsmallest(TOP_K_ANNOTATE, "score")
            for _, r in topk.iterrows():
                lbl = f"{r['CHROM']}:{int(r['WIN_START'])}-{int(r['WIN_END'])} [{int(r['COUNT_FEAT'])}]"
                ax.scatter(r["X_Gb"], r["Y"], s=POINT_SIZE*1.4, c=COLOR_TOP_LABELS, zorder=5)
                ax.text(r["X_Gb"], r["Y"], lbl, fontsize=BASE_FONTSIZE-1, color=COLOR_TOP_LABELS,
                        rotation=45, ha="left", va="bottom")

        ax.set_ylabel(r"$-\log_{10}(q)$", fontsize=BASE_FONTSIZE)
        ax.set_title(f"{comp} — feature: {feature_col}", fontsize=BASE_FONTSIZE+1)
        _axes_common_style(ax)
    axes[-1].set_xlabel("Genome position (Gb, scaffolds concatenated)", fontsize=BASE_FONTSIZE)
    fig.tight_layout()
    outA = outdir / f"manhattan_q_panels_{safe_feat}.png"
    fig.savefig(outA, dpi=220); plt.close(fig)

    # optional GWAS-annotated panel
    outB = None
    if gwas_df is not None and len(gwas_df):
        gw_by_scaf = {s: sub["Pos"].astype(int).to_numpy() for s, sub in gwas_df.groupby("Scaffold")}
        gw_traits  = {s: sub.set_index("Pos")["traits"].astype(str) for s, sub in gwas_df.groupby("Scaffold")}
        fig2, axes2 = plt.subplots(3, 1, figsize=(14, PANEL_ROW_HEIGHT*3), sharex=True)
        for i, comp in enumerate(comps):
            ax = axes2[i]; qcol = qmap[comp]; ycol = f"neglog10_{comp}"
            dfc = base.copy(); dfc["Y"] = dfc[ycol]
            alpha = pd.Series(_alpha_from_q(dfc[qcol]), index=dfc.index)
            _scatter_by_scaffold(ax, dfc, y_col="Y", alpha=alpha)
            hits = []
            for _, r in dfc[dfc[qcol] < Q_CUTOFF].iterrows():
                sc = r["CHROM"]
                if sc not in gw_by_scaf: continue
                pos = gw_by_scaf[sc]
                ok  = (pos >= int(r["WIN_START"])) & (pos <= int(r["WIN_END"]))
                if np.any(ok):
                    first_pos = int(pos[ok][0])
                    trait     = gw_traits[sc].get(first_pos, "")
                    ax.scatter(r["X_Gb"], r["Y"], s=POINT_SIZE*1.4, c=COLOR_TOP_LABELS, zorder=5)
                    ax.text(r["X_Gb"], r["Y"], f"{first_pos}:{trait}", fontsize=BASE_FONTSIZE-1,
                            color=COLOR_TOP_LABELS, rotation=45, ha="left", va="bottom")
            ax.set_ylabel(r"$-\log_{10}(q)$", fontsize=BASE_FONTSIZE)
            ax.set_title(f"{comp} — feature: {feature_col} (GWAS-overlap)", fontsize=BASE_FONTSIZE+1)
            _axes_common_style(ax)
        axes2[-1].set_xlabel("Genome position (Gb, scaffolds concatenated)", fontsize=BASE_FONTSIZE)
        fig2.tight_layout()
        outB = outdir / f"manhattan_q_panels_{safe_feat}_GWASannot.png"
        fig2.savefig(outB, dpi=220); plt.close(fig2)

    return outA, outB

def manhattan_q_panels_GWAS_only(fst_df_with_coords, outdir, gwas_df):
    """Plot Manhattan panels restricted to GWAS hits."""

    comps = ["3vs4","3vs5","4vs5"]; qmap = {c: f"q_poi_{c}" for c in comps}
    base = fst_df_with_coords.copy()
    for c in comps:
        base[f"neglog10_{c}"] = -np.log10(base[qmap[c]].astype(float).clip(lower=1e-300, upper=1.0))
    gw_by_scaf = {}
    gw_traits  = {}
    if gwas_df is not None and len(gwas_df):
        gw_by_scaf = {s: sub["Pos"].astype(int).to_numpy() for s, sub in gwas_df.groupby("Scaffold")}
        gw_traits  = {s: sub.set_index("Pos")["traits"].astype(str) for s, sub in gwas_df.groupby("Scaffold")}
    fig, axes = plt.subplots(3, 1, figsize=(14, PANEL_ROW_HEIGHT*3), sharex=True)
    for i, comp in enumerate(comps):
        ax = axes[i]; qcol = qmap[comp]; ycol = f"neglog10_{comp}"
        dfc = base.copy(); dfc["Y"] = dfc[ycol]
        alpha = pd.Series(np.where(dfc[qcol] < Q_CUTOFF, ALPHA_SIG, ALPHA_NONSIG), index=dfc.index)
        _scatter_by_scaffold(ax, dfc, y_col="Y", alpha=alpha)
        if gw_by_scaf:
            for _, r in dfc[dfc[qcol] < Q_CUTOFF].iterrows():
                sc = r["CHROM"]
                if sc not in gw_by_scaf: continue
                pos = gw_by_scaf[sc]
                ok  = (pos >= int(r["WIN_START"])) & (pos <= int(r["WIN_END"]))
                if np.any(ok):
                    first_pos = int(pos[ok][0])
                    trait     = gw_traits[sc].get(first_pos, "")
                    ax.scatter(r["X_Gb"], r["Y"], s=POINT_SIZE*1.4, c=COLOR_TOP_LABELS, zorder=5)
                    ax.text(r["X_Gb"], r["Y"], f"{first_pos}:{trait}", fontsize=BASE_FONTSIZE-1,
                            color=COLOR_TOP_LABELS, rotation=45, ha="left", va="bottom")
        ax.set_ylabel(r"$-\log_{10}(q)$", fontsize=BASE_FONTSIZE)
        ax.set_title(f"{comp} — GWAS-only overlay of significant windows", fontsize=BASE_FONTSIZE+1)
        _axes_common_style(ax)
    axes[-1].set_xlabel("Genome position (Gb, scaffolds concatenated)", fontsize=BASE_FONTSIZE)
    fig.tight_layout()
    out = outdir / "manhattan_q_panels_GWAS_ONLY.png"
    fig.savefig(out, dpi=220); plt.close(fig)
    return out

def manhattan_counts_panel(df_aug, feature, outdir):
    """Plot Manhattan-style panels showing raw TE counts by window."""

    comps = ["3vs4","3vs5","4vs5"]; qmap = {c: f"q_poi_{c}" for c in comps}
    safe_feat = feature.replace("/", "_").replace(" ", "_")
    base = df_aug[["CHROM","WIN_START","WIN_END","WIN_MID", feature,
                   qmap["3vs4"], qmap["3vs5"], qmap["4vs5"]]].copy()
    base = base.rename(columns={feature:"Y"})
    base, _ = build_genome_coords(base)
    base["Y"] = pd.to_numeric(base["Y"], errors="coerce").fillna(0)
    fig, axes = plt.subplots(3, 1, figsize=(14, PANEL_ROW_HEIGHT*3), sharex=True)
    for i, comp in enumerate(comps):
        ax = axes[i]; qcol = qmap[comp]
        alpha = pd.Series(_alpha_from_q(base[qcol]), index=base.index)
        _scatter_by_scaffold(ax, base, y_col="Y", alpha=alpha)
        ax.set_ylabel(f"{feature} per window", fontsize=BASE_FONTSIZE)
        ax.set_title(f"{comp} — counts per window", fontsize=BASE_FONTSIZE+1)
        _axes_common_style(ax)
    axes[-1].set_xlabel("Genome position (Gb, scaffolds concatenated)", fontsize=BASE_FONTSIZE)
    fig.tight_layout()
    out = outdir / f"manhattan_counts_{safe_feat}.png"
    fig.savefig(out, dpi=220); plt.close(fig)
    return out

# New diagnostics figures
def plot_dispersion_bars(disp_df, outpath):
    """Plot bar charts summarising dispersion ratios for GLM fits."""

    fig, ax = plt.subplots(figsize=(7, 4))
    order = ["3vs4","3vs5","4vs5"]
    vals = [disp_df.loc[disp_df["comparison"]==c, "dispersion"].values[0] for c in order]
    ax.bar(order, vals, color=[COLOR_SCAF_A,COLOR_SCAF_B,COLOR_SCAF_A])
    ax.axhline(1.0, ls="--", color="k", lw=1)
    ax.set_ylabel("Poisson dispersion (Pearson χ² / df)", fontsize=BASE_FONTSIZE)
    ax.set_title("Dispersion by comparison (Poisson GLM)", fontsize=BASE_FONTSIZE+1)
    _axes_common_style(ax)
    fig.tight_layout()
    fig.savefig(outpath, dpi=220); plt.close(fig)

def plot_density_violins(df_aug, comp, outpath):
    """Create violin/box plots comparing TE densities across significance groups."""

    mask = (df_aug[f"q_poi_{comp}"] < Q_CUTOFF)
    dfp = pd.DataFrame({
        "group": np.where(mask, "significant", "not_significant"),
        "density": df_aug["TE_density"].astype(float)
    })
    # simple violin+box using matplotlib
    fig, ax = plt.subplots(figsize=(7, 4))
    data = [dfp.loc[dfp["group"]=="significant","density"].to_numpy(),
            dfp.loc[dfp["group"]=="not_significant","density"].to_numpy()]
    parts = ax.violinplot(data, showmedians=False, showextrema=False)
    for pc in parts['bodies']:
        pc.set_facecolor(COLOR_SIG); pc.set_alpha(0.3)
    ax.boxplot(data, positions=[1,2], widths=0.2, vert=True,
               boxprops=dict(color="black"), medianprops=dict(color="black"),
               whiskerprops=dict(color="black"), capprops=dict(color="black"))
    ax.set_xticks([1,2], ["significant","not"])
    ax.set_ylabel("TE density (TE/base)", fontsize=BASE_FONTSIZE)
    ax.set_title(f"{comp}: TE density by group", fontsize=BASE_FONTSIZE+1)
    _axes_common_style(ax)
    fig.tight_layout()
    fig.savefig(outpath, dpi=220); plt.close(fig)

def plot_count_histograms(df_aug, comp, outpath):
    """Plot histograms of TE counts for significant vs non-significant windows."""

    mask = (df_aug[f"q_poi_{comp}"] < Q_CUTOFF)
    fig, ax = plt.subplots(figsize=(7, 4))
    ax.hist(df_aug.loc[ mask,"TE_count"], bins=50, alpha=0.6, color=COLOR_SIG, label="significant")
    ax.hist(df_aug.loc[~mask,"TE_count"], bins=50, alpha=0.6, color=COLOR_NONSIG, label="not")
    ax.set_yscale("log")
    ax.set_xlabel("TE count per window", fontsize=BASE_FONTSIZE)
    ax.set_ylabel("Frequency (log scale)", fontsize=BASE_FONTSIZE)
    ax.set_title(f"{comp}: TE count distribution", fontsize=BASE_FONTSIZE+1)
    ax.legend()
    _axes_common_style(ax)
    fig.tight_layout()
    fig.savefig(outpath, dpi=220); plt.close(fig)

# =========================
# NEW: Sig vs Not presence & summaries (kept from v8 and reused)
# =========================
def _sig_mask(df, comp, q_cutoff=Q_CUTOFF):
    """Return a boolean mask for windows significant for the comparison."""

    return (df[f"q_poi_{comp}"] < q_cutoff)

def _presence_counts(df, mask_sig, te_col="TE_count"):
    """Count TE presence/absence separated by significance mask."""

    sig_with    = int((df.loc[ mask_sig, te_col] >= 1).sum())
    sig_without = int((df.loc[ mask_sig, te_col] == 0).sum())
    non_with    = int((df.loc[~mask_sig, te_col] >= 1).sum())
    non_without = int((df.loc[~mask_sig, te_col] == 0).sum())
    return sig_with, sig_without, non_with, non_without

def _desc(df, mask):
    """Compute descriptive statistics for TE columns under a mask."""

    x = df.loc[mask, "TE_count"].astype(int)
    return dict(
        n_windows     = int(mask.sum()),
        total_TE      = int(x.sum()),
        mean_TE       = float(x.mean()) if len(x) else 0.0,
        median_TE     = float(x.median()) if len(x) else 0.0,
        p90_TE        = float(x.quantile(0.90)) if len(x) else 0.0,
        max_TE        = int(x.max()) if len(x) else 0,
    )

def write_sig_te_summaries(df_aug, outdir, q_cutoff=Q_CUTOFF):
    """Export summary TSVs describing TE presence and densities by significance."""

    outdir = Path(outdir); outdir.mkdir(parents=True, exist_ok=True)
    comps = ["3vs4","3vs5","4vs5"]

    rowsA = []
    for comp in comps:
        m = _sig_mask(df_aug, comp, q_cutoff=q_cutoff)
        sw, s0, nw, n0 = _presence_counts(df_aug, m, te_col="TE_count")
        rowsA.append(dict(
            comparison=comp,
            q_cutoff=q_cutoff,
            n_sig=int(m.sum()),
            n_nonsig=int((~m).sum()),
            sig_with_TE=sw,   sig_without_TE=s0,
            nonsig_with_TE=nw, nonsig_without_TE=n0
        ))
    pd.DataFrame(rowsA).to_csv(outdir / "te_sig_presence_overview.tsv", sep="\t", index=False)

    rowsB = []
    for comp in comps:
        m = _sig_mask(df_aug, comp, q_cutoff=q_cutoff)
        d_sig = _desc(df_aug, m); d_not = _desc(df_aug, ~m)
        rowsB.append(dict(comparison=comp, group="significant", **d_sig))
        rowsB.append(dict(comparison=comp, group="not_significant", **d_not))
    pd.DataFrame(rowsB).to_csv(outdir / "te_sig_descriptives.tsv", sep="\t", index=False)

    te_class_cols = [c for c in df_aug.columns if c.startswith("TE_") and c not in ("TE_count",)]
    if te_class_cols:
        rowsC = []
        for comp in comps:
            m = _sig_mask(df_aug, comp, q_cutoff=q_cutoff)
            for col in te_class_cols:
                sw, s0, nw, n0 = _presence_counts(df_aug, m, te_col=col)
                rowsC.append(dict(
                    comparison=comp, te_class=col,
                    sig_with_class=sw, sig_without_class=s0,
                    nonsig_with_class=nw, nonsig_without_class=n0
                ))
        pd.DataFrame(rowsC).to_csv(outdir / "te_sig_presence_by_class.tsv", sep="\t", index=False)

# =========================
# MAIN
# =========================
def main():
    # 1) Load
    fst  = read_fst_windows(FST_CSV)
    te   = read_te_table(TE_TSV)
    gwas = read_gwas(GWAS_XLSX)

    # 2) TE overlaps → counts/densities
    te_counts, te_types = count_te_overlaps_per_window(fst, te)
    df_aug = fst.merge(te_counts, on=["CHROM","WIN_START","WIN_END","WIN_LEN","WIN_MID"], how="left")
    num_cols = df_aug.select_dtypes(include=[np.number]).columns
    df_aug[num_cols] = df_aug[num_cols].fillna(0)

    # 3) GWAS overlap columns
    df_aug = add_gwas_overlap_columns(df_aug, gwas)

    # 4) PRESENCE tests (overall + each class) — kept/enriched
    comps = ["3vs4","3vs5","4vs5"]
    comp_masks = {c: (df_aug[f"q_poi_{c}"] < Q_CUTOFF) for c in comps}
    rows_presence = []
    for comp, is_sig in comp_masks.items():
        pres = presence_tests(df_aug, is_sig, "TE_count")
        rows_presence.append({"comparison":comp,"feature":"TE_total",
                              "n_windows_sig":int(is_sig.sum()),"n_windows_not":int((~is_sig).sum()),
                              "total_TE_sig":int(df_aug.loc[is_sig,"TE_count"].sum()),
                              "total_TE_not":int(df_aug.loc[~is_sig,"TE_count"].sum()), **pres})
        for t in te_types:
            col = f"TE_{t}"
            pre_t = presence_tests(df_aug, is_sig, col)
            rows_presence.append({"comparison":comp,"feature":col,
                                  "n_windows_sig":int(is_sig.sum()),"n_windows_not":int((~is_sig).sum()),
                                  "total_TE_sig":int(df_aug.loc[is_sig,col].sum()),
                                  "total_TE_not":int(df_aug.loc[~is_sig,col].sum()), **pre_t})
    pd.DataFrame(rows_presence).to_csv(OUTDIR / "presence_fisher_chisq.tsv", sep="\t", index=False)

    # 5) GLMs (Poisson with dispersion + NB fallback), and diagnostics TSV/plots
    glm_p_rows, glm_nb_rows, disp_rows = [], [], []
    for comp in comps:
        is_sig = comp_masks[comp]
        y = df_aug["TE_count"].astype(int).to_numpy()
        g = np.where(is_sig, 1.0, 0.0)
        offset = np.log(df_aug["WIN_LEN"].astype(float).to_numpy())
        res = glm_compare_counts(y, g, offset_log=offset, prefer_nb_if_overdispersed=True)

        # --- record poisson (pure) results & dispersion
        b  = res["res_p"].params.get("group", np.nan)
        se = res["res_p"].bse.get("group", np.nan)
        z  = (b / se) if (isinstance(se, float) and se>0) else np.nan
        p  = 2 * (1 - norm.cdf(abs(z))) if np.isfinite(z) else np.nan
        rr = math.exp(b) if np.isfinite(b) else np.nan
        ci_lo = clamp_exp(b - 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan
        ci_hi = clamp_exp(b + 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan
        aic_p = res["res_p"].aic if hasattr(res["res_p"], "aic") else np.nan
        glm_p_rows.append(dict(comparison=comp, family="Poisson", RR=rr, RR_CI_lo=ci_lo, RR_CI_hi=ci_hi,
                               z=z, p=p, AIC=aic_p))
        disp_rows.append(dict(comparison=comp, dispersion=float(res["dispersion"])))

        # --- record used model (Poisson or NB fallback)
        aic_u = res["res"].aic if hasattr(res["res"], "aic") else np.nan
        glm_nb_rows.append(dict(comparison=comp, family=res["family"], RR=res["RR"],
                                RR_CI_lo=res["RR_CI_lo"], RR_CI_hi=res["RR_CI_hi"],
                                z=res["z"], p=res["p"], AIC=aic_u))

        # Figures per comparison
        plot_density_violins(df_aug, comp, OUTDIR / f"violin_density_{comp}.png")
        plot_count_histograms(df_aug, comp, OUTDIR / f"hist_counts_{comp}.png")

    pd.DataFrame(glm_p_rows).to_csv(OUTDIR / "glm_poisson_results.tsv", sep="\t", index=False)
    pd.DataFrame(glm_nb_rows).to_csv(OUTDIR / "glm_nb_results.tsv", sep="\t", index=False)
    disp_df = pd.DataFrame(disp_rows)
    disp_df.to_csv(OUTDIR / "dispersion_summary.tsv", sep="\t", index=False)
    plot_dispersion_bars(disp_df, OUTDIR / "dispersion_bars.png")

    # 6) Wilcoxon on TE density (robustness); BH-FDR per comparison family (overall + per class optional)
    wil_rows = []
    p_vals = []
    # overall TE density
    for comp in comps:
        is_sig = comp_masks[comp]
        resw = wilcoxon_density(df_aug, is_sig)
        wil_rows.append(dict(comparison=comp, feature="TE_density", **resw))
        p_vals.append(resw["p"])
    # adjust within the 3 overall tests
    _, qvals, _, _ = multipletests([r["p"] for r in wil_rows if not np.isnan(r["p"])], method="fdr_bh")
    q_iter = iter(qvals)
    for r in wil_rows:
        r["q_bh"] = float(next(q_iter)) if not np.isnan(r["p"]) else np.nan
    pd.DataFrame(wil_rows).to_csv(OUTDIR / "wilcoxon_density.tsv", sep="\t", index=False)

    # 7) Optional: Zero-inflated models (Poisson/NB) to compare AIC if extreme zeros
    if RUN_ZERO_INFLATED and HAS_ZI:
        zi_rows = []
        for comp in comps:
            is_sig = comp_masks[comp].astype(int).to_numpy()
            X = sm.add_constant(is_sig)
            y = df_aug["TE_count"].astype(int).to_numpy()
            exog_infl = X  # same inflation regressor
            # ZIP
            try:
                zip_mod = ZeroInflatedPoisson(endog=y, exog=X, exog_infl=exog_infl, inflation='logit').fit(disp=False)
                zi_rows.append(dict(comparison=comp, model="ZIP", AIC=float(zip_mod.aic),
                                    ll=float(zip_mod.llf)))
            except Exception as e:
                zi_rows.append(dict(comparison=comp, model="ZIP", AIC=np.nan, ll=np.nan))
            # ZINB
            try:
                zinb_mod = ZeroInflatedNegativeBinomialP(endog=y, exog=X, exog_infl=exog_infl, inflation='logit').fit(disp=False)
                zi_rows.append(dict(comparison=comp, model="ZINB", AIC=float(zinb_mod.aic),
                                    ll=float(zinb_mod.llf)))
            except Exception as e:
                zi_rows.append(dict(comparison=comp, model="ZINB", AIC=np.nan, ll=np.nan))
        pd.DataFrame(zi_rows).to_csv(OUTDIR / "zi_model_summary.tsv", sep="\t", index=False)

    # 8) Optional: Block permutation on mean TE density
    if RUN_PERMUTATION:
        perm_rows = []
        for comp in comps:
            is_sig = comp_masks[comp]
            resperm = block_permutation_mean_diff(df_aug, is_sig,
                                                  n_perm=N_PERM, block_size=BLOCK_SIZE)
            perm_rows.append(dict(comparison=comp, **resperm))
        pd.DataFrame(perm_rows).to_csv(OUTDIR / "permutation_summary.tsv", sep="\t", index=False)

    # 9) Manhattan families (unchanged from v8)
    paths = []
    a, b = manhattan_q_panels_by_feature(df_aug, "TE_count", OUTDIR, gwas_df=gwas)
    paths += [p for p in (a,b) if p is not None]
    for t in te_types:
        feat = f"TE_{t}"
        a, b = manhattan_q_panels_by_feature(df_aug, feat, OUTDIR, gwas_df=gwas)
        paths += [p for p in (a,b) if p is not None]
    fst_coords, _ = build_genome_coords(df_aug[["CHROM","WIN_START","WIN_END","WIN_MID",
                                                "q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]])
    paths.append(manhattan_q_panels_GWAS_only(fst_coords, OUTDIR, gwas_df=gwas))
    paths.append(manhattan_counts_panel(df_aug, "TE_count", OUTDIR))
    for t in te_types:
        paths.append(manhattan_counts_panel(df_aug, f"TE_{t}", OUTDIR))

    # 10) Sig/not presence summaries (v8)
    write_sig_te_summaries(df_aug, OUTDIR, q_cutoff=Q_CUTOFF)

    # 11) Save augmented TSV
    out_aug = OUTDIR / "fst_windows_with_TEcounts.tsv"
    df_aug.to_csv(out_aug, sep="\t", index=False)

    # 12) Build PPTX report
    prs = Presentation()
    # Title
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "TE vs FST — v9 Complete Analysis"
    s.placeholders[1].text = (
        f"Q cutoff = {Q_CUTOFF}\n"
        "- Poisson GLM with log-exposure, dispersion diagnostics\n"
        "- Negative Binomial fallback (RR, CI)\n"
        "- Wilcoxon on TE density; Fisher/Chi² presence\n"
        f"- Optional: Zero-inflated & block permutation\n"
        f"Outputs in: {OUTDIR.resolve()}"
    )
    # Dispersion & GLM slides
    for fig in ["dispersion_bars.png",
                "violin_density_3vs4.png","hist_counts_3vs4.png",
                "violin_density_3vs5.png","hist_counts_3vs5.png",
                "violin_density_4vs5.png","hist_counts_4vs5.png"]:
        if (OUTDIR/fig).exists():
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = fig.replace("_"," ").replace(".png","")
            slide.shapes.add_picture(str(OUTDIR/fig), Inches(0.5), Inches(1.2), height=Inches(5.2))
    # Manhattan examples
    for p in paths[:6]:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = Path(p).name
        slide.shapes.add_picture(str(p), Inches(0.5), Inches(1.2), height=Inches(5.2))
    prs.save(str(OUTDIR / "TE_vs_FST_Report_v9.pptx"))

    # 13) Write short console summary
    print("\nDone v9.")
    print(f"- Augmented TSV: {out_aug.name}")
    print(f"- GLM Poisson results: glm_poisson_results.tsv")
    print(f"- GLM NB (or used) results: glm_nb_results.tsv")
    print(f"- Dispersion summary & bars: dispersion_summary.tsv, dispersion_bars.png")
    print(f"- Wilcoxon density: wilcoxon_density.tsv")
    print(f"- Presence tests: presence_fisher_chisq.tsv")
    if RUN_ZERO_INFLATED and HAS_ZI:
        print(f"- ZI model summary: zi_model_summary.tsv")
    if RUN_PERMUTATION:
        print(f"- Permutation summary: permutation_summary.tsv")
    print(f"- Manhattan figures & PPTX in: {OUTDIR.resolve()}")

if __name__ == "__main__":
    main()

