#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
TE overlap & density vs FST-enriched windows

Inputs:
  - FST window CSV (same folder as this script):
      fst_window_counts_tests_win10000000_ov2000000.csv
    Columns (exact list provided by user), we only require:
      CHROM, WIN_START, WIN_END, WIN_MID, and q_poi_3vs4, q_poi_3vs5, q_poi_4vs5

  - TE TSV:
      ~/Desktop/OSU_projects/conifers/LP/lodgepole_pine_assembly/col8_readable.TEanno.cds.scaf1ab.tsv
    Columns (tab-separated):
      seqid, source, sequence_ontology, start, end, score, strand

Outputs (in ./te_vs_fst_out/):
  - fst_windows_with_TEcounts.tsv  (augmented table: TE counts & densities)
  - te_density_group_summary.tsv   (group-level stats & model summaries)
  - te_density_by_type.tsv         (per-TE-type model results with FDR)
  - figures: violin+box, per-type rate-ratio (forest)
  - PowerPoint: TE_vs_FST_report.pptx
"""

import os
import math
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

import statsmodels.api as sm
from statsmodels.api import GLM, families
from statsmodels.genmod.families import links
from statsmodels.stats.multitest import multipletests
from scipy.stats import norm
from pptx import Presentation
from pptx.util import Inches

# =========================
# User knobs (edit here)
# =========================
FST_CSV = "fst_window_counts_tests_win10000000_ov2000000.csv"
TE_TSV  = os.path.expanduser(
    "~/Desktop/OSU_projects/conifers/LP/lodgepole_pine_assembly/col8_readable.TEanno.cds.scaf1ab.tsv"
)

# Which Poisson q-values define “significant FST windows”?
# mode = "any"  -> significant if any of (q_poi_3vs4, q_poi_3vs5, q_poi_4vs5) < 0.05
# mode = "all"  -> significant if all three < 0.05
# mode = "3vs4" / "3vs5" / "4vs5" -> use only that column
SIGNIFICANCE_MODE = "any"
Q_CUTOFF = 0.05

# Output folder
OUTDIR = Path("te_vs_fst_out")
OUTDIR.mkdir(parents=True, exist_ok=True)

# Random seed for reproducible jitter in plots
np.random.seed(1)

# =========================
# Helpers
# =========================
def read_fst_windows(path: str) -> pd.DataFrame:
    df = pd.read_csv(path)
    # Only rely on the exact columns you listed; we need these:
    need = ["CHROM","WIN_START","WIN_END","WIN_MID","q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]
    missing = [c for c in need if c not in df.columns]
    if missing:
        raise ValueError(f"FST file missing required column(s): {missing}")

    # ensure integer coords
    df["CHROM"]     = df["CHROM"].astype(str)
    df["WIN_START"] = pd.to_numeric(df["WIN_START"], errors="coerce").astype("Int64")
    df["WIN_END"]   = pd.to_numeric(df["WIN_END"],   errors="coerce").astype("Int64")
    df["WIN_MID"]   = pd.to_numeric(df["WIN_MID"],   errors="coerce").astype("Int64")
    # window length
    df["WIN_LEN"] = (df["WIN_END"] - df["WIN_START"] + 1).astype("Int64")
    return df

def _has_header(path: str, first_col_name: str) -> bool:
    with open(path, "r") as fh:
        first = fh.readline().strip().split("\t")[0]
        return first == first_col_name

def read_te_table(path: str) -> pd.DataFrame:
    cols = ["seqid","source","sequence_ontology","start","end","score","strand"]
    te = pd.read_csv(path, sep="\t", names=cols,
                     header=0 if _has_header(path, cols[0]) else None)
    # enforce types
    te["seqid"] = te["seqid"].astype(str)
    te["sequence_ontology"] = te["sequence_ontology"].astype(str)
    te["start"] = pd.to_numeric(te["start"], errors="coerce").astype("Int64")
    te["end"]   = pd.to_numeric(te["end"],   errors="coerce").astype("Int64")
    te = te.dropna(subset=["seqid","start","end"]).copy()
    te = te.sort_values(["seqid","start","end"]).reset_index(drop=True)
    return te

def _count_overlaps_one_chrom(wins: pd.DataFrame, tes: pd.DataFrame, te_types: np.ndarray):
    """Count overlaps per window for one chromosome. Two-pointer scan."""
    w_st = wins["WIN_START"].to_numpy(dtype=np.int64)
    w_en = wins["WIN_END"].to_numpy(dtype=np.int64)
    nW   = len(wins)

    total_counts = np.zeros(nW, dtype=np.int64)
    type_to_idx = {t:i for i,t in enumerate(te_types)}
    per_type = np.zeros((nW, len(te_types)), dtype=np.int64)

    if tes.empty:
        return total_counts, per_type

    t_st = tes["start"].to_numpy(dtype=np.int64)
    t_en = tes["end"].to_numpy(dtype=np.int64)
    t_ty = tes["sequence_ontology"].astype(str).to_numpy()

    j = 0  # TE pointer
    for i in range(nW):
        ws, we = w_st[i], w_en[i]
        while j < len(t_st) and t_en[j] < ws:
            j += 1
        k = j
        while k < len(t_st) and t_st[k] <= we:
            if not (t_en[k] < ws or t_st[k] > we):  # overlap
                total_counts[i] += 1
                idx = type_to_idx.get(t_ty[k], None)
                if idx is not None:
                    per_type[i, idx] += 1
            k += 1
    return total_counts, per_type

def count_te_overlaps_per_window(fst_df: pd.DataFrame, te_df: pd.DataFrame):
    """
    Returns: (out_df, te_types)
      out_df columns:
        CHROM, WIN_START, WIN_END, WIN_LEN, TE_count, TE_<type>..., TEdens_<type>..., TE_density
      Built with a single concat to avoid fragmentation.
    """
    te_types = np.array(sorted(te_df["sequence_ontology"].astype(str).unique()))
    base = fst_df[["CHROM","WIN_START","WIN_END","WIN_LEN"]].copy().reset_index(drop=True)

    total_counts_all = np.zeros(len(base), dtype=np.int64)
    per_type_all = np.zeros((len(base), len(te_types)), dtype=np.int64)

    for chrom, wins in fst_df.groupby("CHROM", sort=False):
        rows = wins.index.to_numpy()
        wins_sorted = wins.sort_values("WIN_START")
        # Map sorted back to original order
        order_back = np.argsort(wins_sorted.index.to_numpy())
        tes = te_df[te_df["seqid"] == chrom].sort_values("start")

        total_counts, per_type = _count_overlaps_one_chrom(wins_sorted, tes, te_types)
        total_counts_all[rows] = total_counts[order_back]
        per_type_all[rows, :]  = per_type[order_back, :]

    te_cols = {"TE_count": total_counts_all.astype(np.int64)}
    for j, t in enumerate(te_types):
        te_cols[f"TE_{t}"] = per_type_all[:, j].astype(np.int64)

    te_df_out = pd.DataFrame(te_cols)
    out = pd.concat([base, te_df_out], axis=1)

    # densities
    win_len = out["WIN_LEN"].astype(float)
    out["TE_density"] = out["TE_count"] / win_len
    for t in te_types:
        out[f"TEdens_{t}"] = out[f"TE_{t}"] / win_len

    return out, te_types

def classify_significance(df: pd.DataFrame, mode="any", q_cutoff=0.05) -> pd.Series:
    cols = ["q_poi_3vs4","q_poi_3vs5","q_poi_4vs5"]
    for c in cols:
        if c not in df.columns:
            raise ValueError(f"Expected Poisson q-value column missing: {c}")
    if mode == "any":
        sig = (df[cols] < q_cutoff).any(axis=1)
    elif mode == "all":
        sig = (df[cols] < q_cutoff).all(axis=1)
    elif mode in ("3vs4","3vs5","4vs5"):
        col = f"q_poi_{mode}"
        sig = df[col] < q_cutoff
    else:
        raise ValueError("SIGNIFICANCE_MODE must be 'any', 'all', or one of '3vs4','3vs5','4vs5'.")
    return sig

def dispersion_ratio(y, mu, df_adj=1):
    """Pearson dispersion ≈ sum((y-mu)^2/mu)/df ; df = max(n - df_adj, 1)."""
    y = np.asarray(y, float)
    mu = np.asarray(mu, float)
    ok = mu > 0
    if ok.sum() <= 1:
        return np.nan
    return float(np.sum(((y[ok]-mu[ok])**2)/mu[ok]) / max(ok.sum()-df_adj, 1))

def clamp_exp(x):
    """Safe exp for CI to avoid OverflowError."""
    return float(np.exp(np.clip(x, -700, 700)))

def fit_glm_poisson(y, X, offset):
    fam = families.Poisson(link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    res = model.fit()
    return res

def estimate_nb_alpha(y, mu):
    """
    Method-of-moments estimate of NB alpha from Poisson fit:
      Var ≈ mu + alpha * mu^2  -> alpha ≈ sum((y - mu)^2 - mu) / sum(mu^2)
    """
    y = np.asarray(y, float)
    mu = np.asarray(mu, float)
    numer = np.sum((y - mu) ** 2 - mu)
    denom = np.sum(mu ** 2)
    if denom <= 0:
        return 1.0
    alpha = numer / denom
    return float(max(alpha, 1e-8))

def fit_glm_nb(y, X, offset, alpha):
    fam = families.NegativeBinomial(alpha=alpha, link=links.Log())
    model = GLM(y, X, family=fam, offset=offset)
    res = model.fit()
    return res

def glm_compare_counts(y, group, offset_log, prefer_nb_if_overdispersed=True):
    """
    GLM: count ~ group (0/1), offset = log(exposure)
    Returns dict(coef,se,z,p,RR,CI,dispersion,family,res).
    If overdispersed, switch to NB with alpha estimated from Poisson mean.
    """
    X = pd.DataFrame({"intercept": 1.0, "group": group.astype(float)})

    # Poisson fit
    res_p = fit_glm_poisson(y, X, offset_log)
    mu_p  = res_p.fittedvalues
    disp  = dispersion_ratio(y, mu_p)

    res_used = res_p
    family_used = "Poisson"

    if prefer_nb_if_overdispersed and np.isfinite(disp) and disp > 1.5:
        alpha = estimate_nb_alpha(y, mu_p)
        res_nb = fit_glm_nb(y, X, offset_log, alpha)
        res_used = res_nb
        family_used = f"NegBin(alpha={alpha:.3g})"

    b = res_used.params.get("group", np.nan)
    se = res_used.bse.get("group", np.nan)
    z  = (b / se) if se and np.isfinite(se) and se > 0 else np.nan
    p  = 2 * (1 - norm.cdf(abs(z))) if np.isfinite(z) else np.nan
    rr = math.exp(b) if np.isfinite(b) else np.nan
    ci_lo = clamp_exp(b - 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan
    ci_hi = clamp_exp(b + 1.96*se) if np.isfinite(b) and np.isfinite(se) else np.nan

    return {"family": family_used, "coef": b, "se": se, "z": z, "p": p,
            "RR": rr, "RR_CI_lo": ci_lo, "RR_CI_hi": ci_hi,
            "dispersion": disp, "res": res_used}

def global_rate_ratio_test(k_sig, e_sig, k_nsig, e_nsig):
    """
    Simple Poisson approximate comparison of two rates using counts & exposures.
    """
    lam1 = k_sig / e_sig if e_sig > 0 else np.nan
    lam0 = k_nsig / e_nsig if e_nsig > 0 else np.nan
    rr = lam1/lam0 if (lam0 and lam1) else np.nan
    if k_sig > 0 and k_nsig > 0 and np.isfinite(rr):
        se = math.sqrt(1.0/k_sig + 1.0/k_nsig)
        z  = math.log(rr)/se
        p  = 2*(1 - norm.cdf(abs(z)))
        ci_lo = clamp_exp(math.log(rr) - 1.96*se)
        ci_hi = clamp_exp(math.log(rr) + 1.96*se)
    else:
        p = np.nan; ci_lo = np.nan; ci_hi = np.nan
    return rr, ci_lo, ci_hi, p

def box_and_violin(df_aug: pd.DataFrame, outdir: Path):
    fig, ax = plt.subplots(1,1, figsize=(6,4))
    groups = ["not significant","significant"]
    data = [df_aug.loc[~df_aug["is_sig"], "TE_density"],
            df_aug.loc[df_aug["is_sig"], "TE_density"]]
    ax.violinplot(data, positions=[0,1], showmeans=True, widths=0.8)
    ax.boxplot(data, positions=[0,1], widths=0.2, showcaps=True, showfliers=False)
    ax.set_xticks([0,1], groups)
    ax.set_ylabel("TE density (TEs / bp)")
    ax.set_title("TE density by FST significance group (Poisson q)")
    ax.grid(alpha=0.2, axis="y")
    out = outdir / "te_density_violin_box.png"
    fig.tight_layout(); fig.savefig(out, dpi=200); plt.close(fig)
    return out

def forest_rate_ratio(df_types: pd.DataFrame, outdir: Path):
    d = df_types.sort_values("RR", ascending=False).copy()
    y = np.arange(len(d))[::-1]
    fig, ax = plt.subplots(figsize=(7, 0.35*max(6, len(d))))
    ax.errorbar(d["RR"], y,
                xerr=[d["RR"]-d["RR_CI_lo"], d["RR_CI_hi"]-d["RR"]],
                fmt='o', capsize=3)
    ax.axvline(1.0, color="#999999", lw=1, ls="--")
    ax.set_yticks(y, d["TE_type"])
    ax.set_xscale("log")
    ax.set_xlabel("Rate ratio (significant / not)   [log scale]")
    ax.set_title("Per-TE-type enrichment in significant windows")
    ax.grid(alpha=0.2, axis="x")
    out = outdir / "te_type_rate_ratio_forest.png"
    fig.tight_layout(); fig.savefig(out, dpi=220); plt.close(fig)
    return out

def add_slide(prs, title, bullets=None, image=None):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if bullets:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.6))
        tf = tb.text_frame; tf.clear()
        for i,b in enumerate(bullets):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.text = b; p.level = 0
    if image and Path(image).exists():
        slide.shapes.add_picture(str(image), Inches(0.5), Inches(2.7), height=Inches(3.8))

# =========================
# Main
# =========================
def main():
    # ---- load
    fst = read_fst_windows(FST_CSV)
    te  = read_te_table(TE_TSV)

    # ---- TE overlaps (build with a single concat to avoid fragmentation)
    te_counts, te_types = count_te_overlaps_per_window(fst, te)

    # merge back with FST table
    df_aug = fst.merge(te_counts, on=["CHROM","WIN_START","WIN_END","WIN_LEN"], how="left")
    # safe numeric fill for any missing
    num_cols = df_aug.select_dtypes(include=[np.number]).columns
    df_aug[num_cols] = df_aug[num_cols].fillna(0)

    # ---- classify windows by Poisson q-values (concat to avoid fragmentation warning)
    is_sig = classify_significance(df_aug, mode=SIGNIFICANCE_MODE, q_cutoff=Q_CUTOFF).rename("is_sig")
    df_aug = pd.concat([df_aug, is_sig], axis=1)

    # ---- overall TE density model
    y = df_aug["TE_count"].astype(int).to_numpy()
    group = df_aug["is_sig"].astype(int).to_numpy()
    offset = np.log(df_aug["WIN_LEN"].astype(float).to_numpy())

    res_used = glm_compare_counts(y, group, offset)
    model_used = res_used["family"]

    # global rate-ratio summary
    k_sig   = int(df_aug.loc[df_aug["is_sig"], "TE_count"].sum())
    e_sig   = float(df_aug.loc[df_aug["is_sig"], "WIN_LEN"].sum())
    k_nsig  = int(df_aug.loc[~df_aug["is_sig"], "TE_count"].sum())
    e_nsig  = float(df_aug.loc[~df_aug["is_sig"], "WIN_LEN"].sum())
    rr_glob, ci_lo_glob, ci_hi_glob, p_glob = global_rate_ratio_test(k_sig, e_sig, k_nsig, e_nsig)

    # ---- per-TE-type models; BH-FDR across types
    rows = []
    for t in te_types:
        y_t = df_aug[f"TE_{t}"].astype(int).to_numpy()
        res_t = glm_compare_counts(y_t, group, offset)
        rows.append({
            "TE_type": t,
            "family": res_t["family"],
            "RR": res_t["RR"], "RR_CI_lo": res_t["RR_CI_lo"], "RR_CI_hi": res_t["RR_CI_hi"],
            "p": res_t["p"], "dispersion": res_t["dispersion"]
        })
    df_types = pd.DataFrame(rows)
    df_types["q"] = multipletests(df_types["p"].fillna(1.0), method="fdr_bh")[1]
    df_types = df_types.sort_values("q")

    # ---- write outputs
    OUTDIR.mkdir(exist_ok=True, parents=True)
    out_aug = OUTDIR / "fst_windows_with_TEcounts.tsv"
    df_aug.to_csv(out_aug, sep="\t", index=False)

    out_types = OUTDIR / "te_density_by_type.tsv"
    df_types.to_csv(out_types, sep="\t", index=False)

    # group-level summary table
    summ = []
    for flag, lab in [(True,"significant"), (False,"not significant")]:
        sub = df_aug[df_aug["is_sig"] == flag]
        summ.append({
            "group": lab,
            "n_windows": len(sub),
            "total_TE": int(sub["TE_count"].sum()),
            "total_len_bp": int(sub["WIN_LEN"].sum()),
            "mean_density": float((sub["TE_count"]/sub["WIN_LEN"]).mean()),
            "median_density": float((sub["TE_count"]/sub["WIN_LEN"]).median())
        })
    df_summary = pd.DataFrame(summ)
    df_summary_overall = pd.DataFrame([{
        "model_used": model_used,
        "RR_overall": res_used["RR"], "RR_CI_lo": res_used["RR_CI_lo"],
        "RR_CI_hi": res_used["RR_CI_hi"], "p_model": res_used["p"],
        "dispersion_est": res_used["dispersion"],
        "global_RR": rr_glob, "global_CI_lo": ci_lo_glob, "global_CI_hi": ci_hi_glob, "p_global": p_glob
    }])

    out_summ = OUTDIR / "te_density_group_summary.tsv"
    with open(out_summ, "w") as fh:
        fh.write("# Group-wise TE density summary\n")
    df_summary.to_csv(out_summ, sep="\t", mode="a", index=False)
    with open(out_summ, "a") as fh:
        fh.write("\n# Overall model & global test\n")
    df_summary_overall.to_csv(out_summ, sep="\t", mode="a", index=False)

    # ---- plots
    fig_box = box_and_violin(df_aug, OUTDIR)
    fig_forest = forest_rate_ratio(df_types, OUTDIR)

    # ---- PPTX report
    prs = Presentation()
    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TE density vs FST-enriched windows"
    slide.placeholders[1].text = (
        f"Significance mode: {SIGNIFICANCE_MODE!r} (q_poi < {Q_CUTOFF})\n"
        f"Overall model: {model_used} (auto-switch to NB if overdispersed)\n"
        "Counts modeled with GLM (offset = log window length)\n"
        "Per-TE-type effects include BH-FDR across types"
    )
    # Summary slide
    bullets = [
        f"Windows: total={len(df_aug):,}; significant={int(df_aug['is_sig'].sum()):,}; "
        f"not={int((~df_aug['is_sig']).sum()):,}",
        f"Overall RR (sig/not) = {res_used['RR']:.3f} "
        f"[{res_used['RR_CI_lo']:.3f}, {res_used['RR_CI_hi']:.3f}], p={res_used['p']:.2e}",
        f"Poisson dispersion ≈ {res_used['dispersion']:.2f}  (>1.5 ⇒ overdispersion → NB)",
        f"Global counts rate-ratio: RR={rr_glob:.3f} "
        f"[{ci_lo_glob:.3f}, {ci_hi_glob:.3f}], p={p_glob:.2e}",
        ("Top TE types by FDR (q<0.05): " +
         ", ".join(df_types.loc[df_types['q']<0.05, 'TE_type'].head(10).astype(str)))
            if (df_types['q']<0.05).any() else "No TE types pass FDR<0.05",
        f"Tables: {out_aug.name}, {out_types.name}, {out_summ.name}"
    ]
    # add figures
    def add_slide_local(prs, title, bullets=None, image=None):
        layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if bullets:
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.6))
            tf = tb.text_frame; tf.clear()
            for i,b in enumerate(bullets):
                p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
                p.text = b; p.level = 0
        if image and Path(image).exists():
            slide.shapes.add_picture(str(image), Inches(0.5), Inches(2.7), height=Inches(3.8))

    add_slide_local(prs, "Overview & key results", bullets=bullets, image=str(fig_box))
    add_slide_local(prs, "Per-TE-type rate ratios (GLM)", bullets=None, image=str(fig_forest))
    pptx_path = OUTDIR / "TE_vs_FST_report.pptx"
    prs.save(str(pptx_path))

    print(f"\nDone. Outputs in: {OUTDIR.resolve()}")
    print(f"- Augmented windows: {out_aug.name}")
    print(f"- Group summary:     {out_summ.name}")
    print(f"- Per-type results:  {out_types.name}")
    print(f"- Figures:           {fig_box.name}, {fig_forest.name}")
    print(f"- Report:            {pptx_path.name}")

if __name__ == "__main__":
    main()
