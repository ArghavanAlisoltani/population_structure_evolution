# analyze_tmrca_genes_enrichment.py  (updated)
# ------------------------------------------------------------
# EXACT INPUT COLUMNS (tab-delimited):
# CHROM  start_tmrca  end_tmrca  mean_tmrca  lower_CI  upper_CI
# overlap_genes  overlap_gene_n  nearest_genes
#
# Parts:
# A) Clean + derive nearest-gene distance (absolute bp; 0 if overlap_gene_n>0)
# B) Top 5% vs Background (by mean_tmrca): hist/ECDF of distances + Fisher enrichment at cutoffs
# C) Sliding windows (win=WIN_BP, step=STEP_BP): Binomial/Poisson/Beta-Bin/NegBin tests; Manhattan plots
# D) NEW: Segment-length analysis: summaries, hist (log), ECDF, length vs mean_tmrca hexbin
# E) PPTX report with all panels
# ------------------------------------------------------------

import os
import re
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

from scipy.stats import fisher_exact, binom, poisson, betabinom, nbinom, ks_2samp, mannwhitneyu, spearmanr
from statsmodels.stats.multitest import multipletests
from pptx import Presentation
from pptx.util import Inches

# --------------------------
# USER SETTINGS
# --------------------------
INFILE  = "annotated_tmrca_4_GPT_13columns.tsv"
OUTDIR  = "tmrca_gene_enrichment_report"
os.makedirs(OUTDIR, exist_ok=True)

# Top-% cut for gene-distance part
TOP_PCT = 0.05   # Top 5%

# Distance thresholds (bp) for enrichment tables (Part B)
DIST_THRESH = [1_000, 5_000, 10_000, 50_000]

# Windowing (Part C)
WIN_BP  = 1_000_000    # window size (bp)
STEP_BP = 100_000      # step between windows (overlap = WIN_BP - STEP_BP; must be >0)
THR_HI  = 50_000       # “old” threshold on mean_tmrca (generations); adjust as needed

# Manhattan plotting
TOP_K_GLOBAL = 10
TOP1_LABELS  = True

# --------------------------
# Load and basic cleaning
# --------------------------
need = ["CHROM","start_tmrca","end_tmrca","mean_tmrca","lower_CI","upper_CI",
        "overlap_genes","overlap_gene_n","nearest_genes"]
df = pd.read_csv(INFILE, sep="\t", dtype={"CHROM":str})
missing = [c for c in need if c not in df.columns]
if missing:
    raise ValueError(f"Missing required columns: {missing}\nFound: {list(df.columns)}")

for c in ["start_tmrca","end_tmrca","mean_tmrca","lower_CI","upper_CI","overlap_gene_n"]:
    df[c] = pd.to_numeric(df[c], errors="coerce")

df = df.dropna(subset=["CHROM","start_tmrca","end_tmrca","mean_tmrca"]).copy()
df["start_tmrca"] = df["start_tmrca"].astype(np.int64)
df["end_tmrca"]   = df["end_tmrca"].astype(np.int64)
df["pos_mid"]     = (df["start_tmrca"] + df["end_tmrca"]) / 2.0

# ---------- NEW: segment length (bp)
df["seg_len_bp"] = (df["end_tmrca"] - df["start_tmrca"] + 1).clip(lower=1)

# --------------------------
# Parse nearest-gene distance (absolute)
# --------------------------
int_re = re.compile(r"(-?\d+)")
def parse_distance(s):
    if pd.isna(s): return np.nan
    m = int_re.search(str(s).replace(",", ""))
    return float(m.group(1)) if m else np.nan

df["nearest_gene_dist"] = df["nearest_genes"].apply(parse_distance)
df.loc[df["overlap_gene_n"].fillna(0) > 0, "nearest_gene_dist"] = 0.0
df["nearest_gene_dist"] = df["nearest_gene_dist"].abs()

# --------------------------
# PART D: Segment-length analysis
# --------------------------
# 1) Summary stats
lengths = df["seg_len_bp"].astype(float).values
quantiles = np.quantile(lengths, [0, .01, .05, .25, .5, .75, .95, .99, 1.0])
len_summary = pd.DataFrame({
    "stat": ["min","p01","p05","q25","median","q75","p95","p99","max"],
    "bp":   quantiles.astype(np.float64)
})
len_summary["bp"] = len_summary["bp"].map(lambda x: float(x))
len_summary.to_csv(os.path.join(OUTDIR,"segment_length_summary.csv"), index=False)

# 2) Histogram (log x)
plt.figure(figsize=(8,4))
plt.hist(lengths, bins=np.logspace(np.log10(max(1,lengths.min())), np.log10(lengths.max()), 50),
         color="#6C8EBF", alpha=0.8)
plt.xscale("log")
plt.xlabel("Segment length (bp, log scale)")
plt.ylabel("Count")
plt.title("TMRCA segment lengths (log-scale histogram)")
len_hist_png = os.path.join(OUTDIR, "segment_length_hist_log.png")
plt.tight_layout(); plt.savefig(len_hist_png, dpi=220); plt.close()

# 3) ECDF of lengths
def ecdf(arr):
    a = np.asarray(arr, float)
    a = a[np.isfinite(a)]
    if a.size == 0: return np.array([]), np.array([])
    x = np.sort(a)
    y = np.arange(1, a.size+1)/a.size
    return x, y

xL, yL = ecdf(lengths)
plt.figure(figsize=(8,4))
if xL.size: plt.step(xL, yL, where='post', color="#6C8EBF")
plt.xscale("log")
plt.xlabel("Segment length (bp, log scale)")
plt.ylabel("Empirical CDF")
plt.title("ECDF of TMRCA segment lengths")
len_ecdf_png = os.path.join(OUTDIR, "segment_length_ecdf.png")
plt.tight_layout(); plt.savefig(len_ecdf_png, dpi=220); plt.close()

# 4) Length vs mean TMRCA (hexbin)
plt.figure(figsize=(8,5))
hb = plt.hexbin(df["seg_len_bp"], df["mean_tmrca"], gridsize=60, xscale='log', bins='log', cmap="viridis")
plt.colorbar(hb, label="log10(count)")
plt.xlabel("Segment length (bp, log scale)")
plt.ylabel("Mean TMRCA (generations)")
plt.title("Length vs mean TMRCA (hexbin, log-scaled length)")
len_tmrca_hex_png = os.path.join(OUTDIR, "length_vs_tmrca_hex.png")
plt.tight_layout(); plt.savefig(len_tmrca_hex_png, dpi=220); plt.close()

# Correlation (Spearman, robust to monotonic relationships)
try:
    rho, pval = spearmanr(df["seg_len_bp"], df["mean_tmrca"], nan_policy="omit")
except Exception:
    rho, pval = (np.nan, np.nan)

with open(os.path.join(OUTDIR, "segment_length_notes.txt"), "w") as fh:
    fh.write("Segment length summary (see CSV) and correlation with mean TMRCA\n")
    fh.write(f"Spearman rho={rho:.4f}, p={pval:.3g}\n")

# --------------------------
# PART B: Top 5% vs Background (distances to genes)
# --------------------------
cut = df["mean_tmrca"].quantile(1.0 - TOP_PCT)
df["group"] = np.where(df["mean_tmrca"] >= cut, "TOP5%", "BACKGROUND")

sub = df[df["nearest_gene_dist"].notna()].copy()
top = sub[sub["group"]=="TOP5%"]["nearest_gene_dist"].values
bkg = sub[sub["group"]=="BACKGROUND"]["nearest_gene_dist"].values

ks = ks_2samp(top, bkg, alternative="two-sided", mode="auto") if (len(top)>0 and len(bkg)>0) else None
mw = mannwhitneyu(top, bkg, alternative="two-sided")          if (len(top)>0 and len(bkg)>0) else None

# Histogram
bins = np.unique(np.concatenate([
    np.linspace(0, 5_000, 20),
    np.linspace(5_000, 50_000, 20),
    np.linspace(50_000, 500_000, 20)
]))
plt.figure(figsize=(9,4))
if len(bkg): plt.hist(bkg, bins=bins, color="#B0B0B0", alpha=0.7, label="Background", density=True)
if len(top): plt.hist(top, bins=bins, color="#4C78A8", alpha=0.7, label="Top 5%", density=True)
plt.xlabel("Distance to nearest gene (bp)")
plt.ylabel("Density")
title_stats = f"KS p={ks.pvalue:.2e}, MW p={mw.pvalue:.2e}" if (ks and mw) else ""
plt.title(f"Nearest gene distance: histogram  {title_stats}")
plt.legend()
gene_hist_png = os.path.join(OUTDIR, "gene_dist_hist.png")
plt.tight_layout(); plt.savefig(gene_hist_png, dpi=220); plt.close()

# ECDF
def ecdf_xy(a):
    a = np.asarray(a, float)
    a = a[~np.isnan(a)]
    if a.size == 0: return np.array([]), np.array([])
    x = np.sort(a); y = np.arange(1, a.size+1)/a.size; return x, y

xT,yT = ecdf_xy(top); xB,yB = ecdf_xy(bkg)
plt.figure(figsize=(9,4))
if xB.size: plt.step(xB, yB, where='post', color="#B0B0B0", label="Background")
if xT.size: plt.step(xT, yT, where='post', color="#4C78A8", label="Top 5%")
for D in DIST_THRESH: plt.axvline(D, color="k", alpha=0.12, lw=1)
plt.xlabel("Distance to nearest gene (bp)")
plt.ylabel("Empirical CDF")
plt.title(f"Nearest gene distance: CDF  {title_stats}")
plt.legend()
gene_cdf_png = os.path.join(OUTDIR, "gene_dist_cdf.png")
plt.tight_layout(); plt.savefig(gene_cdf_png, dpi=220); plt.close()

# Enrichment at cutoffs
enr_rows = []
for D in DIST_THRESH:
    a_in  = int((sub["group"].eq("TOP5%") & (sub["nearest_gene_dist"] <= D)).sum())
    a_out = int((sub["group"].eq("TOP5%") & (sub["nearest_gene_dist"] >  D)).sum())
    b_in  = int((sub["group"].eq("BACKGROUND") & (sub["nearest_gene_dist"] <= D)).sum())
    b_out = int((sub["group"].eq("BACKGROUND") & (sub["nearest_gene_dist"] >  D)).sum())
    _, p = fisher_exact([[a_in,a_out],[b_in,b_out]], alternative="two-sided")
    top_rate = a_in / max(1, (a_in+a_out)); bkg_rate = b_in / max(1, (b_in+b_out))
    rr = (top_rate / bkg_rate) if bkg_rate>0 else np.inf
    enr_rows.append(dict(distance_bp=D, TOP_in=a_in, TOP_out=a_out, BKG_in=b_in, BKG_out=b_out,
                         TOP_rate=top_rate, BKG_rate=bkg_rate, rate_ratio=rr, fisher_p=p))
gene_enrich = pd.DataFrame(enr_rows)
gene_enrich_csv = os.path.join(OUTDIR, "gene_enrichment_table.csv")
gene_enrich.to_csv(gene_enrich_csv, index=False)

# --------------------------
# PART C: Window enrichment
# --------------------------
def scaffold_sort_key(chrom: str):
    m = re.search(r"(\d+)([A-Za-z]?)$", str(chrom))
    if not m: return (10**9, 9, str(chrom))
    num = int(m.group(1)); let = m.group(2).lower()
    let_rank = (ord(let)-ord('a')) if let else 9
    return (num, let_rank, str(chrom))

def build_windows_for_scaffold(xmin: int, xmax: int, win_bp: int, step_bp: int):
    if step_bp <= 0: raise ValueError("STEP_BP must be > 0.")
    starts = np.arange(xmin, xmax+1, step_bp, dtype=np.int64)
    ends   = np.minimum(starts + win_bp - 1, xmax)
    return np.vstack([starts, ends]).T

scaffolds = sorted(df["CHROM"].unique().tolist(), key=scaffold_sort_key)
by_scaf = {sc: df[df["CHROM"]==sc].sort_values(["start_tmrca","end_tmrca"]).reset_index(drop=True) for sc in scaffolds}

def coverage_counts_for_scaffold(df_scaf: pd.DataFrame, win_arr: np.ndarray, thr_hi: float):
    s = df_scaf["start_tmrca"].to_numpy(np.int64)
    e = df_scaf["end_tmrca"].to_numpy(np.int64)
    y = df_scaf["mean_tmrca"].to_numpy(float)
    M = win_arr.shape[0]
    G = np.zeros(M, dtype=np.int64)
    K = np.zeros(M, dtype=np.int64)
    j0 = 0
    for i,(ws,we) in enumerate(win_arr):
        while j0 < len(s) and e[j0] <= ws: j0 += 1
        j = j0
        while j < len(s) and s[j] < we+1:
            ov = max(0, min(e[j], we) - max(s[j], ws) + 1)
            if ov > 0:
                G[i] += ov
                if y[j] >= thr_hi: K[i] += ov
            j += 1
    return G, K

win_tables = []
for sc in scaffolds:
    d = by_scaf[sc]
    if d.empty: continue
    xmin, xmax = int(d["start_tmrca"].min()), int(d["end_tmrca"].max())
    win_arr = build_windows_for_scaffold(xmin, xmax, WIN_BP, STEP_BP)
    if win_arr.size == 0: continue
    G, K = coverage_counts_for_scaffold(d, win_arr, THR_HI)
    tab = pd.DataFrame({
        "CHROM": sc,
        "WIN_START": win_arr[:,0],
        "WIN_END":   win_arr[:,1],
        "BP_COVERED": G,
        f"BP_OLD_ge_{THR_HI}": K
    })
    tab["WIN_MID"] = (tab["WIN_START"] + tab["WIN_END"]) // 2
    win_tables.append(tab)

if not win_tables:
    raise RuntimeError("No windows computed — check WIN_BP/STEP_BP and input ranges.")
win_all = pd.concat(win_tables, ignore_index=True)

totG = int(win_all["BP_COVERED"].sum())
totK = int(win_all[f"BP_OLD_ge_{THR_HI}"].sum())
p_glob = (totK / totG) if totG>0 else np.nan

G = win_all["BP_COVERED"].to_numpy(np.int64)
K = win_all[f"BP_OLD_ge_{THR_HI}"].to_numpy(np.int64)

# Binomial (upper tail)
p_bin = np.full_like(G, np.nan, dtype=float)
if np.isfinite(p_glob):
    ok = G > 0
    p_bin[ok]  = 1.0 - binom.cdf(np.maximum(0, K[ok]) - 1, G[ok], p_glob)
    p_bin[~ok] = np.where(K[~ok] > 0, 0.0, 1.0)
q_bin = multipletests(p_bin[np.isfinite(p_bin)], method="fdr_bh")[1] if np.isfinite(p_bin).any() else np.array([])

# Poisson with exposure
mu = p_glob * G if np.isfinite(p_glob) else np.full_like(G, np.nan, dtype=float)
p_poi = np.full_like(G, np.nan, dtype=float)
if np.isfinite(p_glob):
    ok = mu > 0
    p_poi[ok]  = 1.0 - poisson.cdf(np.maximum(0, K[ok]) - 1, mu[ok])
    p_poi[~ok] = np.where(K[~ok] > 0, 0.0, 1.0)
q_poi = multipletests(p_poi[np.isfinite(p_poi)], method="fdr_bh")[1] if np.isfinite(p_poi).any() else np.array([])

def putback_q(p_arr, q_sub):
    q_full = np.full_like(p_arr, np.nan, dtype=float)
    mask = np.isfinite(p_arr); q_full[mask] = q_sub; return q_full

win_all["p_bin"] = p_bin; win_all["q_bin"] = putback_q(p_bin, q_bin)
win_all["p_poi"] = p_poi; win_all["q_poi"] = putback_q(p_poi, q_poi)

# Beta–Binomial
def beta_bin_pvals(G, K, p0):
    ok = (G > 0)
    if not ok.any() or not np.isfinite(p0) or p0<=0 or p0>=1:
        return np.full_like(G, np.nan, dtype=float)
    V = G*p0*(1-p0); chi = np.zeros_like(G, dtype=float)
    chi[ok] = ((K[ok] - G[ok]*p0)**2) / np.maximum(V[ok], 1e-12)
    dfree = max(1, int(ok.sum()) - 1)
    disp_bin = float(np.sum(chi[ok]) / dfree)
    n_bar = G[ok].mean()
    rho = max(0.0, min(0.9999, (disp_bin - 1.0)/max(1e-9, (n_bar - 1.0))))
    if rho <= 1e-10:
        p = np.full_like(G, np.nan, dtype=float)
        p[ok]  = 1.0 - binom.cdf(np.maximum(0, K[ok]) - 1, G[ok], p0)
        p[~ok] = np.where(K[~ok] > 0, 0.0, 1.0)
        return p
    t = (1.0/rho) - 1.0
    a_hat, b_hat = p0*t, (1.0-p0)*t
    p = np.full_like(G, np.nan, dtype=float)
    p[ok]  = betabinom.sf(np.maximum(0, K[ok]) - 1, G[ok], a_hat, b_hat)
    p[~ok] = np.where(K[~ok] > 0, 0.0, 1.0)
    return p

# Negative Binomial
def negbin_pvals(G, K, mu):
    p = np.full_like(G, np.nan, dtype=float)
    ok = np.isfinite(mu)
    if not ok.any(): return p
    resid = (K - mu); numer = np.sum(mu[ok]**2); denom = np.sum((resid[ok]**2) - mu[ok])
    theta = numer/denom if denom > 1e-12 and numer>0 else 1e9
    r = max(theta, 1e-6); p_param = np.where(ok, r/(r+mu), np.nan)
    p[ok] = nbinom.sf(np.maximum(0, K[ok]) - 1, r, p_param[ok])
    return p

win_all["p_bb"] = beta_bin_pvals(G,K,p_glob)
win_all["q_bb"] = putback_q(win_all["p_bb"],
                            multipletests(win_all["p_bb"][np.isfinite(win_all["p_bb"])], method="fdr_bh")[1]) \
                  if np.isfinite(win_all["p_bb"]).any() else np.nan
win_all["p_nb"] = negbin_pvals(G,K,mu)
win_all["q_nb"] = putback_q(win_all["p_nb"],
                            multipletests(win_all["p_nb"][np.isfinite(win_all["p_nb"])], method="fdr_bh")[1]) \
                  if np.isfinite(win_all["p_nb"]).any() else np.nan

# --------------------------
# Manhattan plots
# --------------------------
def build_concat_positions(win_df: pd.DataFrame):
    chroms = sorted(win_df["CHROM"].unique(), key=scaffold_sort_key)
    offsets, xticks, xtlabs = {}, [], []
    cursor = 0
    for ch in chroms:
        sub = win_df[win_df["CHROM"]==ch]
        if sub.empty: continue
        max_mid = int(sub["WIN_MID"].max())
        offsets[ch] = cursor
        xticks.append(cursor + max_mid//2)
        xtlabs.append(ch)
        cursor += max_mid + 1
    return offsets, list(zip(xtlabs, xticks))

offsets, chrom_ticks = build_concat_positions(win_all)

def manhattan_q(win_df, q_col, out_png, title):
    dd = win_df.copy()
    dd["offset"] = dd["CHROM"].map(offsets).fillna(0).astype(np.int64)
    dd["cumpos"] = dd["offset"] + dd["WIN_MID"]
    q = dd[q_col].astype(float)
    dd["ml10"] = -np.log10(np.clip(q, 1e-300, 1.0))

    # top-K global
    topK = dd.nlargest(TOP_K_GLOBAL, "ml10")

    base_colors = {ch: ("#4C78A8" if (i % 2 == 0) else "#A0A0A0")
                   for i, ch in enumerate(sorted(dd["CHROM"].unique(), key=scaffold_sort_key))}
    colors = dd["CHROM"].map(base_colors).values

    fig, ax = plt.subplots(figsize=(18,4.5))
    ax.scatter(dd["cumpos"], dd["ml10"], s=10, c=colors, alpha=0.9, linewidths=0)
    ax.scatter(topK["cumpos"], topK["ml10"], s=28, c="red", alpha=1.0, linewidths=0, zorder=3)
    ax.set_ylabel(r"$-\log_{10}(q)$")
    ax.set_xlabel("Genome (scaffolds concatenated)")
    ax.set_xticks([t for _,t in chrom_ticks])
    ax.set_xticklabels([c for c,_ in chrom_ticks], rotation=90, fontsize=9)
    ax.set_title(title)
    fig.tight_layout(); fig.savefig(out_png, dpi=220); plt.close(fig)

poi_manhattan_png = os.path.join(OUTDIR, f"manhattan_q_poi_win{WIN_BP}_step{STEP_BP}.png")
manhattan_q(win_all, "q_poi", poi_manhattan_png,
            title=f"Poisson FDR (q) — win={WIN_BP:,}, step={STEP_BP:,} — Top {TOP_K_GLOBAL} in red")

def manhattan_top1_per_scaffold(win_df, q_col, out_png, annotate=True):
    dd = win_df.copy()
    dd["offset"] = dd["CHROM"].map(offsets).fillna(0).astype(np.int64)
    dd["cumpos"] = dd["offset"] + dd["WIN_MID"]
    q = dd[q_col].astype(float)
    dd["ml10"] = -np.log10(np.clip(q, 1e-300, 1.0))

    top1 = []
    for ch in sorted(dd["CHROM"].unique(), key=scaffold_sort_key):
        sub = dd[dd["CHROM"]==ch]
        if sub.empty: continue
        top1.append(sub.loc[sub["ml10"].idxmax()])
    top1 = pd.DataFrame(top1)

    base_colors = {ch: ("#4C78A8" if (i % 2 == 0) else "#A0A0A0")
                   for i, ch in enumerate(sorted(dd["CHROM"].unique(), key=scaffold_sort_key))}
    colors = dd["CHROM"].map(base_colors).values

    fig, ax = plt.subplots(figsize=(18,4.8))
    ax.scatter(dd["cumpos"], dd["ml10"], s=8, c=colors, alpha=0.35, linewidths=0)
    ax.scatter(top1["cumpos"], top1["ml10"], s=36, c="red", alpha=1.0, linewidths=0, zorder=3)
    if annotate and not top1.empty:
        for _, r in top1.iterrows():
            label = f"{r['CHROM']}:{int(r['WIN_START'])}-{int(r['WIN_END'])}"
            ax.text(r["cumpos"], r["ml10"]+0.15, label, fontsize=7, color="red", rotation=35,
                    ha="left", va="bottom")

    ax.set_ylabel(r"$-\log_{10}(q)$")
    ax.set_xlabel("Genome (scaffolds concatenated)")
    ax.set_xticks([t for _,t in chrom_ticks])
    ax.set_xticklabels([c for c,_ in chrom_ticks], rotation=90, fontsize=9)
    ax.set_title(f"Top-1 window per scaffold (by −log10 q, {q_col}) — win={WIN_BP:,}, step={STEP_BP:,}")
    fig.tight_layout(); fig.savefig(out_png, dpi=220); plt.close(fig)

poi_top1_png = os.path.join(OUTDIR, f"manhattan_q_poi_TOP1_per_scaffold_win{WIN_BP}_step{STEP_BP}.png")
manhattan_top1_per_scaffold(win_all, "q_poi", poi_top1_png, annotate=TOP1_LABELS)

# --------------------------
# Annotate significant windows (Poisson q<0.05) with overlapped genes
# --------------------------
def split_genes(s):
    if pd.isna(s) or str(s).strip()=="":
        return []
    tmp = re.split(r"[;,|]\s*|\s{2,}", str(s).strip())
    if len(tmp)==1:
        tmp = str(s).replace("|",";").replace(",", ";").split(";")
    return [t for t in (x.strip() for x in tmp) if t]

sig = win_all[np.isfinite(win_all["q_poi"]) & (win_all["q_poi"] < 0.05)].copy().sort_values("q_poi")
by_scaf_rows = {sc: by_scaf[sc] for sc in scaffolds}
rows = []
for _, row in sig.iterrows():
    ch, ws, we = row["CHROM"], int(row["WIN_START"]), int(row["WIN_END"])
    segs = by_scaf_rows[ch]
    ov = segs[(segs["end_tmrca"] >= ws) & (segs["start_tmrca"] <= we)]
    genes = []
    for g in ov["overlap_genes"]:
        genes.extend(split_genes(g))
    genes = sorted(set([g for g in genes if g]))
    rows.append({
        "CHROM": ch, "WIN_START": ws, "WIN_END": we,
        "BP_COVERED": int(row["BP_COVERED"]),
        "BP_OLD": int(row[f"BP_OLD_ge_{THR_HI}"]),
        "p_poi": row["p_poi"], "q_poi": row["q_poi"],
        "overlap_genes_in_window": ";".join(genes)
    })
top_sig_tbl = pd.DataFrame(rows)
top_sig_csv = os.path.join(OUTDIR, "significant_windows_with_genes.csv")
top_sig_tbl.to_csv(top_sig_csv, index=False)

# Save windows table
win_csv = os.path.join(OUTDIR, f"tmrca_window_tests_win{WIN_BP}_step{STEP_BP}_thr{THR_HI}.csv")
win_all.to_csv(win_csv, index=False)

# --------------------------
# PowerPoint report
# --------------------------
prs = Presentation()

# Title
s0 = prs.slides.add_slide(prs.slide_layouts[0])
s0.shapes.title.text = "TMRCA: Gene Proximity, Segment Lengths & Window Enrichment"
s0.placeholders[1].text = (
    f"Input: {os.path.basename(INFILE)}\n"
    f"Top group = top {int(TOP_PCT*100)}% by mean_tmrca (cut={df['mean_tmrca'].quantile(1-TOP_PCT):.3g})\n"
    f"Windows: win={WIN_BP:,} bp, step={STEP_BP:,}; 'old' threshold: {THR_HI} generations\n"
    f"Distances are absolute; length = end - start + 1"
)

# NEW: segment length slides
s_len1 = prs.slides.add_slide(prs.slide_layouts[5])
s_len1.shapes.title.text = "Segment length distribution (log histogram)"
s_len1.shapes.add_picture(len_hist_png, Inches(0.5), Inches(1.2), width=Inches(9))

s_len2 = prs.slides.add_slide(prs.slide_layouts[5])
s_len2.shapes.title.text = "Segment length ECDF (log x)"
s_len2.shapes.add_picture(len_ecdf_png, Inches(0.5), Inches(1.2), width=Inches(9))

s_len3 = prs.slides.add_slide(prs.slide_layouts[5])
s_len3.shapes.title.text = "Length vs mean TMRCA (hexbin)"
s_len3.shapes.add_picture(len_tmrca_hex_png, Inches(0.5), Inches(1.2), width=Inches(9))

# Distance slides
s1 = prs.slides.add_slide(prs.slide_layouts[5])
s1.shapes.title.text = "Nearest gene: histogram (TOP vs BACKGROUND)"
s1.shapes.add_picture(gene_hist_png, Inches(0.5), Inches(1.2), width=Inches(9))

s2 = prs.slides.add_slide(prs.slide_layouts[5])
s2.shapes.title.text = "Nearest gene: ECDF (TOP vs BACKGROUND)"
s2.shapes.add_picture(gene_cdf_png, Inches(0.5), Inches(1.2), width=Inches(9))

# Manhattan slides
s4 = prs.slides.add_slide(prs.slide_layouts[5])
s4.shapes.title.text = "Manhattan (−log10 q, Poisson)"
s4.shapes.add_picture(poi_manhattan_png, Inches(0.5), Inches(1.2), width=Inches(9))

s5 = prs.slides.add_slide(prs.slide_layouts[5])
s5.shapes.title.text = "Manhattan: Top-1 per scaffold (annotated)"
s5.shapes.add_picture(poi_top1_png, Inches(0.5), Inches(1.2), width=Inches(9))

pptx_path = os.path.join(OUTDIR, "tmrca_gene_enrichment_report.pptx")
prs.save(pptx_path)

# Console summary
print("\n=== DONE ===")
print("Output dir:", os.path.abspath(OUTDIR))
print("Segment length summary:", "segment_length_summary.csv")
print("Distance plots:", os.path.basename(gene_hist_png), os.path.basename(gene_cdf_png))
print("Length plots:", os.path.basename(len_hist_png), os.path.basename(len_ecdf_png), os.path.basename(len_tmrca_hex_png))
print("Windows table:", os.path.basename(win_csv))
print("Manhattan plots:", os.path.basename(poi_manhattan_png), os.path.basename(poi_top1_png))
print("Significant windows with genes:", os.path.basename(top_sig_csv))
print("Slides:", os.path.basename(pptx_path))

