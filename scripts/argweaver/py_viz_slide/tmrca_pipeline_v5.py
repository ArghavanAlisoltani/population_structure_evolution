# tmrca_pipeline_v3.py
# ------------------------------------------------------------
# End-to-end TMRCA pipeline for ARGweaver 6-col tracks:
#   • Reads tab-delimited: chrom start end tmrca_mean tmrca_lo tmrca_hi
#   • Summaries (all vs zeros-removed) + per-scaffold
#   • Coverage-weighted windowing with overlap (mean + CI lo/hi)
#   • TWO Manhattan plots (linear & log):
#       - points < THRESH_FADE use low alpha; zeros excluded
#       - Top-K windows (global) highlighted in red
#   • EXTRA Manhattan: highlight **top-1 per scaffold** in red + annotate window name
#   • NEW: Global **line plot** across scaffolds with CI ribbon (using windowed lo/hi)
#   • PowerPoint deck with all figures and tables
# ------------------------------------------------------------

import os, re
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches

# ============ USER SETTINGS ============
INPUT_FILE = "All_tmrca_60_run.txt"   # 6 columns, whitespace/TSV
OUTDIR     = "tmrca_reports"

# Coordinates for display (doesn't affect windowing)
ZERO_BASED_START = True

# Windowing
WIN_BP  = 1_000_000      # window size (bp)
STEP_BP = 100_000        # step (bp). Overlap = WIN_BP - STEP_BP

# Plot controls
TOP_K             = 10        # highlight the top-K windows by mean TMRCA (global)
THRESH_FADE       = 500.0     # points with y_mean < this get a low alpha
ALPHA_FAINT       = 0.12      # alpha for faint points (< threshold)
ALPHA_STRONG      = 0.90      # alpha for points >= threshold
POINT_SIZE        = 12        # non-top points (linear & log)
POINT_SIZE_TOP    = 46        # red top points
ALT_SCAFF_COLORS  = ("#4C78A8", "#A0A0A0")  # alternating scaffold colors (blue/gray)

np.random.seed(13)
os.makedirs(OUTDIR, exist_ok=True)

# ============ HELPERS ============
def parse_scaffold_order(chrom: str):
    """
    Sort key so labels go: 1a, 1b, 2, 3, ..., 27, 31, 44, etc.
    Works with names like 'scaffold_1a', 'scaffold_2', '1b', '20', etc.
    """
    m = re.search(r'(\d+)([A-Za-z]?)$', chrom)
    if not m:  # unknown -> send to end
        return (10**9, 10**9, chrom)
    num = int(m.group(1))
    letter = m.group(2).lower()
    # order a(0), b(1), c(2), ..., ''(9 so number-without-letter after letters)
    letter_rank = (ord(letter) - ord('a')) if letter else 9
    return (num, letter_rank, chrom)

def window_reduce(df_scaf: pd.DataFrame, win_bp: int, step_bp: int) -> pd.DataFrame:
    """
    Coverage-weighted windowing for a single scaffold.
    Returns mean TMRCA plus coverage-weighted mean of lo/hi for CI ribbon.
    """
    start_min = int(max(0, df_scaf["start"].min()))
    end_max   = int(df_scaf["end"].max())
    starts = np.array([start_min], dtype=int) if end_max - start_min < win_bp \
             else np.arange(start_min, end_max - win_bp + 1, step_bp, dtype=int)

    d = df_scaf.sort_values(["start","end"]).reset_index(drop=True)
    s_arr = d["start"].to_numpy()
    e_arr = d["end"].to_numpy()
    y_arr = d["tmrca_mean"].to_numpy()
    lo_arr = d["tmrca_lo"].to_numpy()
    hi_arr = d["tmrca_hi"].to_numpy()

    rows, j0 = [], 0
    for wstart in starts:
        wend = wstart + win_bp
        while j0 < len(d) and e_arr[j0] <= wstart:
            j0 += 1
        j, cov, num_y, num_lo, num_hi = j0, 0, 0.0, 0.0, 0.0
        while j < len(d) and s_arr[j] < wend:
            ovl = max(0, min(e_arr[j], wend) - max(s_arr[j], wstart))
            if ovl > 0:
                cov   += ovl
                num_y += ovl * y_arr[j]
                num_lo += ovl * lo_arr[j]
                num_hi += ovl * hi_arr[j]
            j += 1
        if cov > 0:
            mean_y  = num_y / cov
            mean_lo = num_lo / cov
            mean_hi = num_hi / cov
        else:
            mean_y = mean_lo = mean_hi = np.nan
        rows.append((wstart, wend, mean_y, mean_lo, mean_hi, cov))
    return pd.DataFrame(rows, columns=["w_start","w_end","y_mean","y_lo","y_hi","bp_covered"])

def plot_manhattan(win_all, scaffolds_ordered, chrom_ticks, yscale, out_png, title_suffix):
    """Scatter without bands; fade points below threshold; highlight global Top-K in red."""
    y = win_all["y_mean"].to_numpy()
    alphas = np.where(y >= THRESH_FADE, ALPHA_STRONG, ALPHA_FAINT)

    # Colors alternate by scaffold
    chrom2col = {sc: ALT_SCAFF_COLORS[i % 2] for i, sc in enumerate(scaffolds_ordered)}
    colors = np.array([chrom2col[c] for c in win_all["chrom"]], dtype=object)

    # Top-K by y_mean (positive only)
    topK = win_all[win_all["y_mean"] > 0].nlargest(TOP_K, "y_mean")
    mask_top = win_all.index.isin(topK.index)

    fig, ax = plt.subplots(figsize=(18, 4.5))
    # non-top
    ax.scatter(win_all.loc[~mask_top, "mid_global"],
               win_all.loc[~mask_top, "y_mean"],
               s=POINT_SIZE, c=colors[~mask_top], alpha=alphas[~mask_top], linewidths=0)
    # top-K
    ax.scatter(win_all.loc[mask_top, "mid_global"],
               win_all.loc[mask_top, "y_mean"],
               s=POINT_SIZE_TOP, c="red", alpha=1.0, linewidths=0, zorder=3)

    if yscale == "log":
        ax.set_yscale("log")
        ylab = "Windowed TMRCA mean (generations, log10)"
    else:
        ylab = "Windowed TMRCA mean (generations)"

    ax.set_ylabel(ylab)
    ax.set_xlabel("Genomic position (scaffolds concatenated)")
    ax.set_xticks([t for _, t in chrom_ticks])
    ax.set_xticklabels([c for c, _ in chrom_ticks], rotation=90, fontsize=9)
    ax.set_title(f"Windowed TMRCA Manhattan plot (win={WIN_BP:,} bp, step={STEP_BP:,} bp)\n"
                 f"Top {TOP_K} windows in red {title_suffix}")
    fig.tight_layout()
    fig.savefig(out_png, dpi=220)
    plt.close(fig)

def plot_manhattan_top1_per_scaffold(win_all, scaffolds_ordered, chrom_ticks, yscale, out_png):
    """
    Manhattan variant: color only the **top-1** window in each scaffold in red
    and annotate the point with window name 'chrom:start-end'.
    All others obey the fade rule (alpha).
    """
    y = win_all["y_mean"].to_numpy()
    alphas = np.where(y >= THRESH_FADE, ALPHA_STRONG, ALPHA_FAINT)

    chrom2col = {sc: ALT_SCAFF_COLORS[i % 2] for i, sc in enumerate(scaffolds_ordered)}
    colors = np.array([chrom2col[c] for c in win_all["chrom"]], dtype=object)

    # pick top-1 per scaffold
    toprows = []
    for sc in scaffolds_ordered:
        sub = win_all[(win_all["chrom"] == sc) & (win_all["y_mean"] > 0)]
        if not sub.empty:
            toprows.append(sub.loc[sub["y_mean"].idxmax()])
    top_df = pd.DataFrame(toprows)
    mask_top = win_all.index.isin(top_df.index)

    fig, ax = plt.subplots(figsize=(18, 4.8))
    ax.scatter(win_all.loc[~mask_top, "mid_global"],
               win_all.loc[~mask_top, "y_mean"],
               s=POINT_SIZE, c=colors[~mask_top], alpha=alphas[~mask_top], linewidths=0)

    ax.scatter(win_all.loc[mask_top, "mid_global"],
               win_all.loc[mask_top, "y_mean"],
               s=POINT_SIZE_TOP, c="red", alpha=1.0, linewidths=0, zorder=3)

    # annotate top-1 per scaffold
    for _, r in top_df.iterrows():
        label = f"{r['chrom']}:{int(r['w_start'])}-{int(r['w_end'])}"
        ax.annotate(label,
                    (r["mid_global"], r["y_mean"]),
                    xytext=(0, 6), textcoords="offset points",
                    ha="center", va="bottom", fontsize=7, color="red")

    if yscale == "log":
        ax.set_yscale("log")
        ylab = "Windowed TMRCA mean (generations, log10)"
        title_suf = "(log scale)"
    else:
        ylab = "Windowed TMRCA mean (generations)"
        title_suf = "(linear scale)"

    ax.set_ylabel(ylab)
    ax.set_xlabel("Genomic position (scaffolds concatenated)")
    ax.set_xticks([t for _, t in chrom_ticks])
    ax.set_xticklabels([c for c, _ in chrom_ticks], rotation=90, fontsize=9)
    ax.set_title(f"Top-1 TMRCA window per scaffold highlighted & annotated {title_suf}")
    fig.tight_layout()
    fig.savefig(out_png, dpi=220)
    plt.close(fig)

def plot_line_with_ci(win_all, scaffolds_ordered, chrom_ticks, yscale, out_png):
    """
    Global line plot across concatenated scaffolds:
    - line = windowed mean TMRCA
    - ribbon = windowed CI (lo..hi) (coverage-weighted means)
    - zeros already removed upstream
    """
    fig, ax = plt.subplots(figsize=(18, 4.8))

    for i, sc in enumerate(scaffolds_ordered):
        sub = win_all[win_all["chrom"] == sc].sort_values("mid_global")
        if sub.empty: 
            continue
        y  = sub["y_mean"].to_numpy()
        lo = np.maximum(sub["y_lo"].to_numpy(), 1e-12)  # avoid <=0 for log
        hi = sub["y_hi"].to_numpy()
        x  = sub["mid_global"].to_numpy()

        # CI ribbon then line
        ax.fill_between(x, lo, hi, color="#B0B0B0", alpha=0.25, linewidth=0)
        ax.plot(x, y, color="#000000", linewidth=0.7)

    if yscale == "log":
        ax.set_yscale("log")
        ylab = "Windowed TMRCA mean (generations, log10)"
        title_suf = "(log scale)"
    else:
        ylab = "Windowed TMRCA mean (generations)"
        title_suf = "(linear scale)"

    ax.set_ylabel(ylab)
    ax.set_xlabel("Genomic position (scaffolds concatenated)")
    ax.set_xticks([t for _, t in chrom_ticks])
    ax.set_xticklabels([c for c, _ in chrom_ticks], rotation=90, fontsize=9)
    ax.set_title(f"Windowed TMRCA line plot with CI ribbon {title_suf}")
    fig.tight_layout()
    fig.savefig(out_png, dpi=220)
    plt.close(fig)

# ============ 1) READ DATA ============
cols = ["chrom","start","end","tmrca_mean","tmrca_lo","tmrca_hi"]
tm = pd.read_csv(INPUT_FILE, sep=r"\s+|\t", engine="python", header=None, names=cols)
for c in ["start","end","tmrca_mean","tmrca_lo","tmrca_hi"]:
    tm[c] = pd.to_numeric(tm[c], errors="coerce")
tm = tm.dropna(subset=["chrom","start","end","tmrca_mean"]).reset_index(drop=True)
tm["start_disp"] = tm["start"] + 1 if ZERO_BASED_START else tm["start"]
tm["end_disp"]   = tm["end"]

# ============ 2) SUMMARIES ============
def summarize_series(x: pd.Series, label: str) -> pd.DataFrame:
    x = pd.to_numeric(x, errors="coerce")
    return pd.DataFrame([{
        "n": int(x.size),
        "n_nonzero": int((x>0).sum()),
        "min": float(x.min()),
        "q1": float(x.quantile(0.25)),
        "median": float(x.median()),
        "mean": float(x.mean()),
        "q3": float(x.quantile(0.75)),
        "max": float(x.max()),
        "sd": float(x.std())
    }], index=[label])

summ_all = summarize_series(tm["tmrca_mean"], "all_values")
summ_no0 = summarize_series(tm.loc[tm["tmrca_mean"]>0, "tmrca_mean"], "nonzero_only")
summ_table = pd.concat([summ_all, summ_no0])
summ_table.to_csv(os.path.join(OUTDIR, "summary_global_tmrca.csv"))

per_scaf = tm.groupby("chrom")["tmrca_mean"].agg(
    n="size",
    n_nonzero=lambda s: (s>0).sum(),
    min="min",
    q1=lambda s: s.quantile(0.25),
    median="median",
    mean="mean",
    q3=lambda s: s.quantile(0.75),
    max="max",
    sd="std"
).reset_index()
per_scaf.to_csv(os.path.join(OUTDIR, "summary_by_scaffold.csv"), index=False)

# ============ 3) WINDOWING ============
scaffolds_ordered = sorted(tm["chrom"].unique().tolist(), key=parse_scaffold_order)

win_tables = []
offset = 0
chrom_ticks = []
xpos_min = offset

for chrom in scaffolds_ordered:
    d = tm.loc[tm["chrom"] == chrom, ["start","end","tmrca_mean","tmrca_lo","tmrca_hi"]].copy()
    if d.empty:
        continue
    w = window_reduce(d, WIN_BP, STEP_BP)
    w["chrom"] = chrom
    w["mid"] = (w["w_start"] + w["w_end"]) / 2.0
    w["mid_global"] = w["mid"] + offset
    win_tables.append(w)

    xpos_max = offset + d["end"].max()
    chrom_ticks.append((chrom, (xpos_min + xpos_max) / 2.0))
    offset = xpos_max + 1
    xpos_min = offset

win_all = pd.concat(win_tables, ignore_index=True)

# keep only informative windows: coverage>0, finite mean, and mean>0 (drop zeros)
win_all = win_all[(win_all["bp_covered"] > 0) &
                  (win_all["y_mean"].notna()) & (win_all["y_mean"] > 0) &
                  (win_all["y_lo"].notna()) & (win_all["y_hi"].notna())].copy()

if win_all.empty:
    raise RuntimeError("No windows with positive mean and coverage. Adjust WIN_BP/STEP_BP or check input.")

win_all.to_csv(os.path.join(OUTDIR, f"tmrca_windows_{WIN_BP}_{STEP_BP}.csv"), index=False)

# ============ 4) PLOTS ============
# Manhattan (linear)
out_lin = os.path.join(OUTDIR, f"manhattan_tmrca_LINEAR_win{WIN_BP}_step{STEP_BP}.png")
plot_manhattan(win_all, scaffolds_ordered, chrom_ticks, yscale="linear", out_png=out_lin,
               title_suffix="(linear scale)")
# Manhattan (log)
out_log = os.path.join(OUTDIR, f"manhattan_tmrca_LOG_win{WIN_BP}_step{STEP_BP}.png")
plot_manhattan(win_all, scaffolds_ordered, chrom_ticks, yscale="log", out_png=out_log,
               title_suffix="(log scale)")

# Manhattan: top-1 per scaffold annotated
out_top1_lin = os.path.join(OUTDIR, f"manhattan_TOP1_per_scaffold_LINEAR_win{WIN_BP}_step{STEP_BP}.png")
plot_manhattan_top1_per_scaffold(win_all, scaffolds_ordered, chrom_ticks, yscale="linear", out_png=out_top1_lin)
out_top1_log = os.path.join(OUTDIR, f"manhattan_TOP1_per_scaffold_LOG_win{WIN_BP}_step{STEP_BP}.png")
plot_manhattan_top1_per_scaffold(win_all, scaffolds_ordered, chrom_ticks, yscale="log", out_png=out_top1_log)

# NEW: line plots with CI ribbon
out_line_lin = os.path.join(OUTDIR, f"line_with_CI_LINEAR_win{WIN_BP}_step{STEP_BP}.png")
plot_line_with_ci(win_all, scaffolds_ordered, chrom_ticks, yscale="linear", out_png=out_line_lin)
out_line_log = os.path.join(OUTDIR, f"line_with_CI_LOG_win{WIN_BP}_step{STEP_BP}.png")
plot_line_with_ci(win_all, scaffolds_ordered, chrom_ticks, yscale="log", out_png=out_line_log)

# Optional extra: genome-wide histogram (zeros removed)
vals = tm["tmrca_mean"].replace(0, np.nan).dropna()
fig, ax = plt.subplots(figsize=(7,4))
ax.hist(np.log10(vals), bins=60, color="#747474", edgecolor="none")
ax.set_xlabel("log10(TMRCA mean, generations)")
ax.set_ylabel("Count")
ax.set_title("Genome-wide distribution of local TMRCA (zeros removed)")
fig.tight_layout()
hist_path = os.path.join(OUTDIR, "tmrca_hist_log10.png")
fig.savefig(hist_path, dpi=220)
plt.close(fig)

# ============ 5) TOP-K TABLE ============
topK = win_all.nlargest(TOP_K, "y_mean").copy()
topK["w_start_Mb"] = topK["w_start"]/1e6
topK["w_end_Mb"]   = topK["w_end"]/1e6
topK_csv = os.path.join(OUTDIR, f"top_{TOP_K}_windows.csv")
topK[["chrom","w_start","w_end","w_start_Mb","w_end_Mb","y_mean","y_lo","y_hi","bp_covered"]].to_csv(topK_csv, index=False)

# ============ 6) SLIDES ============
prs = Presentation()

# Title
s0 = prs.slides.add_slide(prs.slide_layouts[0])
s0.shapes.title.text = "ARGweaver TMRCA summary"
s0.placeholders[1].text = (
    f"Input file: {os.path.basename(INPUT_FILE)}\n"
    f"Windows: {WIN_BP:,} bp  |  Step: {STEP_BP:,} bp\n"
    f"Top {TOP_K} windows (global) highlighted in red\n"
    f"Fade rule: alpha={ALPHA_FAINT} if mean < {THRESH_FADE} generations"
)

# Global stats
s1 = prs.slides.add_slide(prs.slide_layouts[5])
s1.shapes.title.text = "Global TMRCA (generations): all vs zeros-removed"
rows, cols = summ_table.shape[0]+1, summ_table.shape[1]+1
tbl = s1.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(1.4)).table
tbl.cell(0,0).text = ""
for j, c in enumerate(summ_table.columns, start=1):
    tbl.cell(0,j).text = c
for i, (idx, row) in enumerate(summ_table.iterrows(), start=1):
    tbl.cell(i,0).text = str(idx)
    for j, c in enumerate(summ_table.columns, start=1):
        v = row[c]
        tbl.cell(i,j).text = f"{v:,.3g}" if isinstance(v,(int,float)) and not pd.isna(v) else str(v)

# Per-scaffold stats table (first 20 rows just to keep slide readable)
s2 = prs.slides.add_slide(prs.slide_layouts[5])
s2.shapes.title.text = "Per-scaffold summary (generations)"
sub = per_scaf.sort_values("chrom").head(20).copy()
rows, cols = len(sub)+1, len(sub.columns)
tbl2 = s2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(3.5)).table
for j,c in enumerate(sub.columns):
    tbl2.cell(0,j).text = c
for i, (_, r) in enumerate(sub.iterrows(), start=1):
    for j, c in enumerate(sub.columns):
        val = r[c]
        tbl2.cell(i,j).text = f"{val:,.3g}" if isinstance(val,(int,float)) and not pd.isna(val) else str(val)

# Manhattan LINEAR
s3 = prs.slides.add_slide(prs.slide_layouts[5])
s3.shapes.title.text = "Windowed TMRCA Manhattan (linear scale)"
s3.shapes.add_picture(out_lin, Inches(0.5), Inches(1.2), width=Inches(9))

# Manhattan LOG
s4 = prs.slides.add_slide(prs.slide_layouts[5])
s4.shapes.title.text = "Windowed TMRCA Manhattan (log scale)"
s4.shapes.add_picture(out_log, Inches(0.5), Inches(1.2), width=Inches(9))

# Manhattan: Top-1 per scaffold (linear)
s5 = prs.slides.add_slide(prs.slide_layouts[5])
s5.shapes.title.text = "Top-1 window per scaffold (linear scale)"
s5.shapes.add_picture(out_top1_lin, Inches(0.5), Inches(1.2), width=Inches(9))

# Line plots with CI
s6 = prs.slides.add_slide(prs.slide_layouts[5])
s6.shapes.title.text = "Windowed TMRCA line plot with CI (linear)"
s6.shapes.add_picture(out_line_lin, Inches(0.5), Inches(1.2), width=Inches(9))

s7 = prs.slides.add_slide(prs.slide_layouts[5])
s7.shapes.title.text = "Windowed TMRCA line plot with CI (log)"
s7.shapes.add_picture(out_line_log, Inches(0.5), Inches(1.2), width=Inches(9))

# Histogram slide
s8 = prs.slides.add_slide(prs.slide_layouts[5])
s8.shapes.title.text = "Genome-wide distribution of local TMRCA (zeros removed)"
s8.shapes.add_picture(os.path.join(OUTDIR, "tmrca_hist_log10.png"), Inches(0.5), Inches(1.2), width=Inches(9))

pptx_path = os.path.join(OUTDIR, "tmrca_summary_slides_v3.pptx")
prs.save(pptx_path)

print("Wrote:")
print("  -", os.path.join(OUTDIR, f"tmrca_windows_{WIN_BP}_{STEP_BP}.csv"))
print("  -", os.path.join(OUTDIR, f"top_{TOP_K}_windows.csv"))
print("  -", out_lin)
print("  -", out_log)
print("  -", out_top1_lin, "and", out_top1_log)
print("  -", out_line_lin, "and", out_line_log)
print("  -", os.path.join(OUTDIR, "tmrca_hist_log10.png"))
print("  -", pptx_path)
