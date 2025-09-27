# tmrca_pipeline.py
# ------------------------------------------------------------
# Read ARGweaver 6-column TMRCA track, summarize (all vs zeros-removed),
# window with overlap, Manhattan-style plot with top-K windows highlighted,
# and build a simple PPTX deck.
# ------------------------------------------------------------

import os
import math
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt

# Optional: pip install python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches
except Exception:
    Presentation = None  # guarded later

# --------------------
# USER SETTINGS
# --------------------
INPUT_FILE = "All_tmrca_60_run.txt"  # tab-delimited, 6 columns
OUTDIR     = "tmrca_reports"

# Coordinates
ZERO_BASED_START = True  # if True, add +1 to start when exporting display columns

# Windowing
WIN_BP  = 1000000        # window size (bp)
STEP_BP = 100000         # step/overlap (bp). Overlap = WIN_BP - STEP_BP

# Plotting / selection
TOP_K   = 10             # highlight top-K windows by mean TMRCA
SEED    = 13

# --------------------
# 0) Setup
# --------------------
os.makedirs(OUTDIR, exist_ok=True)
np.random.seed(SEED)

# --------------------
# 1) Read & name columns
# --------------------
colspec = ["chrom","start","end","tmrca_mean","tmrca_lo","tmrca_hi"]
tm = pd.read_csv(INPUT_FILE, sep=r"\s+|\t", engine="python", header=None, names=colspec)

# Coerce numeric
for c in ["start","end","tmrca_mean","tmrca_lo","tmrca_hi"]:
    tm[c] = pd.to_numeric(tm[c], errors="coerce")

tm = tm.dropna(subset=["chrom","start","end","tmrca_mean"]).reset_index(drop=True)

# Optional display coords
tm["start_disp"] = tm["start"] + 1 if ZERO_BASED_START else tm["start"]
tm["end_disp"]   = tm["end"]

# --------------------
# 2) Global summaries (all vs zeros-removed) + per-scaffold
# --------------------
def summarize_series(x: pd.Series, label: str) -> pd.DataFrame:
    x = pd.to_numeric(x, errors="coerce")
    desc = {
        "n": int(x.size),
        "n_nonzero": int((x>0).sum()),
        "min": float(x.min()),
        "q1": float(x.quantile(0.25)),
        "median": float(x.median()),
        "mean": float(x.mean()),
        "q3": float(x.quantile(0.75)),
        "max": float(x.max()),
        "sd": float(x.std())
    }
    return pd.DataFrame([desc], index=[label])

summ_all = summarize_series(tm["tmrca_mean"], "all_values")
summ_no0 = summarize_series(tm.loc[tm["tmrca_mean"]>0, "tmrca_mean"], "nonzero_only")
summ_table = pd.concat([summ_all, summ_no0])
summ_table.to_csv(os.path.join(OUTDIR, "summary_global_tmrca.csv"))

per_scaf = tm.groupby("chrom")["tmrca_mean"].agg(
    n="size",
    n_nonzero=lambda s: (s>0).sum(),
    min="min",
    q1=lambda s: s.quantile(0.25),
    median="median",
    mean="mean",
    q3=lambda s: s.quantile(0.75),
    max="max",
    sd="std"
).reset_index()
per_scaf.to_csv(os.path.join(OUTDIR, "summary_by_scaffold.csv"), index=False)

# --------------------
# 3) Windowing with overlap (coverage-weighted mean)
# --------------------
scaffolds = sorted(tm["chrom"].unique().tolist())

def window_reduce(df_scaf: pd.DataFrame, win_bp: int, step_bp: int) -> pd.DataFrame:
    start_min = int(max(0, df_scaf["start"].min()))
    end_max   = int(df_scaf["end"].max())
    starts = np.arange(start_min, max(start_min, end_max - win_bp + 1), step_bp, dtype=int)
    rows = []

    d = df_scaf.sort_values(["start","end"]).reset_index(drop=True)
    s_arr = d["start"].to_numpy()
    e_arr = d["end"].to_numpy()
    y_arr = d["tmrca_mean"].to_numpy()

    j0 = 0
    for wstart in starts:
        wend = wstart + win_bp
        while j0 < len(d) and e_arr[j0] <= wstart:
            j0 += 1
        j = j0
        total_cov = 0
        numer = 0.0
        while j < len(d) and s_arr[j] < wend:
            ovl = max(0, min(e_arr[j], wend) - max(s_arr[j], wstart))
            if ovl > 0:
                total_cov += ovl
                numer += ovl * y_arr[j]
            j += 1
        mean_y = numer/total_cov if total_cov>0 else np.nan
        rows.append((wstart, wend, mean_y, total_cov))
    out = pd.DataFrame(rows, columns=["w_start","w_end","y_mean","bp_covered"])
    return out

win_tables = []
offset = 0
chrom_ticks = []
xpos_min = offset

for chrom in scaffolds:
    d = tm.loc[tm["chrom"]==chrom, ["start","end","tmrca_mean"]].copy()
    if d.empty:
        continue
    w = window_reduce(d, WIN_BP, STEP_BP)
    w["chrom"] = chrom
    w["mid"] = (w["w_start"] + w["w_end"])/2.0
    w["mid_global"] = w["mid"] + offset
    win_tables.append(w)

    xpos_max = offset + d["end"].max()
    chrom_ticks.append((chrom, (xpos_min + xpos_max)/2.0))
    offset = xpos_max + 1
    xpos_min = offset

win_all = pd.concat(win_tables, ignore_index=True)

# Drop windows with no coverage or NaN mean (prevents NaN alpha)
win_all = win_all[(win_all["bp_covered"] > 0) & (win_all["y_mean"].notna())].copy()
win_all.to_csv(os.path.join(OUTDIR, f"tmrca_windows_{WIN_BP}_{STEP_BP}.csv"), index=False)

if win_all.empty:
    raise RuntimeError("No windows with coverage—adjust WIN_BP/STEP_BP or check input.")

# --------------------
# 4) Manhattan-style plot (alt scaffold shading; top-K in red; bolder=older)
# --------------------
# rank to find top-K
topK = win_all.nlargest(TOP_K, "y_mean").copy()
top_ids = set(topK.index.tolist())

# visuals scaled by TMRCA (bolder = older)
y_pos = win_all["y_mean"].clip(lower=1e-9).to_numpy()   # >0 for log axis
y_log = np.log10(y_pos)
rng = float(np.nanmax(y_log) - np.nanmin(y_log))
y_scaled = (y_log - np.nanmin(y_log)) / (rng if rng>0 else 1.0)
sizes  = 10 + 40*y_scaled
alphas = np.clip(0.2 + 0.6*y_scaled, 0, 1)

# two colors alternating by scaffold
base_colors = ["#4C78A8", "#A0A0A0"]  # blue / gray
chrom_to_color = {chrom: base_colors[i % 2] for i, chrom in enumerate(scaffolds)}
colors = np.array([chrom_to_color[c] for c in win_all["chrom"]], dtype=object)

fig, ax = plt.subplots(figsize=(14, 4))
mask_top = win_all.index.isin(top_ids)

# non-top
ax.scatter(win_all.loc[~mask_top, "mid_global"],
           y_pos[~mask_top],
           s=sizes[~mask_top],
           c=colors[~mask_top],
           alpha=alphas[~mask_top],
           linewidths=0)

# top-K in red
ax.scatter(win_all.loc[mask_top, "mid_global"],
           y_pos[mask_top],
           s=60, c="red", alpha=1.0, linewidths=0, zorder=3)

ax.set_yscale("log")
ax.set_ylabel("Windowed TMRCA mean (generations, log10)")
ax.set_xlabel("Genomic position (scaffolds concatenated)")

# alternate background bands per scaffold
for i, chrom in enumerate(scaffolds):
    scaf_mask = win_all["chrom"]==chrom
    if not scaf_mask.any():
        continue
    x_min = win_all.loc[scaf_mask, "mid_global"].min() - WIN_BP/2
    x_max = win_all.loc[scaf_mask, "mid_global"].max() + WIN_BP/2
    if i % 2 == 1:
        ax.axvspan(x_min, x_max, color="#000000", alpha=0.05, lw=0, zorder=0)

# x ticks at scaffold centers
ax.set_xticks([t for _, t in chrom_ticks])
ax.set_xticklabels([c for c, _ in chrom_ticks], rotation=90, fontsize=8)
ax.set_title(f"Windowed TMRCA Manhattan plot (win={WIN_BP:,} bp, step={STEP_BP:,} bp)\nTop {TOP_K} windows in red")
fig.tight_layout()
figpath_manh = os.path.join(OUTDIR, f"manhattan_tmrca_win{WIN_BP}_step{STEP_BP}.png")
fig.savefig(figpath_manh, dpi=220)
plt.close(fig)

# --------------------
# 5) Distribution & per-scaffold summaries (figures)
# --------------------
# Histogram of log10 TMRCA (zeros removed)
vals = tm["tmrca_mean"].replace(0, np.nan).dropna()
fig, ax = plt.subplots(figsize=(7,4))
ax.hist(np.log10(vals), bins=60, color="#777777", edgecolor="none")
ax.set_xlabel("log10(TMRCA mean, generations)")
ax.set_ylabel("Count")
ax.set_title("Genome-wide distribution of local TMRCA (zeros removed)")
fig.tight_layout()
figpath_hist = os.path.join(OUTDIR, "tmrca_hist_log10.png")
fig.savefig(figpath_hist, dpi=220)
plt.close(fig)

# Per-scaffold median & IQR on log scale
agg = tm.groupby("chrom")["tmrca_mean"].agg(
    med="median", q1=lambda s: s.quantile(0.25), q3=lambda s: s.quantile(0.75)
).reset_index().sort_values("med")
y = np.arange(len(agg))
fig, ax = plt.subplots(figsize=(10, max(5, len(agg)*0.25)))
ax.hlines(y, agg["q1"], agg["q3"], color="#B0B0B0")
ax.plot(agg["med"], y, "ko", ms=3)
ax.set_xscale("log")
ax.set_yticks(y)
ax.set_yticklabels(agg["chrom"], fontsize=7)
ax.set_xlabel("TMRCA mean (generations, log10)")
ax.set_title("Per-scaffold TMRCA: median and IQR")
fig.tight_layout()
figpath_scafsum = os.path.join(OUTDIR, "tmrca_per_scaffold_summary.png")
fig.savefig(figpath_scafsum, dpi=220)
plt.close(fig)

# --------------------
# 6) Export top-K windows table
# --------------------
topK_out = win_all.nlargest(TOP_K, "y_mean").copy()
topK_out["w_start_Mb"] = topK_out["w_start"]/1e6
topK_out["w_end_Mb"]   = topK_out["w_end"]/1e6
topK_out[["chrom","w_start","w_end","w_start_Mb","w_end_Mb","y_mean","bp_covered"]].to_csv(
    os.path.join(OUTDIR, f"top_{TOP_K}_windows.csv"), index=False
)

# --------------------
# 7) Simple PPTX slide deck (optional)
# --------------------
if Presentation is not None:
    prs = Presentation()
    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "ARGweaver TMRCA summary"
    slide.placeholders[1].text = (
        f"Input: {os.path.basename(INPUT_FILE)}\n"
        f"Columns: chrom, start, end, tmrca_mean, tmrca_lo, tmrca_hi\n"
        f"Windows: {WIN_BP:,} bp; step: {STEP_BP:,} bp; Top-{TOP_K} windows highlighted"
    )

    # Global summaries
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Global TMRCA (generations): all vs zeros-removed"
    rows, cols = summ_table.shape[0]+1, summ_table.shape[1]+1
    table = s2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(9), Inches(1.4)).table
    table.cell(0,0).text = ""
    for j, c in enumerate(summ_table.columns, start=1):
        table.cell(0,j).text = c
    for i, (idx, row) in enumerate(summ_table.iterrows(), start=1):
        table.cell(i,0).text = str(idx)
        for j, c in enumerate(summ_table.columns, start=1):
            v = row[c]
            table.cell(i,j).text = f"{v:,.3g}" if isinstance(v, (int,float)) and not pd.isna(v) else str(v)

    # Manhattan
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.title.text = "Windowed TMRCA Manhattan plot"
    s3.shapes.add_picture(figpath_manh, Inches(0.5), Inches(1.2), width=Inches(9))

    # Per-scaffold summary
    s4 = prs.slides.add_slide(prs.slide_layouts[5])
    s4.shapes.title.text = "Per-scaffold median & IQR"
    s4.shapes.add_picture(figpath_scafsum, Inches(0.5), Inches(1.2), width=Inches(9))

    # Distribution
    s5 = prs.slides.add_slide(prs.slide_layouts[5])
    s5.shapes.title.text = "Genome-wide distribution (zeros removed)"
    s5.shapes.add_picture(figpath_hist, Inches(0.5), Inches(1.2), width=Inches(9))

    pptx_path = os.path.join(OUTDIR, "tmrca_summary_slides.pptx")
    prs.save(pptx_path)
    print(f"Wrote slides: {pptx_path}")
else:
    print("python-pptx not installed; skipping PPTX creation. Install with: pip install python-pptx")

print(f"All outputs in: {OUTDIR}")

