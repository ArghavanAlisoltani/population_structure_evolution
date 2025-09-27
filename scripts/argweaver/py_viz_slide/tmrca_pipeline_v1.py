# pip install pandas numpy matplotlib python-pptx
# tmrca_pipeline.py
# ------------------------------------------------------------
# Read ARGweaver TMRCA track, summarize, window, plot Manhattan,
# mark top-10 windows, and export a PPTX slide deck.
# ------------------------------------------------------------

import os
import math
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.ticker import ScalarFormatter
from collections import defaultdict

# Optional: pip install python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None  # will guard later

# --------------------
# USER SETTINGS
# --------------------
INPUT_FILE = "All_tmrca_60_run.txt"      # tab-delimited
OUTDIR     = "tmrca_reports"
ZERO_BASED_START = True                   # add +1 to start for 1-based display tables if you wish
WIN_BP     = 100_000                      # window size (bp)
STEP_BP    = 10_000                       # overlap control (step between window starts)
TOP_K      = 10                           # highlight top-K windows by mean TMRCA
SEED       = 13

os.makedirs(OUTDIR, exist_ok=True)

# --------------------
# 1) Read & name columns
# --------------------
colspec = ["chrom","start","end","tmrca_mean","tmrca_lo","tmrca_hi"]
tm = pd.read_csv(INPUT_FILE, sep=r"\s+|\t", engine="python", header=None, names=colspec)

# Optional 1-based start for display
tm["start_disp"] = tm["start"] + 1 if ZERO_BASED_START else tm["start"]
tm["end_disp"]   = tm["end"]

# Defensive numeric conversion
for c in ["start","end","tmrca_mean","tmrca_lo","tmrca_hi","start_disp","end_disp"]:
    tm[c] = pd.to_numeric(tm[c], errors="coerce")

# Drop pathological rows
tm = tm.dropna(subset=["chrom","start","end","tmrca_mean"]).reset_index(drop=True)

# --------------------
# 2) Global summaries (with and without zeros)
# --------------------
def summarize_series(x: pd.Series, label: str) -> pd.DataFrame:
    x = pd.to_numeric(x, errors="coerce")
    desc = {
        "n": x.size,
        "n_nonzero": (x!=0).sum(),
        "min": x.min(),
        "q1": x.quantile(0.25),
        "median": x.median(),
        "mean": x.mean(),
        "q3": x.quantile(0.75),
        "max": x.max(),
        "sd": x.std()
    }
    return pd.DataFrame([desc], index=[label])

summ_all   = summarize_series(tm["tmrca_mean"], "all_values")
summ_no0   = summarize_series(tm.loc[tm["tmrca_mean"]>0, "tmrca_mean"], "nonzero_only")
summ_table = pd.concat([summ_all, summ_no0])
summ_table.to_csv(os.path.join(OUTDIR, "summary_global_tmrca.csv"))

# Per-scaffold summaries
per_scaf = tm.groupby("chrom")["tmrca_mean"].agg(
    n="size", n_nonzero=lambda s: (s>0).sum(),
    min="min", q1=lambda s: s.quantile(0.25),
    median="median", mean="mean",
    q3=lambda s: s.quantile(0.75), max="max", sd="std"
).reset_index()
per_scaf.to_csv(os.path.join(OUTDIR, "summary_by_scaffold.csv"), index=False)

# --------------------
# 3) Windowing with overlap (mean TMRCA per window)
# --------------------
# build per-scaffold windows so the x-axis ordering is contiguous by scaffold
np.random.seed(SEED)
scaffolds = sorted(tm["chrom"].unique().tolist())

def window_reduce(df_scaf: pd.DataFrame, win_bp: int, step_bp: int) -> pd.DataFrame:
    # df_scaf contains rows for a single scaffold
    start_min = int(df_scaf["start"].min())
    start_min = max(0, start_min)
    end_max   = int(df_scaf["end"].max())
    starts = np.arange(start_min, end_max - win_bp + 1, step_bp, dtype=int)
    rows = []
    # index intervals for speed
    df_scaf = df_scaf.sort_values(["start","end"]).reset_index(drop=True)
    s_arr = df_scaf["start"].values
    e_arr = df_scaf["end"].values
    y_arr = df_scaf["tmrca_mean"].values

    j0 = 0
    for wstart in starts:
        wend = wstart + win_bp
        # advance left pointer
        while j0 < len(df_scaf) and e_arr[j0] <= wstart:
            j0 += 1
        j = j0
        total_cov = 0
        numer = 0.0
        while j < len(df_scaf) and s_arr[j] < wend:
            ovl = max(0, min(e_arr[j], wend) - max(s_arr[j], wstart))
            if ovl > 0:
                total_cov += ovl
                numer += ovl * y_arr[j]
            j += 1
        mean_y = numer/total_cov if total_cov>0 else np.nan
        rows.append((wstart, wend, mean_y, total_cov))
    out = pd.DataFrame(rows, columns=["w_start","w_end","y_mean","bp_covered"])
    return out

win_tables = []
offset = 0
chrom_ticks = []  # for Manhattan x-axis tick labels
xpos_min = offset

for chrom in scaffolds:
    d = tm.loc[tm["chrom"]==chrom, ["start","end","tmrca_mean"]].copy()
    if d.empty:
        continue
    w = window_reduce(d, WIN_BP, STEP_BP)
    w["chrom"] = chrom
    # global x for Manhattan: cumulative position
    w["mid"] = (w["w_start"] + w["w_end"])/2.0
    w["mid_global"] = w["mid"] + offset
    win_tables.append(w)
    # tick positions
    xpos_max = offset + d["end"].max()
    chrom_ticks.append((chrom, (xpos_min + xpos_max)/2.0))
    offset = xpos_max + 1
    xpos_min = offset

win_all = pd.concat(win_tables, ignore_index=True)
win_all.to_csv(os.path.join(OUTDIR, f"tmrca_windows_{WIN_BP}_{STEP_BP}.csv"), index=False)

# --------------------
# 4) Manhattan-style plot of windowed means
#    - alternating scaffold colors
#    - alpha/size scaled by TMRCA (bolder = older)
#    - top-K windows highlighted in red
# --------------------
# rank to find top-K
topK = win_all.nlargest(TOP_K, "y_mean").copy()
top_ids = set(topK.index.tolist())

# scale marker visuals
y_pos = win_all["y_mean"].copy()
y_pos = y_pos.clip(lower=1e-9)  # avoid <=0 on log axis
y_log = np.log10(y_pos)
y_log = (y_log - np.nanmin(y_log)) / (np.nanmax(y_log) - np.nanmin(y_log) + 1e-12)
sizes = 10 + 40*y_log          # 10..50
alphas = 0.2 + 0.6*y_log       # 0.2..0.8

# two base colors (binary palette), alternate by scaffold index
base_colors = ["#4C78A8", "#A0A0A0"]  # blue vs gray
chrom_to_color = {chrom: base_colors[i % 2] for i, chrom in enumerate(scaffolds)}
colors = [chrom_to_color[c] for c in win_all["chrom"]]

# draw
fig, ax = plt.subplots(figsize=(14, 4))
# non-top points
mask_top = win_all.index.isin(top_ids)
ax.scatter(win_all.loc[~mask_top, "mid_global"],
           y_pos[~mask_top],
           s=sizes[~mask_top],
           c=np.array(colors, dtype=object)[~mask_top],
           alpha=alphas[~mask_top],
           linewidths=0)

# top-K in red (bold)
ax.scatter(win_all.loc[mask_top, "mid_global"],
           y_pos[mask_top],
           s=60,
           c="red",
           alpha=1.0,
           linewidths=0,
           zorder=3)

# cosmetics
ax.set_yscale("log")
ax.set_ylabel("Windowed TMRCA mean (generations, log10)")
ax.set_xlabel("Genomic position (scaffolds concatenated)")
# vertical bands to alternate background like your example
for i, chrom in enumerate(scaffolds):
    # find min/max for this scaffold in global coords
    scaf_mask = win_all["chrom"]==chrom
    if not scaf_mask.any():
        continue
    x_min = win_all.loc[scaf_mask, "mid_global"].min() - WIN_BP/2
    x_max = win_all.loc[scaf_mask, "mid_global"].max() + WIN_BP/2
    if i % 2 == 1:
        ax.axvspan(x_min, x_max, color="#000000", alpha=0.05, lw=0, zorder=0)

# x ticks at scaffold centers
ax.set_xticks([t for _, t in chrom_ticks])
ax.set_xticklabels([c for c, _ in chrom_ticks], rotation=90, fontsize=8)
ax.set_title(f"Windowed TMRCA Manhattan plot (win={WIN_BP:,} bp, step={STEP_BP:,} bp)\nTop {TOP_K} windows in red")
fig.tight_layout()
figpath_manh = os.path.join(OUTDIR, f"manhattan_tmrca_win{WIN_BP}_step{STEP_BP}.png")
fig.savefig(figpath_manh, dpi=200)
plt.close(fig)

# --------------------
# 5) Additional figures: distribution & per-scaffold violin-like (boxen surrogate)
# --------------------
# Histogram (log10)
fig, ax = plt.subplots(figsize=(7,4))
vals = tm["tmrca_mean"].replace(0, np.nan).dropna()
ax.hist(np.log10(vals), bins=60, color="#777777", edgecolor="none")
ax.set_xlabel("log10(TMRCA mean, generations)")
ax.set_ylabel("Count")
ax.set_title("Genome-wide distribution of local TMRCA (zeros removed)")
fig.tight_layout()
figpath_hist = os.path.join(OUTDIR, "tmrca_hist_log10.png")
fig.savefig(figpath_hist, dpi=200)
plt.close(fig)

# Per-scaffold summary dot plot (median & IQR on log scale)
agg = tm.groupby("chrom")["tmrca_mean"].agg(["median", lambda s: s.quantile(0.25), lambda s: s.quantile(0.75)]).reset_index()
agg.columns = ["chrom","med","q1","q3"]
agg = agg.sort_values("med")
y = np.arange(len(agg))
fig, ax = plt.subplots(figsize=(10, 8))
ax.hlines(y, agg["q1"], agg["q3"], color="#B0B0B0")
ax.plot(agg["med"], y, "ko", ms=3)
ax.set_xscale("log")
ax.set_yticks(y)
ax.set_yticklabels(agg["chrom"], fontsize=7)
ax.set_xlabel("TMRCA mean (generations, log10)")
ax.set_title("Per-scaffold TMRCA: median and IQR")
fig.tight_layout()
figpath_scafsum = os.path.join(OUTDIR, "tmrca_per_scaffold_summary.png")
fig.savefig(figpath_scafsum, dpi=200)
plt.close(fig)

# --------------------
# 6) Export top-10 windows table
# --------------------
topK_out = win_all.nlargest(TOP_K, "y_mean").copy()
# add nice Mb columns and scaffold-relative coordinates
topK_out["w_start_Mb"] = topK_out["w_start"]/1e6
topK_out["w_end_Mb"]   = topK_out["w_end"]/1e6
topK_out[["chrom","w_start","w_end","w_start_Mb","w_end_Mb","y_mean","bp_covered"]].to_csv(
    os.path.join(OUTDIR, f"top_{TOP_K}_windows.csv"), index=False
)

# --------------------
# 7) Simple PPTX report (if python-pptx is available)
# --------------------
if Presentation is not None:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "ARGweaver TMRCA summary"
    slide.placeholders[1].text = (
        f"Input: {os.path.basename(INPUT_FILE)}\n"
        f"Windows: {WIN_BP:,} bp, step {STEP_BP:,} bp\n"
        f"Top-{TOP_K} windows highlighted"
    )

    # Slide: global stats
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Global TMRCA summaries (generations)"
    # add table
    rows, cols = summ_table.shape[0]+1, summ_table.shape[1]+1
    table = s2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(1.2)).table
    table.cell(0,0).text = ""
    for j, c in enumerate(summ_table.columns, start=1):
        table.cell(0,j).text = c
    for i, (idx, row) in enumerate(summ_table.iterrows(), start=1):
        table.cell(i,0).text = str(idx)
        for j, c in enumerate(summ_table.columns, start=1):
            val = row[c]
            if isinstance(val, (int,float)) and not pd.isna(val):
                table.cell(i,j).text = f"{val:,.3g}"
            else:
                table.cell(i,j).text = str(val)

    # Slide: Manhattan
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.title.text = "Windowed TMRCA Manhattan plot"
    s3.shapes.add_picture(figpath_manh, Inches(0.5), Inches(1.2), width=Inches(9))

    # Slide: per-scaffold summary
    s4 = prs.slides.add_slide(prs.slide_layouts[5])
    s4.shapes.title.text = "Per-scaffold median & IQR"
    s4.shapes.add_picture(figpath_scafsum, Inches(0.5), Inches(1.2), width=Inches(9))

    # Slide: genome-wide distribution
    s5 = prs.slides.add_slide(prs.slide_layouts[5])
    s5.shapes.title.text = "Genome-wide distribution of local TMRCA (zeros removed)"
    s5.shapes.add_picture(figpath_hist, Inches(0.5), Inches(1.2), width=Inches(9))

    # Slide: top-K table
    s6 = prs.slides.add_slide(prs.slide_layouts[5])
    s6.shapes.title.text = f"Top {TOP_K} windows by mean TMRCA"
    rows = min(12, len(topK_out)+1)
    tbl = s6.shapes.add_table(rows, 6, Inches(0.5), Inches(1.2), Inches(9), Inches(2.5)).table
    headers = ["chrom","w_start (bp)","w_end (bp)","start (Mb)","end (Mb)","mean TMRCA"]
    for j,h in enumerate(headers):
        tbl.cell(0,j).text = h
    for i, (_, r) in enumerate(topK_out.head(rows-1).iterrows(), start=1):
        tbl.cell(i,0).text = r["chrom"]
        tbl.cell(i,1).text = f"{int(r['w_start']):,}"
        tbl.cell(i,2).text = f"{int(r['w_end']):,}"
        tbl.cell(i,3).text = f"{r['w_start_Mb']:.3f}"
        tbl.cell(i,4).text = f"{r['w_end_Mb']:.3f}"
        tbl.cell(i,5).text = f"{r['y_mean']:.3g}"

    pptx_path = os.path.join(OUTDIR, "tmrca_summary_slides.pptx")
    prs.save(pptx_path)
    print(f"Wrote slides: {pptx_path}")
else:
    print("python-pptx not installed; skipping slide generation.")
    print("Install with:  pip install python-pptx")

