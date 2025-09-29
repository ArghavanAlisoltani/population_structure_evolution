#!/usr/bin/env python3
"""
Compare two windowing methods for ARGweaver using:
  (A) classic uniformity (CV of window length and SNP counts)
  (B) density-based uniformity (length-weighted CV of SNP density) and dispersion (chi-square vs length)

Plus: paired Wilcoxon signed-rank tests (A vs B) for classic uniformity, CV_density, and dispersion D.

Inputs: two TSVs with columns (case-insensitive): Scaffold/CHROM, start, end, nsnp
Outputs (files in working directory):
  Tables:
    - scaffold_uniformity_comparison.tsv
    - uniformity_overall_summary.tsv
    - density_overall_summary.tsv
    - decisive_scaffolds.tsv
    - decisive_scaffolds_density.tsv
    - wilcoxon_tests_summary.tsv
  Figures:
    - uniformity_hist.png
    - uniformity_scatter.png
    - alldiff_bar.png
    - density_hist.png
    - density_scatter.png
    - density_alldiff_bar.png
    - dispersion_scatter.png
  Slides:
    - ARGweaver_window_choice_summary.pptx
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
import math

# =======================
# CONFIG (edit as needed)
# =======================
A_PATH = "windows_150Mb_equal_lenght.tsv"
B_PATH = "balanced_scaffold_windows_150Mb.tsv"

# “Decisive” thresholds for per-scaffold differences
DECISIVE_THRESH_CLASSIC = 0.05
DECISIVE_THRESH_DENSITY = 0.05

# If your windows are BED-style (0-based, half-open), set True (length = end - start)
BED_STYLE_INTERVALS = False

# =======================
# Output paths
# =======================
PER_SCAFFOLD_OUT           = "scaffold_uniformity_comparison.tsv"
OVERALL_OUT_CLASSIC        = "uniformity_overall_summary.tsv"
OVERALL_OUT_DENSITY        = "density_overall_summary.tsv"
DECISIVE_OUT_CLASSIC       = "decisive_scaffolds.tsv"
DECISIVE_OUT_DENSITY       = "decisive_scaffolds_density.tsv"
WILCOXON_SUMMARY_OUT       = "wilcoxon_tests_summary.tsv"

HIST_PNG_CLASSIC           = "uniformity_hist.png"
SCATTER_PNG_CLASSIC        = "uniformity_scatter.png"
ALLDIFF_PNG_CLASSIC        = "alldiff_bar.png"

HIST_PNG_DENSITY           = "density_hist.png"
SCATTER_PNG_DENSITY        = "density_scatter.png"
ALLDIFF_PNG_DENSITY        = "density_alldiff_bar.png"
SCATTER_PNG_DISPERSION     = "dispersion_scatter.png"

PPTX_OUT                   = "ARGweaver_window_choice_summary.pptx"

# =======================
# Helpers
# =======================
def load_and_normalize(path: str) -> pd.DataFrame:
    """Load windows TSV and normalize to: Scaffold, start, end, nsnp, length, density."""
    df = pd.read_csv(path, sep="\t")
    cols = {c.lower(): c for c in df.columns}

    def find(name):
        for k, v in cols.items():
            if k == name.lower():
                return v
        return None

    sc_col   = find("scaffold") or find("chrom") or list(df.columns)[0]
    start_col= find("start") or list(df.columns)[1]
    end_col  = find("end") or list(df.columns)[2]
    nsnp_col = find("nsnp") or find("n_snps") or find("count") or find("n")

    df = df.rename(columns={sc_col:"Scaffold", start_col:"start", end_col:"end"})
    if nsnp_col:
        df = df.rename(columns={nsnp_col:"nsnp"})
    else:
        df["nsnp"] = np.nan

    df["start"] = pd.to_numeric(df["start"], errors="coerce")
    df["end"]   = pd.to_numeric(df["end"],   errors="coerce")
    df["nsnp"]  = pd.to_numeric(df["nsnp"],  errors="coerce")

    if BED_STYLE_INTERVALS:
        df["length"] = df["end"] - df["start"]
    else:
        df["length"] = df["end"] - df["start"] + 1

    df["density"] = df["nsnp"] / df["length"]
    return df[["Scaffold","start","end","nsnp","length","density"]]

def cv(series: pd.Series) -> float:
    x = series.to_numpy(dtype=float)
    x = x[np.isfinite(x)]
    if x.size <= 1:
        return 0.0
    m = np.nanmean(x)
    if m == 0 or np.isnan(m):
        return np.nan
    return float(np.nanstd(x, ddof=1) / m)

def weighted_cv_density(lengths: np.ndarray, densities: np.ndarray) -> float:
    """Weighted CV of densities; weights = lengths. Unbiased-ish correction ((k-1)/k)."""
    L = np.asarray(lengths, dtype=float)
    d = np.asarray(densities, dtype=float)
    m = np.isfinite(L) & np.isfinite(d)
    L, d = L[m], d[m]
    k = L.size
    if k == 0:
        return np.nan
    mean_w = np.sum(L * d) / np.sum(L)
    if k == 1 or mean_w == 0 or not np.isfinite(mean_w):
        return 0.0
    num = np.sum(L * (d - mean_w)**2)
    denom = ((k - 1) / k) * np.sum(L)
    s_w = np.sqrt(num / denom) if denom > 0 else np.nan
    return float(s_w / mean_w) if (s_w is not np.nan and mean_w != 0) else np.nan

def dispersion_indices(lengths: np.ndarray, counts: np.ndarray):
    """Pearson chi-square dispersion vs constant density: returns (chi2, D, df)."""
    L = np.asarray(lengths, dtype=float)
    C = np.asarray(counts, dtype=float)
    m = np.isfinite(L) & np.isfinite(C) & (L > 0)
    L, C = L[m], C[m]
    k = L.size
    if k <= 1:
        return np.nan, np.nan, 0
    lam = np.sum(C) / np.sum(L) if np.sum(L) > 0 else np.nan
    if not np.isfinite(lam) or lam == 0:
        return np.nan, np.nan, k - 1
    E = lam * L
    v = E > 0
    if not np.any(v):
        return np.nan, np.nan, k - 1
    chi2 = np.sum((C[v] - E[v])**2 / E[v])
    D = chi2 / (k - 1)
    return float(chi2), float(D), int(k - 1)

def maxmin_ratio(series: pd.Series) -> float:
    x = series.to_numpy(dtype=float)
    x = x[np.isfinite(x)]
    if x.size == 0:
        return np.nan
    mn, mx = np.nanmin(x), np.nanmax(x)
    if mn == 0:
        return np.inf
    return float(mx / mn)

def per_scaffold_metrics(df: pd.DataFrame) -> pd.DataFrame:
    """Compute classic + density metrics per scaffold."""
    g = df.groupby("Scaffold", dropna=False)
    rows = []
    for sc, sub in g:
        L = sub["length"].to_numpy(dtype=float)
        C = sub["nsnp"].to_numpy(dtype=float)
        d = sub["density"].to_numpy(dtype=float)

        cv_len   = cv(sub["length"])
        cv_nsnp  = cv(sub["nsnp"])
        uniform  = cv_len if np.isnan(cv_nsnp) else 0.5 * (cv_len + cv_nsnp)

        cv_dens  = weighted_cv_density(L, d)
        chi2, D, dfree = dispersion_indices(L, C)

        rows.append(dict(
            Scaffold=sc,
            n_windows=int(len(sub)),
            total_len=float(np.nansum(L)),
            total_snps=float(np.nansum(C)),
            median_len=float(np.nanmedian(L)) if len(sub) else np.nan,
            median_nsnp=float(np.nanmedian(C)) if len(sub) else np.nan,
            cv_len=cv_len,
            cv_nsnp=cv_nsnp,
            uniformity_score=uniform,
            mean_density=float(np.nansum(C) / np.nansum(L)) if np.nansum(L) > 0 else np.nan,
            cv_density=cv_dens,
            chi2=chi2,
            dispersion_D=D,
            df=dfree,
            maxmin_len=maxmin_ratio(sub["length"]),
            maxmin_nsnp=maxmin_ratio(sub["nsnp"]),
        ))
    return pd.DataFrame(rows)

def overall_uniformity(met: pd.DataFrame) -> float:
    df2 = met.dropna(subset=["uniformity_score"]).copy()
    if df2.empty: return np.nan
    return float(np.average(df2["uniformity_score"], weights=df2["n_windows"]))

def overall_density_cv(met: pd.DataFrame) -> float:
    df2 = met.dropna(subset=["cv_density", "total_len"]).copy()
    if df2.empty: return np.nan
    return float(np.average(df2["cv_density"], weights=df2["total_len"]))

def overall_dispersion_D(met: pd.DataFrame) -> float:
    df2 = met.dropna(subset=["dispersion_D", "n_windows"]).copy()
    if df2.empty: return np.nan
    return float(np.average(df2["dispersion_D"], weights=df2["n_windows"]))

# ---------- Wilcoxon signed-rank (paired) ----------
def _norm_sf(z):
    # survival function (1 - CDF) for standard normal using erfc
    return 0.5 * math.erfc(z / math.sqrt(2.0))

def wilcoxon_signed_rank(diff):
    """
    Paired Wilcoxon signed-rank test for a 2-sided H0: median(diff)=0.
    Returns dict with n, W_plus, W_minus, Tplus (sum ranks positive), z, p_two_sided.
    Tries SciPy if available; otherwise normal approximation with continuity correction.
    """
    # Mask finite, nonzero diffs
    d = np.asarray(diff, dtype=float)
    d = d[np.isfinite(d)]
    d = d[d != 0]
    n = d.size
    if n == 0:
        return dict(n=0, Tplus=np.nan, z=np.nan, p_two_sided=np.nan, effect_r=np.nan)

    # Try SciPy if present
    try:
        from scipy.stats import wilcoxon
        stat, p = wilcoxon(d, zero_method="wilcox", correction=True, alternative="two-sided", mode="auto")
        # effect size r = z / sqrt(n); SciPy doesn't return z, so approximate from p (two-sided)
        # Convert p to |z| via inverse survival (approx). For small p, z ≈ Φ^{-1}(1 - p/2).
        # We'll compute with erfcinv if available; else skip r.
        try:
            from math import erfcinv
            z_abs = math.sqrt(2.0) * erfcinv(p * 2.0)
            r = z_abs / math.sqrt(n)
        except Exception:
            r = np.nan
        return dict(n=int(n), Tplus=float(stat), z=np.nan, p_two_sided=float(p), effect_r=float(r))
    except Exception:
        pass

    # Manual implementation (no SciPy): rank absolute diffs, sum ranks for positive, normal approx
    absd = np.abs(d)
    # Handle ties by average ranks
    order = np.argsort(absd)
    ranks = np.empty_like(absd, dtype=float)
    i = 0
    R = 1
    while i < n:
        j = i
        while j + 1 < n and absd[order][j+1] == absd[order][i]:
            j += 1
        avg_rank = (R + (R + (j - i))) / 2.0
        ranks[order][i:j+1] = avg_rank
        R += (j - i + 1)
        i = j + 1

    Tplus = float(np.sum(ranks[d > 0]))
    # Expected mean/var under H0
    mean = n * (n + 1) / 4.0
    var  = n * (n + 1) * (2 * n + 1) / 24.0
    # Continuity correction
    cc = 0.5 * np.sign(Tplus - mean)
    z = (Tplus - mean - cc) / math.sqrt(var) if var > 0 else np.nan
    p_two = 2.0 * _norm_sf(abs(z)) if np.isfinite(z) else np.nan
    r = abs(z) / math.sqrt(n) if np.isfinite(z) else np.nan
    return dict(n=int(n), Tplus=Tplus, z=float(z), p_two_sided=float(p_two), effect_r=float(r))

# =======================
# Load & compute metrics
# =======================
A = load_and_normalize(A_PATH)
B = load_and_normalize(B_PATH)

A_met = per_scaffold_metrics(A)
B_met = per_scaffold_metrics(B)

# Merge per-scaffold (classic + density/dispersion)
comp = (A_met.rename(columns=lambda c: f"{c}_A" if c != "Scaffold" else c)
        .merge(B_met.rename(columns=lambda c: f"{c}_B" if c != "Scaffold" else c),
               on="Scaffold", how="outer"))

# Winners (classic / density / dispersion)
comp["better_classic"] = np.where(
    (comp["uniformity_score_A"].notna() & comp["uniformity_score_B"].notna() &
     (comp["uniformity_score_A"] < comp["uniformity_score_B"])),
    "A",
    np.where(
        (comp["uniformity_score_A"].notna() & comp["uniformity_score_B"].notna() &
         (comp["uniformity_score_B"] < comp["uniformity_score_A"])),
        "B",
        "tie/NA"
    )
)
comp["better_density"] = np.where(
    (comp["cv_density_A"].notna() & comp["cv_density_B"].notna() &
     (comp["cv_density_A"] < comp["cv_density_B"])),
    "A",
    np.where(
        (comp["cv_density_A"].notna() & comp["cv_density_B"].notna() &
         (comp["cv_density_B"] < comp["cv_density_A"])),
        "B",
        "tie/NA"
    )
)
comp["better_dispersion"] = np.where(
    (comp["dispersion_D_A"].notna() & comp["dispersion_D_B"].notna() &
     (comp["dispersion_D_A"] < comp["dispersion_D_B"])),
    "A",
    np.where(
        (comp["dispersion_D_A"].notna() & comp["dispersion_D_B"].notna() &
         (comp["dispersion_D_B"] < comp["dispersion_D_A"])),
        "B",
        "tie/NA"
    )
)

# Save per-scaffold comparison
comp.to_csv(PER_SCAFFOLD_OUT, sep="\t", index=False)

# Overall summaries
overall_A_classic = overall_uniformity(A_met)
overall_B_classic = overall_uniformity(B_met)
summary_classic = pd.DataFrame({
    "Dataset": ["A (windows_100Mb_equal_lenght.tsv)", "B (balanced_scaffold_windows_100Mb.tsv)"],
    "Overall_uniformity_score (lower=better)": [overall_A_classic, overall_B_classic]
})
summary_classic.to_csv(OVERALL_OUT_CLASSIC, sep="\t", index=False)

overall_A_density = overall_density_cv(A_met)
overall_B_density = overall_density_cv(B_met)
overall_A_disp    = overall_dispersion_D(A_met)
overall_B_disp    = overall_dispersion_D(B_met)
summary_density = pd.DataFrame({
    "Dataset": ["A (windows_100Mb_equal_lenght.tsv)", "B (balanced_scaffold_windows_100Mb.tsv)"],
    "Overall_CV_density (length-weighted, lower=better)": [overall_A_density, overall_B_density],
    "Overall_dispersion_D (window-weighted, lower≈better)": [overall_A_disp, overall_B_disp],
})
summary_density.to_csv(OVERALL_OUT_DENSITY, sep="\t", index=False)

# Decisive scaffolds (classic & density)
dec_c = comp.copy()
dec_c["diff_classic"] = dec_c["uniformity_score_B"] - dec_c["uniformity_score_A"]
dec_c["abs_diff_classic"] = dec_c["diff_classic"].abs()
dec_c["decisive_winner_classic"] = np.where(
    dec_c["abs_diff_classic"] > DECISIVE_THRESH_CLASSIC,
    np.where(dec_c["diff_classic"] < 0, "A", "B"),
    "none")
decisive_classic = dec_c[dec_c["decisive_winner_classic"] != "none"] \
                     .sort_values("abs_diff_classic", ascending=False)
decisive_classic.to_csv(DECISIVE_OUT_CLASSIC, sep="\t", index=False)

dec_d = comp.copy()
dec_d["diff_density"] = dec_d["cv_density_B"] - dec_d["cv_density_A"]
dec_d["abs_diff_density"] = dec_d["diff_density"].abs()
dec_d["decisive_winner_density"] = np.where(
    dec_d["abs_diff_density"] > DECISIVE_THRESH_DENSITY,
    np.where(dec_d["diff_density"] < 0, "A", "B"),
    "none")
decisive_density = dec_d[dec_d["decisive_winner_density"] != "none"] \
                     .sort_values("abs_diff_density", ascending=False)
decisive_density.to_csv(DECISIVE_OUT_DENSITY, sep="\t", index=False)

# =======================
# Paired Wilcoxon tests
# =======================
def paired_vector(comp_df, colA, colB):
    """Extract paired differences for scaffolds with finite A and B values."""
    m = comp_df[[colA, colB]].dropna()
    a = m[colA].to_numpy(dtype=float)
    b = m[colB].to_numpy(dtype=float)
    return b - a  # B minus A

tests = []
for label, (ca, cb) in {
    "classic_uniformity": ("uniformity_score_A", "uniformity_score_B"),
    "cv_density": ("cv_density_A", "cv_density_B"),
    "dispersion_D": ("dispersion_D_A", "dispersion_D_B"),
}.items():
    d = paired_vector(comp, ca, cb)
    res = wilcoxon_signed_rank(d)
    # Add descriptive stats
    med = float(np.nanmedian(d)) if d.size else np.nan
    mean = float(np.nanmean(d)) if d.size else np.nan
    tests.append(dict(
        metric=label,
        n=res["n"],
        Tplus=res["Tplus"],
        z=res["z"],
        p_two_sided=res["p_two_sided"],
        effect_r=res["effect_r"],     # r ≈ z / sqrt(n)
        median_delta=med,             # B - A
        mean_delta=mean
    ))

wilcoxon_summary = pd.DataFrame(tests)
wilcoxon_summary.to_csv(WILCOXON_SUMMARY_OUT, sep="\t", index=False)

# =======================
# Plots
# =======================
# Classic hist
plt.figure()
plt.hist(comp["uniformity_score_A"].dropna(), bins=30, alpha=0.5, label="A")
plt.hist(comp["uniformity_score_B"].dropna(), bins=30, alpha=0.5, label="B")
plt.xlabel("Classic uniformity score (lower = better)")
plt.ylabel("Count of scaffolds")
plt.title("Classic uniformity distribution")
plt.legend()
plt.savefig(HIST_PNG_CLASSIC, bbox_inches="tight")
plt.close()

# Classic scatter A vs B
plt.figure()
plt.scatter(comp["uniformity_score_A"], comp["uniformity_score_B"], s=10)
vals = np.concatenate([
    comp["uniformity_score_A"].dropna().to_numpy(),
    comp["uniformity_score_B"].dropna().to_numpy()
])
if vals.size > 0:
    lo, hi = float(np.nanmin(vals)), float(np.nanmax(vals))
    plt.plot([lo, hi], [lo, hi])
plt.xlabel("A classic uniformity")
plt.ylabel("B classic uniformity")
plt.title("Per-scaffold: classic A vs B (y = x)")
plt.savefig(SCATTER_PNG_CLASSIC, bbox_inches="tight")
plt.close()

# Classic all-scaffolds bars
alldiff_c = dec_c.sort_values("abs_diff_classic", ascending=False).copy()
x = np.arange(len(alldiff_c))
width = 0.35
fig_w = max(12, 0.3 * len(alldiff_c))
plt.figure(figsize=(fig_w, 6))
plt.bar(x - width/2, alldiff_c["uniformity_score_A"], width, label="A")
plt.bar(x + width/2, alldiff_c["uniformity_score_B"], width, label="B")
plt.xticks(x, alldiff_c["Scaffold"], rotation=90)
plt.ylabel("Classic uniformity")
plt.title("All scaffolds: classic uniformity (sorted by |A − B|)")
plt.legend()
plt.tight_layout()
plt.savefig(ALLDIFF_PNG_CLASSIC, bbox_inches="tight")
plt.close()

# Density hist (CV_density)
plt.figure()
plt.hist(comp["cv_density_A"].dropna(), bins=30, alpha=0.5, label="A")
plt.hist(comp["cv_density_B"].dropna(), bins=30, alpha=0.5, label="B")
plt.xlabel("CV of SNP density (length-weighted; lower = better)")
plt.ylabel("Count of scaffolds")
plt.title("Density uniformity distribution")
plt.legend()
plt.savefig(HIST_PNG_DENSITY, bbox_inches="tight")
plt.close()

# Density scatter A vs B
plt.figure()
plt.scatter(comp["cv_density_A"], comp["cv_density_B"], s=10)
vals_d = np.concatenate([
    comp["cv_density_A"].dropna().to_numpy(),
    comp["cv_density_B"].dropna().to_numpy()
])
if vals_d.size > 0:
    lo, hi = float(np.nanmin(vals_d)), float(np.nanmax(vals_d))
    plt.plot([lo, hi], [lo, hi])
plt.xlabel("A CV_density")
plt.ylabel("B CV_density")
plt.title("Per-scaffold: CV_density A vs B (y = x)")
plt.savefig(SCATTER_PNG_DENSITY, bbox_inches="tight")
plt.close()

# Density all-scaffolds bars
alldiff_d = dec_d.sort_values("abs_diff_density", ascending=False).copy()
x = np.arange(len(alldiff_d))
width = 0.35
fig_w = max(12, 0.3 * len(alldiff_d))
plt.figure(figsize=(fig_w, 6))
plt.bar(x - width/2, alldiff_d["cv_density_A"], width, label="A")
plt.bar(x + width/2, alldiff_d["cv_density_B"], width, label="B")
plt.xticks(x, alldiff_d["Scaffold"], rotation=90)
plt.ylabel("CV_density (lower = better)")
plt.title("All scaffolds: CV_density (sorted by |A − B|)")
plt.legend()
plt.tight_layout()
plt.savefig(ALLDIFF_PNG_DENSITY, bbox_inches="tight")
plt.close()

# Dispersion scatter (D)
plt.figure()
plt.scatter(comp["dispersion_D_A"], comp["dispersion_D_B"], s=10)
vals_disp = np.concatenate([
    comp["dispersion_D_A"].dropna().to_numpy(),
    comp["dispersion_D_B"].dropna().to_numpy()
])
if vals_disp.size > 0:
    lo, hi = float(np.nanmin(vals_disp)), float(np.nanmax(vals_disp))
    plt.plot([lo, hi], [lo, hi])
plt.xlabel("A dispersion D")
plt.ylabel("B dispersion D")
plt.title("Per-scaffold: dispersion D A vs B (y = x)")
plt.savefig(SCATTER_PNG_DISPERSION, bbox_inches="tight")
plt.close()

# =======================
# Slides (python-pptx)
# =======================
# Quick stats (classic)
wins_A_classic = int((comp["uniformity_score_A"] < comp["uniformity_score_B"]).sum())
wins_B_classic = int((comp["uniformity_score_B"] < comp["uniformity_score_A"]).sum())
ties_classic   = int(((comp["uniformity_score_A"] == comp["uniformity_score_B"]) |
                      (comp["uniformity_score_A"].isna()) |
                      (comp["uniformity_score_B"].isna())).sum())
total_scaffolds = int(comp.shape[0])

med_A_classic = float(np.nanmedian(comp["uniformity_score_A"]))
med_B_classic = float(np.nanmedian(comp["uniformity_score_B"]))
mean_A_classic= float(np.nanmean(comp["uniformity_score_A"]))
mean_B_classic= float(np.nanmean(comp["uniformity_score_B"]))
improvement_pct_classic = (overall_A_classic - overall_B_classic) / overall_A_classic * 100.0 \
                          if overall_A_classic else np.nan

# Quick stats (density)
wins_A_density = int((comp["cv_density_A"] < comp["cv_density_B"]).sum())
wins_B_density = int((comp["cv_density_B"] < comp["cv_density_A"]).sum())
med_A_density  = float(np.nanmedian(comp["cv_density_A"]))
med_B_density  = float(np.nanmedian(comp["cv_density_B"]))
mean_A_density = float(np.nanmean(comp["cv_density_A"]))
mean_B_density = float(np.nanmean(comp["cv_density_B"]))
improvement_pct_density = (overall_A_density - overall_B_density) / overall_A_density * 100.0 \
                          if overall_A_density else np.nan

# Pull Wilcoxon rows for slide text
def pull_wilcox(metric):
    row = wilcoxon_summary[wilcoxon_summary["metric"] == metric]
    if row.empty:
        return None
    r = row.iloc[0]
    return dict(n=int(r["n"]), p=float(r["p_two_sided"]), r_eff=float(r["effect_r"]), med=float(r["median_delta"]))

w_classic  = pull_wilcox("classic_uniformity")
w_density  = pull_wilcox("cv_density")
w_disp     = pull_wilcox("dispersion_D")

prs = Presentation()

# Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Choosing a Window Set for ARGweaver"
slide.placeholders[1].text = "A: windows_100Mb_equal_lenght.tsv  vs  B: balanced_scaffold_windows_100Mb.tsv"

# Executive summary (classic)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Executive summary (classic uniformity)"
tf = slide.placeholders[1].text_frame
tf.clear()
for i, line in enumerate([
    f"Overall classic uniformity (lower is better): A = {overall_A_classic:.3f}  |  B = {overall_B_classic:.3f}  (B is {improvement_pct_classic:.1f}% lower).",
    f"Wins by scaffold (classic): A = {wins_A_classic}  |  B = {wins_B_classic}  |  ties/NA = {ties_classic}  (total = {total_scaffolds}).",
    f"Median classic: A = {med_A_classic:.3f}  |  B = {med_B_classic:.3f}",
    f"Mean classic:   A = {mean_A_classic:.3f}  |  B = {mean_B_classic:.3f}",
    f"Wilcoxon (classic): n={w_classic['n'] if w_classic else 'NA'}, p={w_classic['p'] if w_classic else 'NA':.3g}, r={w_classic['r_eff'] if w_classic else float('nan'):.3g}, median Δ(B−A)={w_classic['med'] if w_classic else float('nan'):.3g}",
    "Recommendation: Prefer the set with consistently lower classic AND density metrics."
]):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p.text = line
    p.font.size = Pt(18)

# Classic figures
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Classic uniformity distribution"
slide.shapes.add_picture(HIST_PNG_CLASSIC, Inches(0.8), Inches(1.5), height=Inches(4.5))

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Classic A vs B (y = x)"
slide.shapes.add_picture(SCATTER_PNG_CLASSIC, Inches(0.8), Inches(1.5), height=Inches(4.5))

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "All scaffolds: classic (sorted by |A − B|)"
slide.shapes.add_picture(ALLDIFF_PNG_CLASSIC, Inches(0.5), Inches(1.3), height=Inches(4.8))

# NEW: Density/Dispersion significance slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Density-based uniformity & dispersion: significance"
tf = slide.placeholders[1].text_frame
tf.clear()
for txt in [
    "Per window density d_i = SNPs_i / length_i (SNPs per bp). Length-weighted CV_density: lower = more even density.",
    "Dispersion D = Pearson chi-square / (k−1) vs counts ∝ length; D≈1 good, D≫1 uneven.",
    f"Overall CV_density (length-weighted): A = {overall_A_density:.3f} | B = {overall_B_density:.3f} (B is {improvement_pct_density:.1f}% lower).",
    f"Overall dispersion D (window-weighted): A = {overall_A_disp:.3f} | B = {overall_B_disp:.3f}",
    f"Wilcoxon (CV_density): n={w_density['n'] if w_density else 'NA'}, p={w_density['p'] if w_density else 'NA':.3g}, r={w_density['r_eff'] if w_density else float('nan'):.3g}, median Δ(B−A)={w_density['med'] if w_density else float('nan'):.3g}",
    f"Wilcoxon (Dispersion D): n={w_disp['n'] if w_disp else 'NA'}, p={w_disp['p'] if w_disp else 'NA':.3g}, r={w_disp['r_eff'] if w_disp else float('nan'):.3g}, median Δ(B−A)={w_disp['med'] if w_disp else float('nan'):.3g}",
]:
    p = tf.add_paragraph() if len(tf.paragraphs[0].text) else tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(18)

# Density figures
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "CV_density distribution (A vs B)"
slide.shapes.add_picture(HIST_PNG_DENSITY, Inches(0.8), Inches(1.5), height=Inches(4.5))

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Per-scaffold CV_density: A vs B (y = x)"
slide.shapes.add_picture(SCATTER_PNG_DENSITY, Inches(0.8), Inches(1.5), height=Inches(4.5))

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "All scaffolds: CV_density (sorted by |A − B|)"
slide.shapes.add_picture(ALLDIFF_PNG_DENSITY, Inches(0.5), Inches(1.3), height=Inches(4.8))

# Dispersion figure
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Per-scaffold dispersion D: A vs B (y = x)"
slide.shapes.add_picture(SCATTER_PNG_DISPERSION, Inches(0.8), Inches(1.5), height=Inches(4.5))

# Save PPTX
prs.save(PPTX_OUT)

# Done
print(f"Saved: {PER_SCAFFOLD_OUT}")
print(f"Saved: {OVERALL_OUT_CLASSIC}")
print(f"Saved: {OVERALL_OUT_DENSITY}")
print(f"Saved: {DECISIVE_OUT_CLASSIC}")
print(f"Saved: {DECISIVE_OUT_DENSITY}")
print(f"Saved: {WILCOXON_SUMMARY_OUT}")
print(f"Saved: {HIST_PNG_CLASSIC}")
print(f"Saved: {SCATTER_PNG_CLASSIC}")
print(f"Saved: {ALLDIFF_PNG_CLASSIC}")
print(f"Saved: {HIST_PNG_DENSITY}")
print(f"Saved: {SCATTER_PNG_DENSITY}")
print(f"Saved: {ALLDIFF_PNG_DENSITY}")
print(f"Saved: {SCATTER_PNG_DISPERSION}")
print(f"Saved: {PPTX_OUT}")

