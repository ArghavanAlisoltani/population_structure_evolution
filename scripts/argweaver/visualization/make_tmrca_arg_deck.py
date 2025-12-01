# make_tmrca_arg_deck.py
# Creates a 12-slide PPTX with speaker notes + your three result figures.

from pptx import Presentation
from pptx.util import Inches, Pt

# ---- EDIT THESE PATHS IF YOUR IMAGES ARE IN A DIFFERENT LOCATION ----
IMG_TMRCA = "path/to/your/591d7eb4-455c-4a45-9985-23adafe6cf8f.png"   # genome-wide TMRCA line + CI
IMG_MRNA  = "path/to/your/1e6a57f3-963a-4ae2-b5ed-f00684a9ef2d.png"   # mRNA overlap
IMG_TE    = "path/to/your/1c8dbb6e-0737-4c17-a091-e62fc5362726.png"   # TE overlap
OUTFILE   = "TMRCA_ARG_intro_deck.pptx"
# ---------------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bullet_slide(title, bullets, notes, layout=1):
    s = prs.slides.add_slide(prs.slide_layouts[layout])  # 1 = Title & Content
    s.shapes.title.text = title
    tf = s.placeholders[1].text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(20)
    s.notes_slide.notes_text_frame.text = notes

def add_image_slide(title, image_path, notes, height=4.8):
    s = prs.slides.add_slide(prs.slide_layouts[5])  # 5 = Title Only
    s.shapes.title.text = title
    s.shapes.add_picture(image_path, Inches(0.5), Inches(1.3), height=Inches(height))
    s.notes_slide.notes_text_frame.text = notes

# 1
add_bullet_slide(
    "Reading the genome’s family history in lodgepole pine",
    ["TMRCA & ancestral recombination graphs (ARGs): what they are, why they matter",
     "From local genealogies → breeding decisions"],
    "We read where ancestry is old vs. young across the genome using TMRCA/ARGs, and connect that to SNPs, genes, TEs, and trait signals to inform breeding."
)

# 2
add_bullet_slide(
    "Why should foresters care?",
    ["Climate, drought, fire, insects = moving targets",
     "We already use GWAS & genomic prediction; genealogies add timing",
     "Older vs. younger ancestry can flag durable vs. recent adaptation"],
    "TMRCA adds a time axis: old segments may hold long-standing, robust variation; very young segments can mark recent sweeps/introgression."
)

# 3
add_bullet_slide(
    "Jargon buster (one-liners)",
    ["SNP: single DNA change",
     "Haplotype: a stretch of SNPs inherited together",
     "Recombination: swapping chromosome pieces during meiosis → mosaic genomes",
     "Coalescence: tracing lineages backwards until they merge (share an ancestor)",
     "TMRCA: time to most recent common ancestor of a segment (generations)",
     "Local tree: the family tree for one segment",
     "ARG: all the local trees plus where recombination switches them",
     "Scaffold: long piece of the assembled genome; CI: confidence interval"],
    "If you remember only two: coalescence is lineages merging backward; TMRCA is the age of that merge for a segment."
)

# 4
add_bullet_slide(
    "Coalescence & TMRCA (the idea)",
    ["Think backward: two copies of a segment eventually meet at an ancestor",
     "The time of that meeting = TMRCA",
     "High TMRCA → old, diverse ancestry; Low TMRCA → recent shared ancestor",
     "Recombination means each segment can have a different TMRCA"],
    "Recombination shuffles ancestry, so a chromosome can be ancient in one region but young 50 kb away. We track TMRCA as a genome-long curve."
)

# 5
add_bullet_slide(
    "From recombination to ARG & why ARGweaver",
    ["Recombination creates many local genealogies",
     "ARG = stitched map of those genealogies along the genome",
     "ARGweaver infers local trees + TMRCA per segment with CIs",
     "With a reference genome we map segments and overlay genes/TEs/GWAS"],
    "Our reference enables phasing and smoothed TMRCA windows with clear coordinates for annotation overlays."
)

# 6
add_bullet_slide(
    "Our questions",
    ["Are high-TMRCA regions closer to genes?",
     "Are TEs depleted or enriched in high-TMRCA (and which classes)?",
     "Do GWAS hits fall more often in high-TMRCA segments?",
     "Can we tell region-level stories that guide breeding?"],
    "These tie timing (TMRCA) to function (genes/TEs) and utility (GWAS peaks)."
)

# 7
add_bullet_slide(
    "Objectives",
    ["Infer local TMRCA segments genome-wide; smooth in windows",
     "Annotate overlap with mRNA and TE catalogs",
     "Compare top-percentile TMRCA vs. background with Fisher tests",
     "Quantify GWAS coverage across TMRCA thresholds (5–25%)",
     "Create interpretation-first plots for decision makers"],
    "Top “x% TMRCA” is our operational definition of “old.” Confidence intervals communicate uncertainty."
)

# 8
add_bullet_slide(
    "Hypotheses",
    ["H1: High-TMRCA regions are gene-rich (older, maintained variation)",
     "H2: Any TE is less frequent in high-TMRCA (constraint near genes)",
     "H3: Gypsy more depleted than Copia near high-TMRCA (size/epigenetic costs)",
     "H4: More GWAS hits lie within top-TMRCA than expected by chance"],
    "These align with conifer genome architecture (large, TE-dense genomes with islands of constraint)."
)

# 9
add_bullet_slide(
    "Data & pipeline (summary)",
    ["Samples: ~1,489 lodgepole pine; GBS SNPs aligned to our reference",
     "TMRCA/ARG: segments with mean TMRCA + CIs; windowed smoothing",
     "Annotations: mRNA (InterPro IDs) and EDTA TE classes/superfamilies",
     "Stats: Fisher tests (top vs background) + GWAS coverage across thresholds"],
    "We merge short adjacent segments on the same scaffold to avoid double-counting mRNA overlaps and weight by overlap length when smoothing."
)

# 10
add_image_slide(
    "Result A — Genome-wide TMRCA profile (windowed mean with CI)",
    IMG_TMRCA,
    "Peaks = older ancestry; valleys = recent. X = scaffolds concatenated; Y = generations. CI ribbons convey uncertainty from inference and data coverage."
)

# 11 (two images on one slide)
s = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
s.shapes.title.text = "Result B — Feature overlap by TMRCA group"
s.shapes.add_picture(IMG_MRNA, Inches(0.3), Inches(1.5), height=Inches(4.0))
s.shapes.add_picture(IMG_TE,   Inches(6.8), Inches(1.5), height=Inches(4.0))
s.notes_slide.notes_text_frame.text = (
    "mRNA overlap: higher in top-TMRCA than background (~36% → ~39%). "
    "Any TE overlap: lower in top-TMRCA (~19.9% → ~17.3%). "
    "Interpretation: Older segments are modestly more genic and less TE-occupied, consistent with constraint near genes."
)

# 12
add_bullet_slide(
    "Result C — GWAS, TMRCA & implications",
    ["Coverage example: ~22.4% of hits in top 5% TMRCA; ~35.2% in top 25%",
     "Within top 25%, ~55.6% WWD and ~38.9% HT+δ¹³C",
     "Take-home: many trait-associated loci sit in older ancestry blocks",
     "Breeding links: prioritize stable old haplotypes; design crosses; add TMRCA features to GP",
     "Limits & next steps: right-skewed ages, ARG uncertainty; zoom to candidates; validate across trials"],
    "A WWD hit inside an old, gene-rich, TE-poor block is a strong candidate for broad deployment; young+hit may signal local adaptation worth targeted deployment."
)

prs.save(OUTFILE)
print(f"Wrote {OUTFILE}")
