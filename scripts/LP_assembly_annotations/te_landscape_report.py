# te_landscape_report.py
# Create figures + a PowerPoint slide deck from sequence_ontology.tsv

import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
import os

# ---------- CONFIG ----------
INPUT_TSV = "sequence_ontology.tsv"   # change if needed
OUT_DIR   = "."                         # change if needed

FIG1_PATH = os.path.join(OUT_DIR, "te_composition_grouped.png")
FIG2_PATH = os.path.join(OUT_DIR, "top15_so_terms.png")
PPTX_PATH = os.path.join(OUT_DIR, "Conifer_TE_landscape_slides.pptx")

# ---------- LOAD ----------
df = pd.read_csv(INPUT_TSV, sep="\t")
df.columns = [c.strip() for c in df.columns]
# Expecting columns: Var1, Freq
df["Var1"] = df["Var1"].astype(str).str.strip()
df["Freq"] = pd.to_numeric(df["Freq"], errors="coerce").fillna(0).astype(int)
total = int(df["Freq"].sum())

# ---------- GROUPING ----------
def classify(name: str) -> str:
    n = name.lower()
    # LTR retrotransposons
    if "ltr" in n or "gypsy" in n or "copia" in n:
        return "LTR retrotransposons (Gypsy/Copia/LTR)"
    # TIR DNA transposons
    if any(tag in n for tag in ["mutator", "cacta", "pif", "harbinger", "hat", "tc1", "mariner"]):
        return "TIR DNA transposons (Mutator/CACTA/PIF/hAT/Tc1-Mariner)"
    # Helitrons
    if "helitron" in n:
        return "Helitron (RC/Helitrons)"
    # LINEs
    if any(tag in n for tag in ["line_", " line", "l1", "l2", "rte", "jockey", "cr1", "tad1"]) or n.startswith("line"):
        return "LINEs"
    # SINEs
    if "sine" in n:
        return "SINEs"
    # Other repeat features
    if any(tag in n for tag in ["repeat_region", "repeat_fragment", "target_site_duplication"]):
        return "Other repeat features"
    return "Other/Unclassified"

df["Group"] = df["Var1"].map(classify)
grouped = (
    df.groupby("Group", as_index=False)["Freq"]
      .sum()
      .sort_values("Freq", ascending=False)
)
grouped["Percent"] = grouped["Freq"] / grouped["Freq"].sum() * 100

# ---------- FIGURE 1: grouped composition ----------
plt.figure(figsize=(9, 6))
plt.bar(grouped["Group"], grouped["Freq"])
plt.title("TE Composition (Grouped)")
plt.xlabel("Class")
plt.ylabel("Count")
plt.xticks(rotation=25, ha="right")
plt.tight_layout()
plt.savefig(FIG1_PATH, dpi=300, bbox_inches="tight")
plt.close()

# ---------- FIGURE 2: top 15 SO terms ----------
topn = df.sort_values("Freq", ascending=False).head(15)
plt.figure(figsize=(9, 6))
plt.bar(topn["Var1"], topn["Freq"])
plt.title("Top 15 SO Terms by Count")
plt.xlabel("SO Term")
plt.ylabel("Count")
plt.xticks(rotation=45, ha="right")
plt.tight_layout()
plt.savefig(FIG2_PATH, dpi=300, bbox_inches="tight")
plt.close()

# ---------- HEADLINE COUNTS ----------
headline_counts = {
    "Total annotated features": total,
    "LTR retrotransposons": int(grouped.loc[grouped["Group"].str.startswith("LTR"), "Freq"].sum()),
    "TIR DNA transposons": int(grouped.loc[grouped["Group"].str.startswith("TIR"), "Freq"].sum()),
    "Helitrons": int(grouped.loc[grouped["Group"].str.startswith("Helitron"), "Freq"].sum()),
    "LINEs": int(grouped.loc[grouped["Group"].str.startswith("LINE"), "Freq"].sum()),
    "SINEs": int(grouped.loc[grouped["Group"].str.startswith("SINE"), "Freq"].sum()),
}

# ---------- POWERPOINT ----------
prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "TE Landscape & Implications for Conifer Breeding"
slide.placeholders[1].text = "Auto-generated from sequence_ontology.tsv"

def add_textbox(slide, left_in, top_in, width_in, height_in, text, font_size=14, bold=False):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    return box

# Slide: grouped composition + stats (blank layout)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(slide, 0.5, 0.3, 9, 0.6, "Grouped TE Composition", font_size=24, bold=True)
stats_lines = [
    f"Total annotated features: {total}",
    f"LTR retrotransposons: {headline_counts['LTR retrotransposons']}",
    f"TIR DNA transposons: {headline_counts['TIR DNA transposons']}",
    f"Helitrons: {headline_counts['Helitrons']}",
    f"LINEs: {headline_counts['LINEs']}",
    f"SINEs: {headline_counts['SINEs']}",
]
add_textbox(slide, 0.6, 1.1, 4.6, 3.0, "\n".join(stats_lines), font_size=14)
slide.shapes.add_picture(FIG1_PATH, Inches(5.2), Inches(1.1), height=Inches(4.5))

# Slide: top 15 SO terms
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(slide, 0.5, 0.3, 9, 0.6, "Top 15 SO Terms", font_size=24, bold=True)
add_textbox(slide, 0.6, 1.0, 9, 0.5, "Most abundant TE-related SO terms by count.", font_size=14)
slide.shapes.add_picture(FIG2_PATH, Inches(0.6), Inches(1.6), height=Inches(4.5))

# Slide: interpretation
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(slide, 0.5, 0.3, 9, 0.6, "Why this matters (Conifers & Breeding)", font_size=24, bold=True)
interp = [
    "• Conifer genomes are dominated by LTR retrotransposons (esp. Ty3/Gypsy), consistent with this profile.",
    "• LTR-dense regions can inflate genome size and suppress recombination, broadening QTLs and shaping LD.",
    "• DNA transposons (Mutator, CACTA, PIF/Harbinger, hAT) often insert near/within genes—useful polymorphic markers.",
    "• Helitrons can capture gene fragments, generating novel transcripts and regulatory variation.",
    "• TE insertion polymorphisms (TIPs) can be trait-causal—include TIPs alongside SNPs in GWAS/GS."
]
add_textbox(slide, 0.6, 1.0, 9, 4.5, "\n".join(interp), font_size=16)

# Slide: practical next steps
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(slide, 0.5, 0.3, 9, 0.6, "Practical next steps", font_size=24, bold=True)
steps = [
    "• Curate a species-specific TE library; annotate with RepeatModeler2 + LTR_retriever; use SO terms in GFF3.",
    "• Detect TE insertion polymorphisms (TIPs) and intersect with promoters/UTRs of candidate genes.",
    "• Use TE proximity/density/age as covariates in GWAS/GS; validate high-value TIPs with long reads.",
    "• Apply repeat-aware mapping/filters; flag variants in repeat_region/repeat_fragment for extra caution."
]
add_textbox(slide, 0.6, 1.0, 9, 4.5, "\n".join(steps), font_size=16)

# Slide: citations
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(slide, 0.5, 0.3, 9, 0.6, "Selected references", font_size=24, bold=True)
cites = [
    "Nystedt et al. (2013) Nature – Norway spruce genome & LTR-RT accumulation. https://www.nature.com/articles/nature12211",
    "De La Torre et al. (2014) Genome Biology – Conifer genomics & TE-driven expansion. https://genomebiology.biomedcentral.com/articles/10.1186/gb-2013-14-6-122",
    "Zhang et al. (2020) Evolutionary Bioinformatics – TE applications in conifer breeding. https://journals.sagepub.com/doi/10.1177/1176934320930263",
    "Thomas et al. (2024) Trends in Genetics – Helitrons and developmental novelties. https://www.sciencedirect.com/science/article/pii/S0168952524000295",
    "Li et al. (2012) Genetics – Helitron gene-fragment capture in maize. https://academic.oup.com/genetics/article/190/3/965/6062435",
    "Studer et al. (2011) tb1/Hopscotch domestication example. https://peerj.com/articles/900.pdf",
    "Jiang et al. (2009) Rider LTR-RT & tomato traits; plus 2019 overview. https://onlinelibrary.wiley.com/doi/pdf/10.1111/j.1365-313X.2009.03946.x ; https://journals.plos.org/plosgenetics/article?id=10.1371/journal.pgen.1008370",
    "Sequence Ontology – official site. https://www.sequenceontology.org/",
    "Ensembl GFF3 uses SO terms. https://grch37.ensembl.org/info/website/upload/gff3.html",
]
add_textbox(slide, 0.6, 1.0, 9, 4.5, "\n".join(cites), font_size=12)

# ---------- SAVE ----------
prs.save(PPTX_PATH)

print("Saved:")
print("  ", FIG1_PATH)
print("  ", FIG2_PATH)
print("  ", PPTX_PATH)

